"""Complete PPTX visual audit script — run against generated report files."""
import sys, os, glob
# Force UTF-8 output on Windows
sys.stdout.reconfigure(encoding="utf-8", errors="replace")
sys.path.insert(0, os.path.join(os.path.dirname(__file__), "..", "backend"))

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.shapes import MSO_SHAPE_TYPE
import io


def emu2in(emu):
    return round(emu / 914400, 2) if emu else 0.0


def get_fill(shape):
    try:
        fill = shape.fill
        if fill is None:
            return "no fill"
        ft = str(fill.type) if fill.type is not None else ""
        if "SOLID" in ft:
            try:
                return "#" + str(fill.fore_color.rgb)
            except Exception:
                return "solid(?)"
        if "GRADIENT" in ft:
            return "gradient"
        if "PATTERN" in ft:
            return "pattern"
        return "no fill"
    except Exception:
        return "no fill"


def get_border(shape):
    try:
        line = shape.line
        if line is None:
            return "none"
        ft = str(line.fill.type) if line.fill and line.fill.type is not None else ""
        if "BACKGROUND" in ft or not ft:
            return "none"
        w = round(line.width / 12700, 1) if line.width else 0
        try:
            c = "#" + str(line.color.rgb)
        except Exception:
            c = "?"
        return f"{c} {w}pt"
    except Exception:
        return "none"


def get_bg(slide):
    try:
        fill = slide.background.fill
        if fill.type is not None and "SOLID" in str(fill.type):
            return "#" + str(fill.fore_color.rgb)
    except Exception:
        pass
    return "default(white)"


def stype_str(shape):
    mapping = {1: "SHAPE", 5: "FREEFORM", 6: "GROUP", 13: "PICTURE",
               14: "PLACEHOLDER", 17: "TEXTBOX", 19: "TABLE"}
    try:
        return mapping.get(shape.shape_type, f"TYPE_{shape.shape_type}")
    except Exception:
        return "UNKNOWN"


def get_text(shape):
    if not shape.has_text_frame:
        return None
    full = shape.text_frame.text.strip()
    fs = None; color = None; bold = False; align_s = None
    for para in shape.text_frame.paragraphs:
        try:
            align_s = str(para.alignment) if para.alignment else None
        except Exception:
            pass
        for run in para.runs:
            try:
                if run.font.size:
                    fs = round(run.font.size / 12700, 1)
            except Exception:
                pass
            try:
                if run.font.color and run.font.color.rgb:
                    color = "#" + str(run.font.color.rgb)
            except Exception:
                pass
            try:
                bold = bool(run.font.bold)
            except Exception:
                pass
            break
        break
    return {"text": full[:60], "full_len": len(full), "fs": fs,
            "color": color or "inherit", "bold": bold, "align": align_s}


def get_shadow(shape):
    try:
        xml = shape._element.xml
        if "<a:outerShdw" in xml:
            return "outer"
        if "<a:innerShdw" in xml:
            return "inner"
    except Exception:
        pass
    return "no"


def check_overlaps(sdata):
    out = []
    for i, s1 in enumerate(sdata):
        for j, s2 in enumerate(sdata):
            if j <= i:
                continue
            ox = max(0, min(s1["x"]+s1["w"], s2["x"]+s2["w"]) - max(s1["x"], s2["x"]))
            oy = max(0, min(s1["y"]+s1["h"], s2["y"]+s2["h"]) - max(s1["y"], s2["y"]))
            oa = ox * oy
            mn = min(s1["w"]*s1["h"], s2["w"]*s2["h"])
            if mn < 0.01:
                continue
            pct = oa / mn
            if pct > 0.5 and oa > 0.5 and min(s1["h"], s2["h"]) > 0.1:
                out.append(f"Shapes id={s1['id']} & id={s2['id']} overlap {pct*100:.0f}% "
                           f"({s1['type']} vs {s2['type']})")
    return out


SLIDE_NAMES = ["Cover", "Executive Summary", "KPI Scorecard", "Website Performance",
               "Meta Ads Performance", "Key Wins", "Concerns & Recommendations",
               "Next Steps", "Custom Section"]

files = sorted(glob.glob("F:/Sapienbotics/ClaudeCode/reportpilot/test_generated_pptx/*.pptx"))
print(f"Found {len(files)} files to audit\n")

all_verdicts = {}

for fpath in files:
    fname = os.path.basename(fpath)
    prs = Presentation(fpath)
    issues = []

    print("=" * 80)
    print(f"=== {fname} ===")
    print(f"Slide size: {emu2in(prs.slide_width)}\" x {emu2in(prs.slide_height)}\"  |  Slides: {len(prs.slides)}")
    print("=" * 80)

    footer_ys = []
    title_ys = []

    for si, slide in enumerate(prs.slides):
        sname = SLIDE_NAMES[si] if si < len(SLIDE_NAMES) else f"Extra {si+1}"
        bg = get_bg(slide)
        print(f"\n--- Slide {si+1} ({sname}) ---")
        print(f"Background: {bg}   Shapes: {len(slide.shapes)}")

        sdata = []
        pic_count = 0
        has_footer = False
        has_cta = False
        kpi_l = kpi_v = kpi_c = 0
        kpi_green = kpi_red = kpi_gray = 0
        placeholders = []

        for shape in slide.shapes:
            x = emu2in(shape.left)
            y = emu2in(shape.top)
            w = emu2in(shape.width)
            h = emu2in(shape.height)
            st = stype_str(shape)
            fl = get_fill(shape)
            bd = get_border(shape)
            sh = get_shadow(shape)
            sdata.append({"id": shape.shape_id, "x": x, "y": y, "w": w, "h": h, "type": st})

            sh_tag = f" shadow={sh}" if sh != "no" else ""

            if st == "PICTURE":
                pic_count += 1
                try:
                    ct = shape.image.content_type
                except Exception:
                    ct = "?"
                print(f"  [{st:10s}] x={x:5.2f} y={y:5.2f} w={w:5.2f} h={h:5.2f}  — {ct}")
                continue

            ti = get_text(shape)
            if ti:
                txt = ti["text"]
                fs_str = f"{ti['fs']}pt" if ti["fs"] else "?pt"
                b_str = " bold" if ti["bold"] else ""

                if "{{" in txt and "}}" in txt:
                    placeholders.append(txt)
                if "Confidential" in txt or "Page " in txt:
                    has_footer = True
                    footer_ys.append(round(y, 1))
                if ti["fs"] and ti["fs"] >= 20 and si > 0:
                    title_ys.append(round(y, 1))
                if "Questions?" in txt or "schedule a call" in txt.lower():
                    has_cta = True

                # KPI tracking on slide 3
                if si == 2:
                    if ti["fs"] and ti["fs"] <= 10 and txt == txt.upper() and len(txt) > 2:
                        kpi_l += 1
                    elif ti["fs"] and ti["fs"] >= 20:
                        kpi_v += 1
                    if "%" in txt or txt == "N/A":
                        kpi_c += 1
                        # Color check
                        for para in shape.text_frame.paragraphs:
                            for run in para.runs:
                                t = run.text.strip()
                                try:
                                    rgb = str(run.font.color.rgb) if run.font.color and run.font.color.rgb else None
                                except Exception:
                                    rgb = None
                                if t.startswith("+") and "%" in t:
                                    if rgb == "059669":
                                        kpi_green += 1
                                    elif rgb:
                                        issues.append(f"WARNING: S3 \"{t}\" color={rgb} (want #059669)")
                                elif t.startswith("-") and "%" in t:
                                    if rgb == "E11D48":
                                        kpi_red += 1
                                    elif rgb:
                                        issues.append(f"WARNING: S3 \"{t}\" color={rgb} (want #E11D48)")
                                elif t == "N/A":
                                    if rgb == "94A3B8":
                                        kpi_gray += 1

                # Overflow risk
                if ti["full_len"] > 200 and w < 4.0 and ti["fs"] and ti["fs"] >= 12:
                    issues.append(f"WARNING: S{si+1} text overflow risk: {ti['full_len']} chars in {w}\" box at {fs_str}")

                print(f"  [{st:10s}] x={x:5.2f} y={y:5.2f} w={w:5.2f} h={h:5.2f}  fill={fl} bd={bd}{sh_tag}")
                if txt:
                    print(f"              \"{txt}\" {fs_str} {ti['color']}{b_str}")
            else:
                print(f"  [{st:10s}] x={x:5.2f} y={y:5.2f} w={w:5.2f} h={h:5.2f}  fill={fl} bd={bd}{sh_tag}")

            # Edge margin check (skip full-width decorative elements)
            if w < 12.0:
                if x + w > 13.1:
                    issues.append(f"WARNING: S{si+1} shape id={shape.shape_id} right edge at {x+w:.2f}\" (>13.1\")")
                if y + h > 7.35 and h < 6.0:
                    issues.append(f"WARNING: S{si+1} shape id={shape.shape_id} bottom edge at {y+h:.2f}\" (>7.35\")")

        # Overlap check
        for oi in check_overlaps(sdata):
            issues.append(f"WARNING: S{si+1} {oi}")

        # Placeholder leftovers
        for pl in placeholders:
            issues.append(f"CRITICAL: S{si+1} leftover placeholder: \"{pl}\"")

        # Chart check (slides 4 & 5)
        if si in (3, 4):
            if pic_count >= 2:
                print(f"  CHARTS: {pic_count} images embedded")
            else:
                issues.append(f"CRITICAL: S{si+1} ({sname}) only {pic_count} chart images (need 2)")

        # KPI check (slide 3)
        if si == 2:
            print(f"  KPI: {kpi_l} labels, {kpi_v} values, {kpi_c} changes")
            print(f"  KPI colors: {kpi_green} green, {kpi_red} red, {kpi_gray} gray")
            if kpi_l < 6:
                issues.append(f"WARNING: S3 only {kpi_l} KPI labels (want 6)")
            if kpi_v < 6:
                issues.append(f"WARNING: S3 only {kpi_v} KPI values (want 6)")

        # Footer check (slides 2-8)
        if 1 <= si <= 7 and not has_footer:
            issues.append(f"WARNING: S{si+1} ({sname}) missing footer")

        # Cover checks
        if si == 0:
            all_txt = " ".join(shape.text_frame.text for shape in slide.shapes if shape.has_text_frame)
            if len(all_txt.strip()) < 5:
                issues.append(f"CRITICAL: Cover has almost no text content")

        # CTA on slide 8
        if si == 7 and not has_cta:
            issues.append(f"INFO: S8 missing CTA band")

    # Cross-slide consistency
    if len(set(footer_ys)) > 3:
        issues.append(f"WARNING: Footers at {len(set(footer_ys))} different y positions: {sorted(set(footer_ys))}")

    # ── Print issues ────────────────────────────────────────────────
    print("\n" + "-" * 60)
    crit = [i for i in issues if i.startswith("CRITICAL")]
    warn = [i for i in issues if i.startswith("WARNING")]
    info = [i for i in issues if i.startswith("INFO")]
    if not issues:
        print("  No issues found")
    for i in crit:
        print(f"  CRIT  {i}")
    for i in warn:
        print(f"  WARN  {i}")
    for i in info:
        print(f"  INFO  {i}")

    v = "PASS" if not crit and len(warn) <= 2 else ("MAJOR ISSUES" if crit else "NEEDS FIXES")
    print(f"\n  VERDICT: {v}  ({len(crit)} critical, {len(warn)} warnings, {len(info)} info)")
    all_verdicts[fname] = (v, len(crit), len(warn), len(info))
    print()

print("\n" + "=" * 80)
print("FINAL AUDIT SUMMARY")
print("=" * 80)
for fn, (v, c, w, i) in all_verdicts.items():
    tag = fn.split("(")[-1].rstrip(").pptx")
    print(f"  [{v:14s}] {c}C/{w}W/{i}I  {fn}")
