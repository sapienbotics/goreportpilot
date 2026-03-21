#!/usr/bin/env python3
"""
Adapt Canva PPTX templates for ReportPilot.

Creates 3 template files in backend/templates/pptx/ with {{placeholder}}
markers that the template-based report generator can populate at runtime.

Run once:  python scripts/adapt_canva_templates.py
"""
import os
import sys
import copy
import logging
from lxml import etree
from pptx import Presentation
from pptx.util import Inches, Pt, Emu

logging.basicConfig(level=logging.INFO, format="%(levelname)s: %(message)s")
log = logging.getLogger(__name__)

BASE_DIR = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
SRC_DIR = os.path.join(BASE_DIR, "PPTX_template")
OUT_DIR = os.path.join(BASE_DIR, "backend", "templates", "pptx")

# ── Slide section definitions (in our standard 9-slide order) ────────────────
# Each section defines the placeholder text to insert per shape role.
SECTIONS = {
    "cover": {
        "title": "{{client_name}}",
        "subtitle": "{{report_type}}",
        "body": "{{report_period}}",
        "footer": "Prepared by {{agency_name}}",
    },
    "executive_summary": {
        "title": "Executive Summary",
        "body": "{{executive_summary}}",
    },
    "kpi_scorecard": {
        "title": "Key Performance Indicators",
        # KPI placeholders are handled specially — 6 cards
    },
    "website_traffic": {
        "title": "Website Performance",
        "chart_1": "{{chart_sessions}}",
        "chart_2": "{{chart_traffic}}",
        "body": "{{website_narrative}}",
    },
    "meta_ads": {
        "title": "Paid Advertising — Meta Ads",
        "chart_1": "{{chart_spend}}",
        "chart_2": "{{chart_campaigns}}",
        "body": "{{ads_narrative}}",
    },
    "key_wins": {
        "title": "Key Wins",
        "body": "{{key_wins}}",
    },
    "concerns": {
        "title": "Concerns & Recommendations",
        "body": "{{concerns}}",
    },
    "next_steps": {
        "title": "Next Steps",
        "body": "{{next_steps}}",
    },
    "custom_section": {
        "title": "{{custom_section_title}}",
        "body": "{{custom_section_text}}",
    },
}

SECTION_ORDER = [
    "cover", "executive_summary", "kpi_scorecard", "website_traffic",
    "meta_ads", "key_wins", "concerns", "next_steps", "custom_section",
]


# ── Helpers ──────────────────────────────────────────────────────────────────

def _delete_slide(prs, slide_index):
    """Delete a slide at a given 0-based index."""
    rId = prs.slides._sldIdLst[slide_index].rId
    prs.part.drop_rel(rId)
    del prs.slides._sldIdLst[slide_index]


def _reorder_slides(prs, new_order):
    """Reorder slides. new_order is a list of current 0-based indices."""
    sldIdLst = prs.slides._sldIdLst
    elements = list(sldIdLst)
    ordered = [elements[i] for i in new_order]
    for child in list(sldIdLst):
        sldIdLst.remove(child)
    for elem in ordered:
        sldIdLst.append(elem)


def _replace_shape_text(shape, new_text, preserve_first_run=True):
    """Replace all text in a shape's text frame with new_text.
    Preserves the formatting of the first run."""
    if not shape.has_text_frame:
        return
    tf = shape.text_frame
    if not tf.paragraphs:
        return
    # Get formatting from first run of first paragraph
    first_para = tf.paragraphs[0]
    first_run_fmt = None
    if first_para.runs:
        first_run_fmt = first_para.runs[0]._r  # XML element

    # Clear all paragraphs except the first one
    p_elements = tf._txBody.findall('{http://schemas.openxmlformats.org/drawingml/2006/main}p')
    for p_elem in p_elements[1:]:
        tf._txBody.remove(p_elem)

    # Handle multi-line text
    lines = new_text.split('\n')
    for i, line in enumerate(lines):
        if i == 0:
            # Reuse first paragraph
            para = tf.paragraphs[0]
            # Clear existing runs
            r_elements = para._p.findall('{http://schemas.openxmlformats.org/drawingml/2006/main}r')
            for r in r_elements:
                para._p.remove(r)
            # Add new run with preserved formatting
            run = para.add_run()
            run.text = line
            if first_run_fmt is not None and preserve_first_run:
                # Copy run properties from original
                rPr = first_run_fmt.find('{http://schemas.openxmlformats.org/drawingml/2006/main}rPr')
                if rPr is not None:
                    new_rPr = copy.deepcopy(rPr)
                    run._r.insert(0, new_rPr)
        else:
            # Add new paragraph
            new_p = copy.deepcopy(para._p)
            # Clear runs in the copy
            nsmap = {'a': 'http://schemas.openxmlformats.org/drawingml/2006/main'}
            for r in new_p.findall('a:r', nsmap):
                new_p.remove(r)
            tf._txBody.append(new_p)
            new_para = tf.paragraphs[-1]
            run = new_para.add_run()
            run.text = line
            if first_run_fmt is not None and preserve_first_run:
                rPr = first_run_fmt.find('{http://schemas.openxmlformats.org/drawingml/2006/main}rPr')
                if rPr is not None:
                    new_rPr = copy.deepcopy(rPr)
                    run._r.insert(0, new_rPr)


def _get_text_shapes(slide):
    """Get all shapes with text, sorted by top position then left."""
    text_shapes = []
    for shape in slide.shapes:
        if shape.has_text_frame and shape.text_frame.text.strip():
            text_shapes.append(shape)
    text_shapes.sort(key=lambda s: (s.top or 0, s.left or 0))
    return text_shapes


def _classify_shapes(slide):
    """Classify shapes into title, subtitle, body, small based on font size and position."""
    shapes = _get_text_shapes(slide)
    classified = {"title": None, "subtitle": None, "body": [], "small": []}

    for shape in shapes:
        text = shape.text_frame.text.strip()
        if not text:
            continue
        # Get the max font size in this shape
        max_size = 0
        for para in shape.text_frame.paragraphs:
            for run in para.runs:
                if run.font.size:
                    sz = run.font.size
                    if sz > max_size:
                        max_size = sz

        text_len = len(text)

        if max_size >= Pt(24) and text_len < 100:
            if classified["title"] is None:
                classified["title"] = shape
            elif classified["subtitle"] is None:
                classified["subtitle"] = shape
            else:
                classified["body"].append(shape)
        elif max_size >= Pt(14) and text_len < 60:
            if classified["subtitle"] is None:
                classified["subtitle"] = shape
            else:
                classified["small"].append(shape)
        elif text_len > 80:
            classified["body"].append(shape)
        elif text_len < 30:
            classified["small"].append(shape)
        else:
            classified["body"].append(shape)

    return classified


def _insert_kpi_placeholders(slide):
    """Replace text in a KPI-style slide with {{kpi_N_*}} placeholders.
    Finds shape groups and replaces them with KPI templates."""
    shapes = _get_text_shapes(slide)

    # Find the title shape (biggest font / topmost)
    title_shape = None
    metric_shapes = []

    for shape in shapes:
        text = shape.text_frame.text.strip()
        max_size = 0
        for para in shape.text_frame.paragraphs:
            for run in para.runs:
                if run.font.size and run.font.size > max_size:
                    max_size = run.font.size

        if max_size >= Pt(20) and len(text) < 80:
            if title_shape is None:
                title_shape = shape
            else:
                metric_shapes.append(shape)
        else:
            metric_shapes.append(shape)

    # Replace title
    if title_shape:
        _replace_shape_text(title_shape, "Key Performance Indicators")

    # Replace metric shapes with KPI placeholders
    # Group by position — try to identify 6 KPI regions
    kpi_idx = 0
    for shape in metric_shapes:
        text = shape.text_frame.text.strip()
        if kpi_idx >= 6:
            # Clear excess shapes
            _replace_shape_text(shape, "")
            continue

        # Check if this looks like a label (short, uppercase or small font)
        # or a value (number-like) or a change (has % or +/-)
        if len(text) < 30:
            if any(c in text for c in ['$', '%', '+', '-']) or any(c.isdigit() for c in text):
                # Value or change
                _replace_shape_text(shape, f"{{{{kpi_{kpi_idx}_value}}}}")
            else:
                # Label
                _replace_shape_text(shape, f"{{{{kpi_{kpi_idx}_label}}}}")
                kpi_idx += 1
        elif '%' in text or text.startswith('+') or text.startswith('-'):
            _replace_shape_text(shape, f"{{{{kpi_{max(0, kpi_idx-1)}_change}}}}")
        else:
            _replace_shape_text(shape, f"{{{{kpi_{kpi_idx}_value}}}}")


def _apply_section_placeholders(slide, section_key):
    """Apply placeholders to a slide based on its section type."""
    section = SECTIONS[section_key]

    if section_key == "kpi_scorecard":
        _insert_kpi_placeholders(slide)
        return

    classified = _classify_shapes(slide)

    # Replace title
    if classified["title"] and "title" in section:
        _replace_shape_text(classified["title"], section["title"])

    # Replace subtitle
    if classified["subtitle"] and "subtitle" in section:
        _replace_shape_text(classified["subtitle"], section["subtitle"])
    elif classified["subtitle"] and "body" in section and not classified["body"]:
        # If no body shapes found, use subtitle for body
        _replace_shape_text(classified["subtitle"], section["body"])

    # Replace body shapes
    if classified["body"]:
        if section_key in ("website_traffic", "meta_ads"):
            # For chart slides: first body = chart placeholder, second = narrative
            if len(classified["body"]) >= 2:
                _replace_shape_text(classified["body"][0], section.get("chart_1", ""))
                if len(classified["body"]) >= 3:
                    _replace_shape_text(classified["body"][1], section.get("chart_2", ""))
                    _replace_shape_text(classified["body"][2], section.get("body", ""))
                else:
                    _replace_shape_text(classified["body"][1], section.get("body", ""))
            elif len(classified["body"]) == 1:
                _replace_shape_text(classified["body"][0], section.get("body", ""))
        else:
            # Standard section: use first body for main content
            _replace_shape_text(classified["body"][0], section.get("body", ""))
            # Clear remaining body shapes or use for footer
            for extra in classified["body"][1:]:
                if "footer" in section:
                    _replace_shape_text(extra, section["footer"])
                    section.pop("footer")  # Only use once
                else:
                    _replace_shape_text(extra, "")

    # Handle footer in small shapes
    for small_shape in classified.get("small", []):
        text = small_shape.text_frame.text.strip().lower()
        # Replace website URLs, brand names, dates with footer placeholder
        if any(kw in text for kw in ["www.", ".com", "reallygreat", "thynk", "arowwai",
                                       "january", "august", "2024", "2023", "2030", "2045"]):
            _replace_shape_text(small_shape, "{{footer_text}}")


# ── Template-specific adaptations ────────────────────────────────────────────

def adapt_monochrome_minimalist():
    """Adapt 'Monochrome Minimalist Social Media Marketing Report' → modern_clean.pptx"""
    src = os.path.join(SRC_DIR, "Monochrome Minimalist Social Media Marketing Report Presentation.pptx")
    prs = Presentation(src)
    log.info("Loaded Monochrome Minimalist: %d slides", len(prs.slides))

    # Current slides (0-indexed):
    # 0: Cover (Social Media Marketing Report)
    # 1: Summary (Monthly Objectives, Key Results)
    # 2: Content Overview (with images)
    # 3: Content Overview (table - top posts)
    # 4: Summary (Organic Content metrics table)
    # 5: Performance Overview (Paid Advertising table)
    # 6: Charts/Insights (visual)
    # 7: Insights text
    # 8: Next Steps text
    # 9: Conclusion

    # Delete slide 3 (Content table - not needed) — delete in reverse
    _delete_slide(prs, 3)  # After this: indices shift
    # Now 9 slides remain, indices 0-8
    # New mapping: 0=Cover, 1=Summary, 2=Content(images), 3=Organic, 4=Paid, 5=Charts, 6=Insights, 7=NextSteps, 8=Conclusion

    # Desired order: Cover, ExecSummary, KPI, Website, MetaAds, KeyWins, Concerns, NextSteps, Custom
    # Map: 0→0(Cover), 1→1(Exec), 3→2(KPI/Organic), 2→3(Website/images), 4→4(MetaAds/Paid),
    #       5→5(KeyWins/Charts), 6→6(Concerns/Insights), 7→7(NextSteps), 8→8(Custom/Conclusion)
    _reorder_slides(prs, [0, 1, 3, 2, 4, 5, 6, 7, 8])

    # Apply placeholders to each slide
    for i, section_key in enumerate(SECTION_ORDER):
        slide = prs.slides[i]
        _apply_section_placeholders(slide, section_key)
        log.info("  Slide %d → %s", i + 1, section_key)

    out_path = os.path.join(OUT_DIR, "modern_clean.pptx")
    prs.save(out_path)
    log.info("Saved: %s (%d slides)", out_path, len(prs.slides))
    return out_path


def adapt_black_blue():
    """Adapt 'Black Blue Simple Modern Monthly Report' → dark_executive.pptx"""
    src = os.path.join(SRC_DIR, "Black Blue Simple Modern Monthly Report Presentation.pptx")
    prs = Presentation(src)
    log.info("Loaded Black Blue: %d slides", len(prs.slides))

    # Current slides (0-indexed):
    # 0: Cover (MONTHLY REPORT)
    # 1: Executive Summary
    # 2: Objectives
    # 3: Key Metrics Overview
    # 4: Achievements
    # 5: Challenges & Issues
    # 6: Break Slides (unwanted)
    # 7: Department Breakdown
    # 8: Insights & Learnings
    # 9: Thank You

    # Delete break slide (index 6)
    _delete_slide(prs, 6)
    # Now 9 slides: 0=Cover, 1=ExecSummary, 2=Objectives, 3=KeyMetrics, 4=Achievements,
    #               5=Challenges, 6=DeptBreakdown, 7=Insights, 8=ThankYou

    # Desired: Cover, ExecSummary, KPI, Website, MetaAds, KeyWins, Concerns, NextSteps, Custom
    # Map: 0→0, 1→1, 3→2(KPI), 2→3(Website/Objectives), 6→4(MetaAds/Dept),
    #       4→5(KeyWins/Achievements), 5→6(Concerns/Challenges), 7→7(NextSteps/Insights), 8→8(Custom/ThankYou)
    _reorder_slides(prs, [0, 1, 3, 2, 6, 4, 5, 7, 8])

    for i, section_key in enumerate(SECTION_ORDER):
        slide = prs.slides[i]
        _apply_section_placeholders(slide, section_key)
        log.info("  Slide %d → %s", i + 1, section_key)

    out_path = os.path.join(OUT_DIR, "dark_executive.pptx")
    prs.save(out_path)
    log.info("Saved: %s (%d slides)", out_path, len(prs.slides))
    return out_path


def adapt_white_orange():
    """Adapt 'White Black Orange Modern Social Media Performance' → colorful_agency.pptx"""
    src = os.path.join(SRC_DIR, "White Black Orange Modern Social Media Performance Report Presentation.pptx")
    prs = Presentation(src)
    log.info("Loaded White/Orange: %d slides", len(prs.slides))

    # Current slides (0-indexed):
    # 0: Cover (SOCIAL MEDIA)
    # 1: Monthly Summary (Overview with KPI cards)
    # 2: Platform Performance Overview (table-like)
    # 3: Audience Insights (with image)
    # 4: Best Performing Posts (engagement rates)
    # 5: Content Performance Breakdown
    # 6: Monthly Campaign Summary
    # 7: Ad Campaign Overview
    # 8: Key Learnings & Insights
    # 9: Next Month's Focus

    # All 10 slides → need exactly 9
    # Delete slide 4 (Best Performing Posts) as least relevant
    _delete_slide(prs, 4)
    # Now 9 slides: 0=Cover, 1=Summary, 2=PlatformPerf, 3=Audience, 4=ContentPerf,
    #               5=Campaign, 6=AdCampaign, 7=KeyLearnings, 8=NextMonth

    # Desired: Cover, ExecSummary, KPI, Website, MetaAds, KeyWins, Concerns, NextSteps, Custom
    # Map: 0→0, 1→1(ExecSummary), 2→2(KPI), 3→3(Website), 4→4(MetaAds),
    #       5→5(KeyWins), 6→6(Concerns), 7→7(NextSteps), 8→8(Custom)
    # Already in decent order! No reorder needed.

    for i, section_key in enumerate(SECTION_ORDER):
        slide = prs.slides[i]
        _apply_section_placeholders(slide, section_key)
        log.info("  Slide %d → %s", i + 1, section_key)

    out_path = os.path.join(OUT_DIR, "colorful_agency.pptx")
    prs.save(out_path)
    log.info("Saved: %s (%d slides)", out_path, len(prs.slides))
    return out_path


# ── Main ─────────────────────────────────────────────────────────────────────

def main():
    os.makedirs(OUT_DIR, exist_ok=True)
    log.info("Adapting Canva templates → %s", OUT_DIR)

    results = []
    results.append(adapt_monochrome_minimalist())
    results.append(adapt_black_blue())
    results.append(adapt_white_orange())

    log.info("\n✓ All 3 templates created:")
    for r in results:
        log.info("  %s", r)


if __name__ == "__main__":
    main()
