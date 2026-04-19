"""
Cover Page Presets — Phase 3 (rewritten for bug-fixes).

Five preset visual styles applied as IN-PLACE modifications to the
template's existing cover slide. We NEVER add new shapes on top of the
template — all colour + text changes target the template's pre-designed
shapes. Only the `hero` preset adds two net-new elements (the hero
picture + a semi-dark overlay), carefully z-ordered behind the text.

Shape identification uses placeholder TEXT content ({{client_name}},
{{report_period}}, {{report_type}}, {{agency_name}}, {{agency_email}}).
Because of this, `apply_cover_preset()` MUST run BEFORE the generator's
text-replacement pass — otherwise the tokens have already been
substituted and cannot be used as identifiers.

Presets:
  * default    — no-op; use the template's native cover.
  * minimal    — white bg, dark type, thin accent bar.
  * bold       — full-bleed brand colour with large white title.
  * corporate  — dark-navy header band, white title, brand accent.
  * hero       — full-bleed hero image with 40% dark overlay.
  * gradient   — solid brand-colour header over dark body.
"""
from __future__ import annotations

import io
import logging
from typing import Any, Optional

import httpx
from lxml import etree

logger = logging.getLogger(__name__)


# ---------------------------------------------------------------------------
# Preset style configs
# ---------------------------------------------------------------------------

PRESETS: dict[str, Optional[dict[str, Any]]] = {
    "default": None,

    "minimal": {
        "header_fill":    "FFFFFF",
        "page_bg":        "FFFFFF",
        "headline_color": "0F172A",
        "subtitle_color": "64748B",
        "use_hero_image": False,
    },

    "bold": {
        "header_fill":    "brand",
        "page_bg":        "brand",
        "headline_color": "FFFFFF",
        "subtitle_color": "F1F5F9",
        "use_hero_image": False,
    },

    "corporate": {
        "header_fill":    "1E293B",
        "page_bg":        "FFFFFF",
        "headline_color": "FFFFFF",
        "subtitle_color": "CBD5E1",
        "use_hero_image": False,
    },

    "hero": {
        # Hero image replaces visible header/body — colours are only used
        # for text recolouring so titles stay readable on the image.
        "header_fill":    None,
        "page_bg":        None,
        "headline_color": "FFFFFF",
        "subtitle_color": "F1F5F9",
        "use_hero_image": True,
    },

    "gradient": {
        "header_fill":    "brand",
        "page_bg":        "0F172A",
        "headline_color": "FFFFFF",
        "subtitle_color": "E2E8F0",
        "use_hero_image": False,
    },
}


# Placeholder groups used for SUBSTITUTION of custom headline/subtitle.
# v7 — logs confirm templates have no {{report_date}} token on the cover
# (only {{client_name}}, {{report_period}}, {{report_type}}, {{agency_name}},
# plus logo placeholders and labels). So the v2 choice to target
# {{report_date}} meant custom subtitles were silently discarded. Revert
# to {{report_period}} per the phase-3-fix-v2 note: "if client.cover_subtitle
# is set, replace the period text; if not, keep the date". When the user
# provides a subtitle, the reporting-period date is intentionally replaced.
_HEADLINE_SUB_TOKENS = ("{{client_name}}",)
_SUBTITLE_SUB_TOKENS = ("{{report_period}}",)

# Placeholder groups used for COLORING text on the cover. Broader than
# the substitution groups — every cover text shape should pick up the
# preset palette, even slots we don't overwrite.
#   HEADLINE_COLOR tokens → get `headline_color` from the preset config.
#   SUBTITLE_COLOR tokens → get `subtitle_color`.
_HEADLINE_COLOR_TOKENS = ("{{client_name}}", "{{report_type}}")
_SUBTITLE_COLOR_TOKENS = (
    "{{report_period}}", "{{report_date}}",
    "{{agency_name}}", "{{agency_email}}",
)


# ---------------------------------------------------------------------------
# Public entry point
# ---------------------------------------------------------------------------


def apply_cover_preset(
    prs: Any,
    *,
    preset: str,
    headline: Optional[str],
    subtitle: Optional[str],
    hero_image_url: Optional[str],
    brand_color: Optional[str] = None,
    accent_color: Optional[str] = None,
) -> None:
    """
    Apply a cover preset to slide[0]. Must be called BEFORE the generator's
    text-replacement pass so placeholder tokens are still intact and can be
    used to identify shapes.

    Safe to call with preset='default' — that skips colour + hero logic but
    still honours headline/subtitle text overrides.
    """
    try:
        cover = prs.slides[0]
    except Exception as exc:
        logger.warning("Cover preset: no cover slide: %s", exc)
        return

    brand_hex = _normalise_hex(brand_color) or "4338CA"
    accent_hex = _normalise_hex(accent_color)
    config = PRESETS.get(preset or "default")

    # Order matters:
    #   1. Header/page colours first — pure shape-fill work, unrelated to text.
    #   2. Hero image + overlay — add net-new shapes at the back of z-order.
    #   3. RECOLOUR before SUBSTITUTE. Recolour uses placeholder tokens to
    #      identify shape roles; substitute replaces those tokens with
    #      custom text. Doing recolour second would mean the {{client_name}}
    #      token is already gone by the time recolour tries to find it.
    #   4. Substitute headline/subtitle text (narrow token set — only
    #      {{client_name}} and {{report_date}}, preserving {{report_type}}
    #      and {{report_period}} so they get their real values).
    #   5. Accent bar last — a small decorative overlay.

    if config:
        # v7 — aggressive strip for EVERY non-default preset. The CSS
        # preview renders 3 text roles + header band only; the PPTX was
        # carrying the template's full chrome (PERFORMANCE REPORT label,
        # Prepared by line, report_type, logo placeholder text, divider).
        # For hero, the strip is done inside _insert_hero_image after the
        # image is placed (and deletes the header band too).
        is_hero = bool(config.get("use_hero_image")) and bool(hero_image_url)

        if not is_hero:
            # Non-hero preset — strip chrome but KEEP the header band so
            # _apply_header_and_bg_colours can recolour it.
            try:
                _strip_cover_for_preset(prs, cover, keep_header_band=True)
            except Exception as exc:
                logger.debug("Cover preset: strip failed: %s", exc)

        try:
            _apply_header_and_bg_colours(prs, cover, config, brand_hex)
        except Exception as exc:
            logger.debug("Cover preset: colour apply failed: %s", exc)

        if is_hero:
            try:
                _insert_hero_image(prs, cover, hero_image_url)
            except Exception as exc:
                logger.debug("Cover preset: hero-image apply failed: %s", exc)
            try:
                _ensure_hero_headline_size(cover, min_pt=36)
            except Exception as exc:
                logger.debug("Cover preset: headline size boost failed: %s", exc)

        # Reposition title + subtitle into the vertical centre for EVERY
        # non-default preset (including hero). Template positions sit in
        # the lower half because they were designed for a chrome-heavy
        # layout — without the chrome around them they look visually
        # bottom-weighted. v9 fix.
        try:
            _reposition_cover_text_for_preset(prs, cover)
        except Exception as exc:
            logger.debug("Cover preset: reposition failed: %s", exc)

        try:
            _recolour_cover_text(
                cover,
                headline_hex=config.get("headline_color"),
                subtitle_hex=config.get("subtitle_color"),
            )
        except Exception as exc:
            logger.debug("Cover preset: text recolour failed: %s", exc)

    # Text substitution always runs (honours user overrides even on the
    # 'default' preset).
    _substitute_cover_text(cover, headline=headline, subtitle=subtitle)

    # Accent bar renders for every preset EXCEPT hero — the hero's
    # full-bleed image gets cluttered by a stripe bisecting it.
    if config and accent_hex and preset != "hero":
        try:
            _draw_accent_bar(prs, cover, accent_hex)
        except Exception as exc:
            logger.debug("Cover preset: accent bar failed: %s", exc)


# ---------------------------------------------------------------------------
# Step 1 — header + page background colours
# ---------------------------------------------------------------------------


def _apply_header_and_bg_colours(
    prs: Any, cover: Any, config: dict, brand_hex: str,
) -> None:
    """Recolour the header-band shape and slide background in place."""
    from pptx.dml.color import RGBColor  # noqa: PLC0415

    page_bg = config.get("page_bg")
    if page_bg:
        rgb = _hex_to_rgb(_resolve_brand(page_bg, brand_hex))
        try:
            bg = cover.background
            bg.fill.solid()
            bg.fill.fore_color.rgb = RGBColor(*rgb)
        except Exception as exc:
            logger.debug("Cover: slide background fill failed: %s", exc)

    header_fill = config.get("header_fill")
    if header_fill:
        rgb = _hex_to_rgb(_resolve_brand(header_fill, brand_hex))
        header_shape = _find_header_band(cover, prs)
        if header_shape is not None:
            try:
                header_shape.fill.solid()
                header_shape.fill.fore_color.rgb = RGBColor(*rgb)
                try:
                    header_shape.line.fill.background()
                except Exception:
                    pass
            except Exception as exc:
                logger.debug("Cover: header-band recolour failed: %s", exc)


def _find_header_band(cover: Any, prs: Any) -> Optional[Any]:
    """
    Find the header band shape: topmost non-picture shape that spans at
    least 60% of slide width and sits in the top third. Present in every
    visual template's cover. NEVER returns a text-bearing shape — those
    are title boxes, not the coloured header band.
    """
    slide_w = prs.slide_width
    slide_h = prs.slide_height
    min_w   = int(slide_w * 0.6)
    top_cap = int(slide_h * 0.33)

    candidates = []
    for shape in cover.shapes:
        if shape.shape_type == 13:  # MSO_SHAPE_TYPE.PICTURE
            continue
        # Skip shapes with text — they're title/subtitle boxes, not a band.
        if shape.has_text_frame and (shape.text_frame.text or "").strip():
            continue
        try:
            w = int(shape.width or 0)
            t = int(shape.top or 0)
            if w >= min_w and t <= top_cap:
                candidates.append(shape)
        except Exception:
            continue
    if not candidates:
        return None
    candidates.sort(key=lambda s: (int(s.top or 0), -int(s.width or 0)))
    return candidates[0]


# ---------------------------------------------------------------------------
# Step 2 — hero image (z-ordered at back, 40% dark overlay above it)
# ---------------------------------------------------------------------------


def _insert_hero_image(prs: Any, cover: Any, image_url: str) -> None:
    """
    Insert a hero image full-bleed behind all existing shapes.

    v6 fixes (per user spec + screenshot evidence):
      * Strip is now a WHITELIST. Only three roles survive: the title
        ({{client_name}}), the subtitle ({{report_period}}), and the
        footer agency line ({{agency_name}} in the bottom half). Plus
        pictures. Everything else on the cover — "PERFORMANCE REPORT"
        label, "Prepared by" line, report-type text, divider shapes,
        logo placeholder text — gets deleted.
      * Hero headline font boosted to 36pt minimum so it reads at the
        scale the CSS mockup shows.
      * Accent bar skipped for hero preset (no green stripe bisecting
        the cover).
    """
    from pptx.util import Emu  # noqa: PLC0415

    img_bytes = _download_image(image_url)
    if not img_bytes:
        logger.debug("Cover preset: hero image could not be downloaded")
        return

    _log_cover_shapes(cover, "Hero preset — BEFORE strip")
    _strip_cover_for_preset(prs, cover, keep_header_band=False)
    _log_cover_shapes(cover, "Hero preset — AFTER strip")

    pic = cover.shapes.add_picture(
        io.BytesIO(img_bytes), Emu(0), Emu(0),
        width=prs.slide_width, height=prs.slide_height,
    )
    logger.info(
        "Hero pic inserted: left=%d top=%d width=%d height=%d "
        "(slide is %dx%d EMU = %.2fx%.2f inches)",
        int(pic.left or 0), int(pic.top or 0),
        int(pic.width or 0), int(pic.height or 0),
        int(prs.slide_width), int(prs.slide_height),
        prs.slide_width / 914400, prs.slide_height / 914400,
    )

    # Hero goes to the back of z-order so the surviving text shapes
    # still render on top of it.
    _reorder_to_back(cover, pic)
    logger.info("Hero pic moved to back of spTree")


def _log_cover_shapes(cover: Any, label: str) -> None:
    """
    Dump every shape on the cover slide to the log, one line per shape.
    Phase-3 hero-preset diagnostic — shows what decoration survives the
    strip pass so we can iteratively zero in on stray template chrome.

    Output line format (per shape):
      [idx] name='Name' type=<int> pos=(L,T)" size=(WxH)" fill=<type>
             has_text=<bool> text='first 60 chars'
    """
    try:
        lines = []
        for idx, shape in enumerate(cover.shapes):
            name = getattr(shape, "name", "?")
            st = getattr(shape, "shape_type", "?")
            try:
                l_in = (int(shape.left or 0))  / 914400 if shape.left  is not None else 0.0
                t_in = (int(shape.top or 0))   / 914400 if shape.top   is not None else 0.0
                w_in = (int(shape.width or 0)) / 914400 if shape.width is not None else 0.0
                h_in = (int(shape.height or 0))/ 914400 if shape.height is not None else 0.0
            except Exception:
                l_in = t_in = w_in = h_in = 0.0

            # Fill type — 1=solid, 2=picture, 3=gradient, etc. str() is safest.
            fill_desc = "?"
            try:
                fill_desc = str(shape.fill.type)
            except Exception:
                pass

            has_tf = False
            text_preview = ""
            try:
                if getattr(shape, "has_text_frame", False):
                    has_tf = True
                    text_preview = (shape.text_frame.text or "")[:60].replace("\n", " ")
            except Exception:
                pass

            lines.append(
                f"  [{idx}] name={name!r} type={st} "
                f"pos=({l_in:.2f},{t_in:.2f})\" size=({w_in:.2f}x{h_in:.2f})\" "
                f"fill={fill_desc} has_text={has_tf} text={text_preview!r}"
            )
        logger.info("%s — %d shapes:\n%s", label, len(lines), "\n".join(lines))
    except Exception as exc:
        logger.warning("_log_cover_shapes failed: %s", exc)


def _strip_cover_for_preset(prs: Any, cover: Any, *, keep_header_band: bool) -> int:
    """
    Aggressive whitelist strip applied to EVERY non-default preset.
    Keep only:
      * Pictures (existing logos or other art).
      * The header band shape, if `keep_header_band` — non-hero presets
        recolour it via _apply_header_and_bg_colours.
      * A shape whose text contains {{client_name}} (headline slot).
      * A shape whose text contains {{report_period}} (subtitle slot).
      * A shape whose text contains {{agency_name}} AND sits in the
        bottom half of the slide (the footer agency line — NOT the
        "Prepared by {{agency_name}}" line near the top).

    Everything else is removed from the cover's spTree:
      * Coloured header band (only when keep_header_band=False)
      * "PERFORMANCE REPORT" label
      * "Prepared by ..." line
      * {{report_type}} / {{report_date}} text
      * {{agency_logo}} / {{client_logo}} placeholder text shapes —
        _embed_logos places logos from explicit coords or defaults so
        these text placeholders are safe to delete.
      * Divider lines, accent stripes, any other decoration.

    Returns the count of shapes deleted.
    """
    try:
        from pptx.enum.shapes import MSO_SHAPE_TYPE  # noqa: PLC0415
        PICTURE = MSO_SHAPE_TYPE.PICTURE
    except Exception:
        PICTURE = 13

    sp_tree = cover.shapes._spTree
    slide_h = int(prs.slide_height)
    bottom_threshold = int(slide_h * 0.55)

    # Resolve the header band shape once — reference-compare inside the loop.
    header_band_shape = _find_header_band(cover, prs) if keep_header_band else None

    deleted = 0
    kept: list[str] = []

    for shape in list(cover.shapes):
        st = getattr(shape, "shape_type", None)
        name = getattr(shape, "name", "?")

        # Pictures always stay.
        if st == PICTURE:
            kept.append(f"{name} (picture)")
            continue

        # Keep the header band if asked.
        if header_band_shape is not None and shape is header_band_shape:
            kept.append(f"{name} (header band)")
            continue

        text = ""
        try:
            if getattr(shape, "has_text_frame", False):
                text = (shape.text_frame.text or "")
        except Exception:
            text = ""

        top = 0
        try:
            top = int(shape.top or 0)
        except Exception:
            pass

        keep_reason: Optional[str] = None
        if "{{client_name}}" in text:
            keep_reason = "client_name"
        elif "{{report_period}}" in text:
            keep_reason = "report_period"
        elif "{{agency_name}}" in text and top >= bottom_threshold:
            keep_reason = "footer agency_name"

        if keep_reason:
            kept.append(f"{name} ({keep_reason})")
            continue

        try:
            el = shape._element
            if el.getparent() is sp_tree:
                sp_tree.remove(el)
                deleted += 1
                logger.info(
                    "Preset strip DELETE: name=%r type=%s top_in=%.2f text=%r",
                    name, st, top / 914400, text[:60].replace("\n", " "),
                )
        except Exception as exc:
            logger.debug("Preset strip failed for %r: %s", name, exc)

    logger.info(
        "Preset strip complete (keep_header_band=%s) — deleted=%d kept=%d; kept: %s",
        keep_header_band, deleted, len(kept), kept,
    )
    return deleted


def _reposition_cover_text_for_preset(prs: Any, cover: Any) -> None:
    """
    Move the {{client_name}} and {{report_period}} shapes into the
    vertical centre of the cover for preset-overridden layouts.

    Why: the template positions those shapes for its NATIVE cover
    (with PERFORMANCE REPORT label, Prepared-by line, report_type
    placeholder, etc. surrounding them). Once we strip that chrome,
    the title is left at its original y≈5.0" — dropping it into the
    lower-middle of the slide visually, far below where a title should
    sit on a clean preset cover.

    Repositioning rule:
      * {{client_name}} → top = 35% of slide height     (≈ 2.6" on 7.5")
      * {{report_period}} → top = 50% of slide height    (≈ 3.8")
      * Width extended to full slide − 1" margin, horizontally centred,
        so custom headlines of any length stay inside the cover.

    Footer agency_name is NOT moved — its template position is already
    in the bottom band which works fine.
    """
    from pptx.util import Inches  # noqa: PLC0415

    slide_w = int(prs.slide_width)
    slide_h = int(prs.slide_height)
    margin  = int(Inches(0.5))
    target_w = slide_w - 2 * margin   # full width minus 0.5" margins each side

    for shape in cover.shapes:
        if not getattr(shape, "has_text_frame", False):
            continue
        text = shape.text_frame.text or ""
        try:
            if "{{client_name}}" in text:
                shape.top   = int(slide_h * 0.35)
                shape.left  = margin
                shape.width = target_w
                logger.info(
                    "Reposition: client_name → top=%.2f\" left=%.2f\" w=%.2f\"",
                    shape.top / 914400, shape.left / 914400, shape.width / 914400,
                )
            elif "{{report_period}}" in text:
                shape.top   = int(slide_h * 0.50)
                shape.left  = margin
                shape.width = target_w
                logger.info(
                    "Reposition: report_period → top=%.2f\" left=%.2f\" w=%.2f\"",
                    shape.top / 914400, shape.left / 914400, shape.width / 914400,
                )
        except Exception as exc:
            logger.debug("Cover reposition failed for shape: %s", exc)


def _ensure_hero_headline_size(cover: Any, min_pt: int = 36) -> None:
    """
    Ensure the {{client_name}} headline on the cover renders at ≥ `min_pt`
    points for the hero preset — hero templates often have a smaller
    title because they rely on negative space, but against a hero image
    we want the title to dominate. Runs already ≥ min_pt are left alone.
    """
    from pptx.util import Pt  # noqa: PLC0415
    for shape in cover.shapes:
        if not getattr(shape, "has_text_frame", False):
            continue
        if "{{client_name}}" not in (shape.text_frame.text or ""):
            continue
        for para in shape.text_frame.paragraphs:
            for run in para.runs:
                try:
                    cur = run.font.size
                    if cur is None or cur.pt < min_pt:
                        run.font.size = Pt(min_pt)
                except Exception:
                    continue
        return  # only one title shape


def _reorder_to_back(slide: Any, *shapes: Any) -> None:
    """
    Move one or more shapes to the start of spTree (right after the
    group-property children). Order is preserved — the first shape passed
    ends up deepest at back, the next sits right above it, and so on.
    All other template shapes remain on top.
    """
    if not shapes:
        return
    sp_tree = slide.shapes._spTree

    # Remove each element from wherever it currently sits.
    elements = [s._element for s in shapes]
    for el in elements:
        if el.getparent() is sp_tree:
            sp_tree.remove(el)

    # Find the first index after the last group-property child.
    insert_at = 0
    for i, child in enumerate(sp_tree):
        tag = etree.QName(child).localname
        if tag in ("nvGrpSpPr", "grpSpPr"):
            insert_at = i + 1
        else:
            break

    for offset, el in enumerate(elements):
        sp_tree.insert(insert_at + offset, el)


def _set_shape_fill_alpha(shape: Any, alpha_percent: int) -> None:
    """
    Apply alpha (0-100) to a solid-fill shape. python-pptx doesn't expose
    this on FillFormat; we patch the <a:srgbClr> XML directly.
    """
    spPr = shape.fill._xPr
    nsmap = {"a": "http://schemas.openxmlformats.org/drawingml/2006/main"}
    srgb = spPr.find(".//a:solidFill/a:srgbClr", nsmap)
    if srgb is None:
        return
    for child in srgb.findall("a:alpha", nsmap):
        srgb.remove(child)
    alpha = etree.SubElement(
        srgb,
        "{http://schemas.openxmlformats.org/drawingml/2006/main}alpha",
    )
    alpha.set("val", str(max(0, min(100, alpha_percent)) * 1000))


# ---------------------------------------------------------------------------
# Step 3 — text overrides (run BEFORE generator's replacement pass)
# ---------------------------------------------------------------------------


def _substitute_cover_text(
    cover: Any,
    *,
    headline: Optional[str],
    subtitle: Optional[str],
) -> None:
    """
    Apply custom headline/subtitle overrides on the cover slide.

    * Headline REPLACES {{client_name}} — full substitution. The user's
      custom title takes over the primary slot.

    * Subtitle APPENDS to {{report_period}} — the period token stays in
      place so the real date range still shows after the generator's
      substitution pass runs. Example:
          Original run text:  "{{report_period}}"
          After subtitle 'Test Sub' applied here:
                              "{{report_period}} — Test Sub"
          After generator substitutes {{report_period}} with "April 2026":
                              "April 2026 — Test Sub"

      This matches the v7 feedback: "report time and all are missing
      now" — which was caused by subtitle REPLACING the period rather
      than appending. No subtitle → period shows as-is.

    Token-level operations preserve the run's font/size/colour so
    recoloured or resized runs carry their formatting forward.
    """
    if not headline and not subtitle:
        return

    for shape in cover.shapes:
        if not shape.has_text_frame:
            continue
        for para in shape.text_frame.paragraphs:
            for run in para.runs:
                txt = run.text or ""
                # Headline: replace {{client_name}} with custom text.
                if headline:
                    for token in _HEADLINE_SUB_TOKENS:
                        if token in txt:
                            run.text = txt.replace(token, headline)
                            txt = run.text or ""
                # Subtitle: append after {{report_period}} so the real
                # period is preserved.
                if subtitle:
                    for token in _SUBTITLE_SUB_TOKENS:
                        if token in txt:
                            run.text = txt.replace(token, f"{token} — {subtitle}")
                            txt = run.text or ""


def _recolour_cover_text(
    cover: Any,
    *,
    headline_hex: Optional[str],
    subtitle_hex: Optional[str],
) -> None:
    """
    Walk the cover's text frames and recolour runs containing the known
    placeholder tokens. Runs are recoloured IN PLACE — no new shapes, no
    text rewriting. Because colouring is applied before the generator's
    text-substitution pass, the colour survives when the token is
    replaced with the real value.

    Heuristic:
      • runs containing {{client_name}} / {{report_type}} → headline colour
      • runs containing {{report_period}} / {{report_date}} / {{agency_name}}
        / {{agency_email}} → subtitle colour
      • runs that have ALREADY been substituted (custom headline present)
        are matched by the override text instead, so the colour still
        applies.
    """
    if not headline_hex and not subtitle_hex:
        return

    from pptx.dml.color import RGBColor  # noqa: PLC0415

    head_rgb = _hex_to_rgb(headline_hex) if headline_hex else None
    sub_rgb  = _hex_to_rgb(subtitle_hex) if subtitle_hex else None

    for shape in cover.shapes:
        if not shape.has_text_frame:
            continue
        frame_text = (shape.text_frame.text or "")
        if not frame_text.strip():
            continue
        # Skip logo placeholders handled by _embed_logos.
        if "agency_logo" in frame_text.lower() or "client_logo" in frame_text.lower():
            continue
        # Determine which colour applies based on the tokens present.
        # Uses the broader *_COLOR_TOKENS groups so every cover text
        # shape gets recoloured, even slots that aren't substituted.
        target_rgb = None
        if head_rgb and any(t in frame_text for t in _HEADLINE_COLOR_TOKENS):
            target_rgb = head_rgb
        elif sub_rgb and any(t in frame_text for t in _SUBTITLE_COLOR_TOKENS):
            target_rgb = sub_rgb
        elif head_rgb:
            # Fallback: any remaining cover text gets the headline colour so
            # titles/subheadings stay readable against the new background.
            target_rgb = head_rgb

        if target_rgb is None:
            continue

        for para in shape.text_frame.paragraphs:
            for run in para.runs:
                try:
                    run.font.color.rgb = RGBColor(*target_rgb)
                except Exception:
                    continue


# ---------------------------------------------------------------------------
# Step 4 — accent bar (optional)
# ---------------------------------------------------------------------------


def _draw_accent_bar(prs: Any, cover: Any, accent_hex: str) -> None:
    """
    Draw a thin horizontal accent bar at ~33% of slide height (bottom of
    the header band). Exists only when the user picked an accent colour;
    kept intentionally minimal so it doesn't collide with the template's
    native geometry.
    """
    from pptx.dml.color import RGBColor  # noqa: PLC0415
    from pptx.util import Emu            # noqa: PLC0415

    from pptx.util import Inches  # noqa: PLC0415

    rgb = _hex_to_rgb(accent_hex)
    slide_w = prs.slide_width
    slide_h = prs.slide_height
    # v2 fix: 0.06" was ~4px, effectively invisible. 0.10" reads as a
    # deliberate design element while still staying out of title area.
    bar_h   = Inches(0.10)
    bar_top = int(slide_h * 0.33) - int(bar_h)
    bar     = cover.shapes.add_shape(1, Emu(0), Emu(bar_top), slide_w, bar_h)
    try:
        bar.fill.solid()
        bar.fill.fore_color.rgb = RGBColor(*rgb)
        bar.line.fill.background()
    except Exception as exc:
        logger.debug("Cover: accent bar fill failed: %s", exc)


# ---------------------------------------------------------------------------
# Small helpers
# ---------------------------------------------------------------------------


def _normalise_hex(value: Optional[str]) -> Optional[str]:
    if not value:
        return None
    v = value.strip()
    if v.startswith("#"):
        v = v[1:]
    if len(v) == 6:
        return v.upper()
    return None


def _hex_to_rgb(hex_str: Optional[str]) -> tuple[int, int, int]:
    h = _normalise_hex(hex_str) or "4338CA"
    return int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16)


def _resolve_brand(value: str, brand_hex: str) -> str:
    return brand_hex if value == "brand" else value


def _download_image(url: str, max_bytes: int = 5 * 1024 * 1024) -> Optional[bytes]:
    if not url:
        return None
    try:
        with httpx.Client(timeout=20.0, follow_redirects=True) as client:
            resp = client.get(url)
            resp.raise_for_status()
            data = resp.content
            if len(data) > max_bytes:
                logger.warning("Hero image exceeds %d bytes — skipping", max_bytes)
                return None
            return data
    except Exception as exc:
        logger.debug("Hero image download failed (%s): %s", url, exc)
        return None
