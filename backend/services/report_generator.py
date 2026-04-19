"""
Report generation orchestrator.
Template-based PPTX generation + ReportLab PDF fallback.
Both functions are synchronous — call via asyncio.to_thread() in async endpoints.
"""
import io
import os
import re
import logging
import subprocess
import tempfile
from datetime import datetime
from typing import Dict, Any

from services.translations import t, translate_kpi_label

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN  # type: ignore[attr-defined]

from reportlab.lib.pagesizes import A4
from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
from reportlab.lib.units import inch
from reportlab.lib.enums import TA_LEFT, TA_CENTER  # noqa: F401
from reportlab.platypus import (
    SimpleDocTemplate,
    Paragraph,
    Spacer,
    Table,
    TableStyle,
    HRFlowable,
)
from reportlab.platypus import Image as RLImage
from reportlab.platypus import PageBreak
from reportlab.lib import colors as rl_colors
from reportlab.lib.utils import ImageReader

logger = logging.getLogger(__name__)


# ── Register DejaVu Sans for PDF Unicode rendering (₹ € £ ¥ etc.) ───────────
_PDF_BODY_FONT = "Helvetica"
_PDF_BOLD_FONT = "Helvetica-Bold"
try:
    import matplotlib.font_manager as _fm
    from reportlab.pdfbase import pdfmetrics as _pdfmetrics
    from reportlab.pdfbase.ttfonts import TTFont as _TTFont

    _djv_regular = _fm.findfont(_fm.FontProperties(family="DejaVu Sans", style="normal", weight="regular"))
    _djv_bold    = _fm.findfont(_fm.FontProperties(family="DejaVu Sans", style="normal", weight="bold"))
    _pdfmetrics.registerFont(_TTFont("DejaVu",      _djv_regular))
    _pdfmetrics.registerFont(_TTFont("DejaVu-Bold", _djv_bold))
    _PDF_BODY_FONT = "DejaVu"
    _PDF_BOLD_FONT = "DejaVu-Bold"
    logger.info("DejaVu Sans registered for PDF — Unicode currency symbols will render correctly")
except Exception as _font_err:
    logger.warning(
        "Could not register DejaVu Sans for PDF (%s) — "
        "non-ASCII currency symbols (₹ € £) may render as black squares",
        _font_err,
    )

# ── Static file root (logos served by FastAPI StaticFiles mount at /static) ──
_STATIC_DIR = os.path.join(os.path.dirname(os.path.dirname(os.path.abspath(__file__))), "static")

# ── Branding helpers ─────────────────────────────────────────────────────────

def _hex_to_rgb(hex_color: str) -> RGBColor:
    """Convert a #RRGGBB hex string to RGBColor. Returns default indigo on parse error."""
    try:
        h = hex_color.lstrip("#")
        if len(h) == 6:
            return RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))
    except (ValueError, AttributeError):
        pass
    return RGBColor(0x43, 0x38, 0xCA)  # default indigo


def _delete_shape(slide: Any, shape: Any) -> None:
    """Remove a shape element from a slide entirely."""
    try:
        sp = shape._element
        sp.getparent().remove(sp)
    except Exception as exc:
        logger.debug("Could not delete shape: %s", exc)


def _hex_to_rl(hex_color: str):
    """Convert a #RRGGBB hex string to a ReportLab HexColor."""
    try:
        h = hex_color.lstrip("#")
        if len(h) == 6:
            return rl_colors.HexColor(f"#{h}")
    except (ValueError, AttributeError):
        pass
    return _RL_INDIGO


def _download_image(url: str) -> io.BytesIO | None:
    """Return image bytes as BytesIO from a URL or local static path."""
    if not url:
        return None
    local_subpath: str | None = None
    if url.startswith("/static/"):
        local_subpath = url[len("/static/"):]
    else:
        import re as _re
        from config import settings as _settings
        _backend = _re.escape(_settings.BACKEND_URL.rstrip("/"))
        _m = _re.match(rf"(?:https?://localhost(?::\d+)?|{_backend})/static/(.*)", url)
        if _m:
            local_subpath = _m.group(1)
    if local_subpath is not None:
        file_path = os.path.join(_STATIC_DIR, local_subpath)
        try:
            with open(file_path, "rb") as _fh:
                return io.BytesIO(_fh.read())
        except Exception as exc:
            logger.warning("Could not read local logo file %s: %s", file_path, exc)
            return None
    try:
        import httpx as _httpx
        response = _httpx.get(url, timeout=10, follow_redirects=True)
        if response.status_code == 200:
            return io.BytesIO(response.content)
        logger.warning("Logo download returned status %s for URL: %s", response.status_code, url)
    except Exception as exc:
        logger.warning("Logo download failed (%s): %s", url, exc)
    return None


# ── Currency symbol lookup ───────────────────────────────────────────────────
_CURRENCY_SYMBOLS: dict[str, str] = {
    "USD": "$",   "EUR": "€",   "GBP": "£",   "INR": "₹",
    "AUD": "A$",  "CAD": "C$",  "JPY": "¥",   "CNY": "¥",
    "BRL": "R$",  "MXN": "Mex$","SGD": "S$",  "HKD": "HK$",
    "CHF": "CHF ","SEK": "kr",  "NOK": "kr",  "DKK": "kr",
    "ZAR": "R",   "AED": "AED ","SAR": "SAR ","MYR": "RM",
}


def _currency_symbol(data: Dict[str, Any]) -> str:
    code = data.get("meta_ads", {}).get("currency", "USD") or "USD"
    return _CURRENCY_SYMBOLS.get(code.upper(), code + " ")


def _section_enabled(enabled_sections: dict | None, key: str) -> bool:
    if enabled_sections is None:
        return True
    return bool(enabled_sections.get(key, True))


def _to_lines(value: Any) -> list[str]:
    if isinstance(value, list):
        lines: list[str] = []
        for item in value:
            if isinstance(item, str):
                lines.extend(item.splitlines())
            else:
                lines.append(str(item))
        return [l for l in lines if l.strip()]
    if isinstance(value, str):
        return [l for l in value.splitlines() if l.strip()]
    return []


# ── Colour constants ────────────────────────────────────────────────────────
# NOTE: emerald upgraded from #059669 to #047857 (emerald-700, WCAG AA 5.1:1)
# to pass contrast ratio checks for small trend labels on white backgrounds.
_INDIGO        = RGBColor(0x43, 0x38, 0xCA)
_WHITE         = RGBColor(0xFF, 0xFF, 0xFF)
_SLATE_900     = RGBColor(0x0F, 0x17, 0x2A)
_SLATE_700     = RGBColor(0x33, 0x41, 0x55)
_SLATE_500     = RGBColor(0x64, 0x74, 0x8B)
_SLATE_400     = RGBColor(0x94, 0xA3, 0xB8)
_SLATE_200     = RGBColor(0xE2, 0xE8, 0xF0)
_SLATE_100     = RGBColor(0xF1, 0xF5, 0xF9)
_EMERALD       = RGBColor(0x04, 0x78, 0x57)   # emerald-700, WCAG AA
_AMBER         = RGBColor(0xD9, 0x77, 0x06)
_ROSE          = RGBColor(0xE1, 0x1D, 0x48)

_RL_INDIGO     = rl_colors.HexColor("#4338CA")
_RL_SLATE_900  = rl_colors.HexColor("#0F172A")
_RL_SLATE_700  = rl_colors.HexColor("#334155")
_RL_SLATE_500  = rl_colors.HexColor("#64748B")
_RL_SLATE_400  = rl_colors.HexColor("#94A3B8")
_RL_SLATE_200  = rl_colors.HexColor("#E2E8F0")
_RL_SLATE_100  = rl_colors.HexColor("#F1F5F9")
_RL_SLATE_50   = rl_colors.HexColor("#F8FAFC")
_RL_EMERALD    = rl_colors.HexColor("#047857")   # emerald-700, WCAG AA
_RL_AMBER      = rl_colors.HexColor("#D97706")
_RL_ROSE       = rl_colors.HexColor("#E11D48")


# ═══════════════════════════════════════════════════════════════════════════════
# ── Template-based PowerPoint report generator ─────────────────────────────
# ═══════════════════════════════════════════════════════════════════════════════

from services.slide_selector import (
    select_slides, select_kpis, get_slides_to_delete,
    SLIDE_INDEX, TOTAL_TEMPLATE_SLIDES,
)
from services.text_formatter import parse_structured_text

_TEMPLATES_DIR = os.path.join(os.path.dirname(os.path.dirname(os.path.abspath(__file__))),
                              "templates", "pptx")

VISUAL_TEMPLATES = {
    "modern_clean":    os.path.join(_TEMPLATES_DIR, "modern_clean.pptx"),
    "dark_executive":  os.path.join(_TEMPLATES_DIR, "dark_executive.pptx"),
    "colorful_agency": os.path.join(_TEMPLATES_DIR, "colorful_agency.pptx"),
    "bold_geometric":  os.path.join(_TEMPLATES_DIR, "bold_geometric.pptx"),
    "minimal_elegant": os.path.join(_TEMPLATES_DIR, "minimal_elegant.pptx"),
    "gradient_modern": os.path.join(_TEMPLATES_DIR, "gradient_modern.pptx"),
}

# Legacy 9-slide mapping (for backwards compatibility with existing templates)
SLIDE_MAP = {
    "cover": 0,
    "executive_summary": 1,
    "kpi_scorecard": 2,
    "website_traffic": 3,
    "meta_ads": 4,
    "key_wins": 5,
    "concerns": 6,
    "next_steps": 7,
    "custom_section": 8,
}

# New 19-slide mapping from slide_selector
SLIDE_MAP_V2 = SLIDE_INDEX


def _format_period(start_str: str, end_str: str) -> str:
    """
    Convert ISO period dates to a human-readable string for the cover slide
    and `{{report_period}}` placeholder.

    Examples
    --------
    ``_format_period("2026-03-01", "2026-03-31")`` -> ``"March 1 — 31, 2026"``
    ``_format_period("2026-01-15", "2026-03-15")`` -> ``"January 15 — March 15, 2026"``
    ``_format_period("2025-01-01", "2026-03-30")`` -> ``"January 2025 — March 2026"``

    On any parse error the raw "start to end" string is returned as a
    graceful fallback.
    """
    if not start_str or not end_str:
        return f"{start_str or ''} to {end_str or ''}".strip()
    try:
        start = datetime.strptime(start_str, "%Y-%m-%d")
        end = datetime.strptime(end_str, "%Y-%m-%d")
    except (ValueError, TypeError):
        return f"{start_str} to {end_str}"

    if start.year == end.year and start.month == end.month:
        # Same month — "March 1 — 31, 2026"
        return f"{start.strftime('%B')} {start.day} \u2014 {end.day}, {end.year}"
    if start.year == end.year:
        # Same year — "January 15 — March 15, 2026"
        return (
            f"{start.strftime('%B')} {start.day} \u2014 "
            f"{end.strftime('%B')} {end.day}, {end.year}"
        )
    # Multi-year — "January 2025 — March 2026"
    return f"{start.strftime('%B %Y')} \u2014 {end.strftime('%B %Y')}"


def _fmt_num(val: Any) -> str:
    """Format a number with comma separators (full precision)."""
    try:
        return f"{int(val):,}"
    except (ValueError, TypeError):
        return str(val) if val else "0"


def _fmt_compact(val: Any) -> str:
    """
    Format a number in compact K/M/B notation for KPI card big-number slots.
    Keeps full precision with commas for values under 1,000.

    Examples:
        847       -> "847"
        1234      -> "1.2K"
        12_456    -> "12.5K"
        1_234_567 -> "1.2M"
        2_300_000_000 -> "2.3B"

    Do NOT use this for detailed tables or narrative text — only for the
    large hero numbers on scorecard slides (see Phase 5 gap analysis item 4).
    """
    try:
        num = float(val)
    except (ValueError, TypeError):
        return str(val) if val else "0"
    abs_num = abs(num)
    if abs_num < 1_000:
        # Keep integer precision for small numbers
        return f"{int(num):,}" if num == int(num) else f"{num:,.1f}"
    if abs_num < 1_000_000:
        return f"{num / 1_000:.1f}K"
    if abs_num < 1_000_000_000:
        return f"{num / 1_000_000:.1f}M"
    return f"{num / 1_000_000_000:.1f}B"


def _fmt_change(val: float | None) -> str:
    """
    Format a percentage change with a leading direction glyph:
        positive (> +1%)   -> "▲ +12.3%"
        negative (< −1%)   -> "▼ -5.1%"
        neutral (|val|<1%) -> "▬ +0.3%"
        None               -> "N/A"

    The ±1% neutral band prevents noise-level changes (±0.x%) from lighting up
    as real trends. Color is applied later by ``_colorize_kpi_changes``, which
    also handles inverse metrics (CPC, CPA, bounce rate, etc.).
    """
    if val is None:
        return "N/A"
    # Neutral / dead-zone band: −1% < val < +1%
    if abs(val) < 1.0:
        sign = "+" if val >= 0 else ""
        return f"\u25AC {sign}{val:.1f}%"
    if val >= 0:
        return f"\u25B2 +{val:.1f}%"
    return f"\u25BC {val:.1f}%"


def _narrative_to_text(content: Any) -> str:
    """Convert narrative content (str or list) to plain text."""
    if isinstance(content, list):
        return "\n\n".join(str(item) for item in content if str(item).strip())
    return str(content) if content else ""


def _list_to_text(content: Any) -> str:
    """Convert list content to newline-separated text."""
    if isinstance(content, list):
        return "\n".join(str(item) for item in content if str(item).strip())
    if isinstance(content, str):
        return content
    return ""


def _build_replacements(
    data: Dict[str, Any],
    narrative: Dict[str, str],
    client_info: Dict[str, Any],
    branding: dict | None,
    custom_section: dict | None,
    template: str = "full",
    language: str = "en",
) -> dict[str, str]:
    """Build the {{placeholder}} → value replacement dictionary."""
    br = branding or {}
    cur_sym = _currency_symbol(data)

    # Smart KPI selection — pick the 6 best KPIs from available data
    smart_kpis = select_kpis(data, currency_symbol=cur_sym)
    kpis = [
        (k["label"], k["value"], k["change"] or "")
        for k in smart_kpis
    ]
    # Pad to 6 if fewer available
    while len(kpis) < 6:
        kpis.append(("", "", ""))

    agency_name = (br.get("agency_name") or "").strip() or (client_info.get("agency_name") or "").strip() or "Your Agency"
    report_date = datetime.now().strftime("%B %d, %Y")
    p_start = data.get("period_start", "")
    p_end   = data.get("period_end", "")
    report_period_label = _format_period(p_start, p_end)

    # Template-specific label — avoids "Performance Report" appearing twice on cover
    _template_labels = {
        "summary": t(language, "summary_report"),
        "brief":   t(language, "executive_brief"),
    }
    report_type_label = _template_labels.get(template, t(language, "monthly_performance_report"))

    replacements = {
        "{{client_name}}":           client_info.get("name", "Client"),
        "{{report_period}}":         report_period_label,
        "{{agency_name}}":           agency_name,
        # When no agency email is configured, substitute a "Confidential • Page N"
        # token so the next_steps footer renders properly.  _renumber_slide_footers
        # will replace "Page 99" with the actual sequential page number.
        "{{agency_email}}":          br.get("agency_email") or f"{t(language, 'confidential')}  \u2022  {t(language, 'page')} 99",
        "{{report_type}}":           report_type_label,
        "{{report_date}}":           report_date,
        "{{executive_summary}}":     _narrative_to_text(narrative.get("executive_summary", "")),
        "{{website_narrative}}":     _narrative_to_text(narrative.get("website_performance", "")),
        "{{engagement_narrative}}":  _narrative_to_text(narrative.get("engagement_analysis", "")),
        "{{ads_narrative}}":         _narrative_to_text(narrative.get("paid_advertising", "")),
        "{{key_wins}}":              _list_to_text(narrative.get("key_wins", "")),
        "{{concerns}}":              _list_to_text(narrative.get("concerns", "")),
        "{{next_steps}}":            _list_to_text(narrative.get("next_steps", "")),
        "{{custom_section_title}}":  (custom_section or {}).get("title", "Additional Notes"),
        "{{custom_section_text}}":   (custom_section or {}).get("text", ""),
        "{{footer_text}}":           (
            f"{agency_name} \u2022 Powered by GoReportPilot \u2022 {report_date}"
            if br.get("powered_by_badge", True)
            else f"{agency_name} \u2022 {t(language, 'confidential')} \u2022 {report_date}"
        ),
        "{{agency_logo}}":           "",  # Cleared — actual image embedded by _embed_logos()
        "{{client_logo}}":           "",  # Cleared — actual image embedded by _embed_logos()
    }

    # KPI placeholders — translate labels to target language
    for i, (label, value, change) in enumerate(kpis):
        replacements[f"{{{{kpi_{i}_label}}}}"]  = translate_kpi_label(label, language)
        replacements[f"{{{{kpi_{i}_value}}}}"]  = value
        replacements[f"{{{{kpi_{i}_change}}}}"] = change

    # CSV slide placeholders are intentionally NOT included here.
    # Each CSV source gets its own duplicate of the csv_data template slide,
    # populated individually by _populate_csv_slide() after slide deletion.
    # Any leftover {{csv_*}} tokens on non-CSV slides are cleared by the
    # leftover-placeholder cleanup pass in generate_pptx_report().

    return replacements


def _embed_kpi_sparklines(
    prs: Any,
    charts: Dict[str, str],
    selected_kpi_labels: list[str],
) -> None:
    """
    Embed small sparkline PNGs on the KPI scorecard slide.

    Must run **before** ``_replace_placeholders_in_slide`` so the
    ``{{kpi_i_value}}`` / ``{{kpi_i_change}}`` tokens are still present
    for shape identification. Handles both separate-shape and combined-
    shape card layouts.

    Picture height is fixed at 0.25" so sparklines read clearly on every
    template. Earlier gap-proportional sizing collapsed them to ~0.08"
    on tight cards. Separate-shape cards with a gap tighter than
    ``_MIN_GAP_EMU`` are skipped entirely rather than overlapping the
    change indicator.

    Labels with no matching sparkline PNG in ``charts`` are silently
    skipped.
    """
    if not charts or not selected_kpi_labels:
        logger.info(
            "KPI sparklines: skipping (charts=%s, labels=%s)",
            bool(charts), len(selected_kpi_labels or []),
        )
        return

    # Collect all available sparkline keys for diagnostics.
    _spark_keys = [k for k in charts if k.startswith("sparkline__")]
    if not _spark_keys:
        logger.info("KPI sparklines: no sparkline files were generated")
        return
    logger.info(
        "KPI sparklines: %d files available, KPI labels: %s",
        len(_spark_keys), selected_kpi_labels,
    )

    _PIC_HEIGHT  = Inches(0.25)
    _MIN_GAP_EMU = Inches(0.05)
    _PIC_GAP     = Inches(0.02)
    _MIN_PIC_W   = Inches(1.2)
    _slide_w     = prs.slide_width

    _processed_any = False

    for slide_idx, slide in enumerate(prs.slides):
        # Only process slides that contain KPI card tokens.
        _slide_text = ""
        for shape in slide.shapes:
            if shape.has_text_frame:
                _slide_text += shape.text_frame.text
        if "{{kpi_0_value}}" not in _slide_text:
            continue

        _processed_any = True
        logger.info(
            "KPI sparklines: processing slide %d for %d KPIs",
            slide_idx, len(selected_kpi_labels),
        )

        # Build per-slot shape lookup by token.
        value_shapes: dict[int, Any] = {}
        change_shapes: dict[int, Any] = {}
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            text = shape.text_frame.text
            for i in range(6):
                if f"{{{{kpi_{i}_value}}}}" in text:
                    value_shapes[i] = shape
                if f"{{{{kpi_{i}_change}}}}" in text:
                    change_shapes[i] = shape

        _embedded = 0
        _skipped: list[str] = []
        for i in range(min(6, len(selected_kpi_labels))):
            label = selected_kpi_labels[i]
            spark_path = charts.get(f"sparkline__{label}")
            if not spark_path or not os.path.exists(spark_path):
                _skipped.append(f"{i}:{label}=no-file")
                continue

            v_shape = value_shapes.get(i)
            c_shape = change_shapes.get(i)
            if v_shape is None:
                _skipped.append(f"{i}:{label}=no-value-shape")
                continue

            try:
                v_top    = int(v_shape.top or 0)
                v_left   = int(v_shape.left or 0)
                v_width  = int(v_shape.width or 0)
                v_height = int(v_shape.height or 0)
                v_bottom = v_top + v_height

                _same_shape = c_shape is not None and c_shape is v_shape

                # Separate-shape cards with a too-tight gap are skipped
                # outright — the sparkline would overlap the change glyph.
                if c_shape is not None and not _same_shape:
                    gap = int(c_shape.top or 0) - v_bottom
                    if gap < _MIN_GAP_EMU:
                        _skipped.append(
                            f"{i}:{label}=tight-gap({gap / 914_400:.3f}in)"
                        )
                        continue

                pic_t = v_bottom + _PIC_GAP
                pic_h = _PIC_HEIGHT
                pic_w = max(int(v_width * 0.9), _MIN_PIC_W)
                pic_l = v_left + (v_width - pic_w) // 2
                # Clamp horizontally so the picture never lands off-slide.
                if pic_l < 0:
                    pic_l = 0
                if pic_l + pic_w > _slide_w:
                    pic_l = max(0, _slide_w - pic_w)

                slide.shapes.add_picture(
                    spark_path, pic_l, pic_t, width=pic_w, height=pic_h,
                )
                _embedded += 1
                logger.info(
                    "KPI sparkline embedded: slot=%d label=%s mode=%s",
                    i, label, "combined" if _same_shape else "separate",
                )
            except Exception as exc:
                logger.warning(
                    "Sparkline embed failed for KPI %d (%s): %s",
                    i, label, exc,
                )
                _skipped.append(f"{i}:{label}=error")

        if _skipped:
            logger.info("KPI sparklines skipped: %s", ", ".join(_skipped))
        logger.info(
            "KPI sparklines: slide %d → embedded=%d skipped=%d",
            slide_idx, _embedded, len(_skipped),
        )

    if not _processed_any:
        logger.info(
            "KPI sparklines: no slide contained {{kpi_0_value}} (template "
            "may not declare scorecard tokens)"
        )


def _replace_placeholders_in_slide(slide: Any, replacements: dict[str, str]) -> None:
    """Find and replace all {{placeholder}} text in a slide, preserving formatting."""
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        for paragraph in shape.text_frame.paragraphs:
            # Combine all runs to find placeholders that may span runs
            full_text = "".join(run.text for run in paragraph.runs)
            matched = False
            for placeholder, value in replacements.items():
                if placeholder in full_text:
                    full_text = full_text.replace(placeholder, str(value))
                    matched = True
            if matched and paragraph.runs:
                # Put all text in the first run (preserving its formatting), clear others
                paragraph.runs[0].text = full_text
                for run in paragraph.runs[1:]:
                    run.text = ""


def _renumber_slide_footers(prs: Any, language: str = "en") -> None:
    """
    After slide deletion the footer 'Page N' values reflect the original template
    numbering. Walk every remaining slide and overwrite 'Page N' with the correct
    sequential number.  Supports translated page labels.
    """
    page_word = t(language, "page")
    # Match both English "Page" and the translated word followed by digits
    _page_re = re.compile(r"(?:Page|" + re.escape(page_word) + r")\s+\d+")
    for idx, slide in enumerate(prs.slides, 1):
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            for para in shape.text_frame.paragraphs:
                full_text = "".join(run.text for run in para.runs)
                if _page_re.search(full_text):
                    for run in para.runs:
                        if _page_re.search(run.text):
                            run.text = _page_re.sub(f"{page_word} {idx}", run.text)


def _remove_static_email_cta(prs: Any) -> None:
    """
    Remove shapes that contain generic email-only call-to-action text such as
    'Questions? Reply to this email' or 'schedule a call'.
    These are static template elements that should only appear in emailed reports.
    """
    _patterns = [
        r"questions\?",
        r"reply to this email",
        r"schedule a call",
        r"have questions\?",
    ]
    combined = re.compile("|".join(_patterns), re.IGNORECASE)
    for slide in prs.slides:
        to_remove = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                text = shape.text_frame.text
                if combined.search(text):
                    to_remove.append(shape)
        for shape in to_remove:
            _delete_shape(slide, shape)


def _replace_charts(prs: Any, charts: Dict[str, str]) -> None:
    """Replace chart placeholder text boxes with actual chart PNG images.

    When a chart file exists: embeds the PNG and clears the placeholder text.
    When no chart is available: deletes the placeholder shape.

    Dual-chart slide handling:
    When a slide has two chart placeholders (e.g. device_breakdown + top_pages)
    and only ONE chart has data, the available chart is automatically expanded
    to fill the combined bounding box of both placeholder areas so no blank
    half-slide is left visible.

    CSV chart placeholders ({{chart_csv_data}} / {{chart_csv_*}}) are
    intentionally excluded — they are embedded per-source in _populate_csv_slide()
    after slide duplication.
    """
    chart_mapping: dict[str, str | None] = {
        # GA4 charts
        "{{chart_sessions}}":           charts.get("sessions"),
        "{{chart_traffic}}":            charts.get("traffic_sources"),
        "{{chart_device_breakdown}}":   charts.get("device_breakdown"),
        "{{chart_top_pages}}":          charts.get("top_pages"),
        "{{chart_new_vs_returning}}":   charts.get("new_vs_returning"),
        "{{chart_top_countries}}":      charts.get("top_countries"),
        "{{chart_bounce_rate}}":        charts.get("bounce_rate_trend"),
        "{{chart_conversion_funnel}}":  charts.get("conversion_funnel"),
        # Meta Ads charts
        "{{chart_spend}}":              charts.get("spend_conversions"),
        "{{chart_campaigns}}":          charts.get("campaigns"),
        "{{chart_demographics}}":       charts.get("audience_demographics"),
        "{{chart_placements}}":         charts.get("placements"),
        # Google Ads charts
        "{{chart_gads_spend}}":         charts.get("gads_spend_conversions"),
        "{{chart_gads_campaigns}}":     charts.get("gads_campaigns"),
        "{{chart_search_terms}}":       charts.get("search_terms_bar"),
        # Search Console charts
        "{{chart_seo_trend}}":          charts.get("seo_clicks_trend"),
        "{{chart_top_queries}}":        charts.get("top_queries"),
    }

    # Known chart pairs that share a slide side-by-side.
    # If one is missing, the other gets centered on the slide.
    _DUAL_PAIRS: dict[str, str] = {
        "{{chart_device_breakdown}}": "{{chart_top_pages}}",
        "{{chart_top_pages}}":        "{{chart_device_breakdown}}",
    }

    # Pie/donut-style charts — keep original width when solo, just center.
    # All other chart types may expand to a moderate max width.
    _PIE_CHART_HINTS = ("device", "breakdown", "pie", "donut", "returning")
    _SLIDE_W = prs.slide_width                         # 13.33" in EMU
    _MAX_BAR_W = int(8.5 * 914400)                     # bar/line charts: 8.5" max

    logger.debug("_replace_charts: charts available: %s", list(charts.keys()))

    replaced = 0
    for slide in prs.slides:
        # ── Collect all chart placeholder shapes on this slide ────────────────
        shapes_on_slide: list[tuple] = []   # (shape, chart_path, placeholder)
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            text = shape.text_frame.text.strip()
            for ph, cp in chart_mapping.items():
                if ph in text:
                    shapes_on_slide.append((shape, cp, ph))
                    break

        if not shapes_on_slide:
            continue

        # ── Dual-chart centering pre-pass ─────────────────────────────────────
        # Build a lookup of placeholder → (shape, chart_path) for this slide
        ph_map: dict[str, tuple] = {ph: (sh, cp) for sh, cp, ph in shapes_on_slide}

        # When one chart of a known pair is absent, compute a *centered*
        # geometry for the surviving chart rather than stretching it to
        # the full combined width (which distorts pie charts).
        solo_geom: dict[str, tuple[int, int, int, int]] = {}
        for ph, (sh, cp) in ph_map.items():
            partner_ph = _DUAL_PAIRS.get(ph)
            if partner_ph and partner_ph in ph_map:
                _has      = bool(cp and os.path.exists(cp))
                _psh, _pcp = ph_map[partner_ph]
                _p_has    = bool(_pcp and os.path.exists(_pcp))
                if _has and not _p_has:
                    # Determine chart category from placeholder name
                    ph_lower = ph.lower()
                    is_pie = any(hint in ph_lower for hint in _PIE_CHART_HINTS)

                    combined_top    = min(sh.top,  _psh.top)
                    combined_bottom = max(sh.top + sh.height, _psh.top + _psh.height)
                    combined_h      = combined_bottom - combined_top

                    if is_pie:
                        # Pie/donut: keep original single-slot width, center it
                        w = sh.width
                    else:
                        # Bar/line: expand up to _MAX_BAR_W, but not full slide
                        w = min(sh.width + _psh.width, _MAX_BAR_W)

                    l = (_SLIDE_W - w) // 2   # horizontally centered
                    solo_geom[ph] = (l, combined_top, w, combined_h)
                    logger.info(
                        "Solo chart centering: %s (%s) → %.2f\"×%.2f\" centered",
                        ph, "pie" if is_pie else "bar",
                        w / 914400, combined_h / 914400,
                    )

        # ── Embed or delete each chart placeholder ────────────────────────────
        for shape, chart_path, placeholder in shapes_on_slide:
            if chart_path and os.path.exists(chart_path):
                if placeholder in solo_geom:
                    s_l, s_t, s_w, s_h = solo_geom[placeholder]
                    slide.shapes.add_picture(chart_path, s_l, s_t, s_w, s_h)
                else:
                    slide.shapes.add_picture(
                        chart_path, shape.left, shape.top, shape.width, shape.height,
                    )
                for para in shape.text_frame.paragraphs:
                    for run in para.runs:
                        run.text = ""
                replaced += 1
                logger.debug("Replaced %s with image", placeholder)
            else:
                logger.info("No chart for %s — removing placeholder shape", placeholder)
                _delete_shape(slide, shape)

    logger.info("_replace_charts: %d chart(s) embedded", replaced)


# KPI label keywords where LOWER is better. A positive change on these
# metrics is a BAD outcome (rising costs, worsening ranking, rising
# bounce, higher frequency) and must render in the negative color.
_INVERSE_METRIC_KEYWORDS: tuple[str, ...] = (
    "BOUNCE",
    "COST",             # COST / CONV., CPC, CPA
    "CPC",
    "CPA",
    "CPM",
    "FREQ",             # FREQUENCY
    "POSITION",         # SEO avg position — lower rank is better
    "UNSUBSCRIBE",
    "SPEND",            # AD SPEND / META SPEND / GOOGLE SPEND / SEARCH SPEND
)


def _is_inverse_label(text: str) -> bool:
    """Return True when the given label text describes a 'lower is better' metric."""
    t = (text or "").upper()
    return any(kw in t for kw in _INVERSE_METRIC_KEYWORDS)


def _colorize_kpi_changes(prs: Any, data: Dict[str, Any]) -> None:
    """
    Colorize KPI trend indicators produced by ``_fmt_change``.

    Recognised run formats:
        "▲ +12.3%"  -> positive direction
        "▼ -5.1%"   -> negative direction
        "▬ +0.3%"   -> neutral (inside the ±1% dead zone)
        "+12.3%"    -> legacy positive (pre-glyph callers)
        "-5.1%"     -> legacy negative
        "N/A"       -> missing data

    Color mapping is INVERTED for cost / bounce / position / frequency labels
    (the ``_INVERSE_METRIC_KEYWORDS`` set): on those cards, ▲ is red and
    ▼ is green because lower values are better.

    The inverse determination is best-effort: we look first at the shape's
    own text frame (label + value + change are often in the same box), then
    fall back to the spatially-nearest label shape on the same slide.
    """
    for slide in prs.slides:
        # Pre-scan all shapes on the slide so we can find the nearest label
        # for any change run that lives in its own shape.
        shapes_info: list[tuple[Any, str]] = []
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            shapes_info.append((shape, shape.text_frame.text or ""))

        for shape, shape_full_text in shapes_info:
            shape_has_inverse = _is_inverse_label(shape_full_text)
            for para in shape.text_frame.paragraphs:
                for run in para.runs:
                    text = run.text.strip()
                    if not text:
                        continue

                    # Classify direction by leading glyph or sign.
                    is_up = text.startswith("\u25B2") or (
                        text.startswith("+") and "%" in text
                    )
                    is_down = text.startswith("\u25BC") or (
                        text.startswith("-") and "%" in text
                    )
                    is_flat = text.startswith("\u25AC")

                    if text == "N/A":
                        run.font.color.rgb = _SLATE_400
                        continue

                    if not (is_up or is_down or is_flat):
                        continue

                    # Determine inverse-metric status for this run.
                    is_inverse = shape_has_inverse
                    if not is_inverse:
                        try:
                            s_top = int(shape.top or 0)
                            s_left = int(shape.left or 0)

                            def _dist(other: tuple[Any, str]) -> int:
                                osh, _txt = other
                                if osh is shape:
                                    return 10 ** 12
                                return abs(int(osh.top or 0) - s_top) + abs(
                                    int(osh.left or 0) - s_left
                                )

                            nearest = min(shapes_info, key=_dist, default=None)
                            if nearest and _is_inverse_label(nearest[1]):
                                is_inverse = True
                        except Exception:
                            pass

                    # Apply color.
                    if is_flat:
                        run.font.color.rgb = _SLATE_500
                    elif is_up:
                        run.font.color.rgb = _ROSE if is_inverse else _EMERALD
                    else:  # is_down
                        run.font.color.rgb = _EMERALD if is_inverse else _ROSE
                    run.font.bold = True


def _get_slides_to_delete(
    enabled_sections: dict | None,
    template: str,
    narrative: Dict[str, str],
    custom_section: dict | None,
) -> set[int]:
    """Determine which slide indices to delete based on config and content."""
    to_delete: set[int] = set()
    es = enabled_sections or {}

    section_slide_map = {
        "executive_summary": 1,
        "kpi_scorecard": 2,
        "website_traffic": 3,
        "meta_ads": 4,
        "key_wins": 5,
        "concerns": 6,
        "next_steps": 7,
        "custom_section": 8,
    }

    for section_key, slide_idx in section_slide_map.items():
        # Delete if explicitly disabled
        if section_key in es and not es[section_key]:
            to_delete.add(slide_idx)
            continue

        # Delete custom section if no content
        if section_key == "custom_section":
            cs = custom_section or {}
            if not cs.get("title") or not cs.get("text", "").strip():
                to_delete.add(slide_idx)
        # Delete narrative sections if empty
        elif section_key in ("key_wins", "concerns", "next_steps",
                             "executive_summary"):
            content = narrative.get(section_key, "")
            lines = _to_lines(content)
            if not lines:
                to_delete.add(slide_idx)

    # Template-level filtering
    if template == "summary":
        # Summary keeps: cover, exec summary, KPI, key_wins/next_steps
        to_delete.update({3, 4, 6, 8})  # website, meta_ads, concerns, custom
    elif template == "brief":
        # Brief keeps: cover, exec summary, KPI only
        to_delete.update({3, 4, 5, 6, 7, 8})

    return to_delete


def _fit_image_to_box(
    img_width: int,
    img_height: int,
    max_width: int,
    max_height: int,
) -> tuple[int, int]:
    """
    Scale ``img_width × img_height`` to fit inside ``max_width × max_height``
    while preserving aspect ratio. Returns ``(width, height)`` in whatever
    units the caller passed in (EMU in practice).

    Also protects against degenerate (zero or missing) source dimensions by
    returning the max box as a conservative fallback.
    """
    if img_width <= 0 or img_height <= 0:
        return max_width, max_height
    ratio = min(max_width / img_width, max_height / img_height)
    return int(img_width * ratio), int(img_height * ratio)


def _measure_image(image_stream: "io.BytesIO") -> tuple[int, int]:
    """
    Return the native ``(width_px, height_px)`` of an in-memory image.
    Falls back to ``(0, 0)`` on any PIL / decoding error — callers treat
    that as "unknown size, use the max box".
    """
    try:
        from PIL import Image  # noqa: PLC0415
        image_stream.seek(0)
        with Image.open(image_stream) as img:
            w, h = img.size
        image_stream.seek(0)
        return int(w), int(h)
    except Exception as exc:
        logger.debug("Could not measure logo image: %s", exc)
        return 0, 0


# Phase-3 logo placement helpers ---------------------------------------------
_LOGO_MARGIN = Inches(0.3)


def _logo_max_box(size: str, *, kind: str) -> tuple[int, int]:
    """Return (max_w, max_h) EMU for a given size preset and logo kind."""
    size = (size or "default").lower()
    # Agency logos sit in the header band — relatively narrow.
    # Client logos can be larger because they're prominent on the cover.
    if kind == "agency":
        box = {
            "small":   (Inches(1.5), Inches(0.5)),
            "medium":  (Inches(2.5), Inches(0.8)),   # default
            "large":   (Inches(3.5), Inches(1.2)),
            "default": (Inches(2.5), Inches(0.8)),
        }
    else:
        box = {
            "small":   (Inches(1.5), Inches(1.0)),
            "medium":  (Inches(3.0), Inches(2.0)),   # default
            "large":   (Inches(4.5), Inches(3.0)),
            "default": (Inches(3.0), Inches(2.0)),
        }
    return box.get(size, box["default"])


def _clamp_logo_for_position(position: str, max_w: int, max_h: int, *, kind: str) -> tuple[int, int]:
    """
    Shrink the max bounding box for logos placed in footer / bottom slots.

    Why: a "medium" client logo defaults to 2" tall. Dropping a 2"-tall
    logo at the bottom of a 7.5" slide means its top edge sits at 5.2" —
    looking visually mid-slide, not "footer". Clamp to a small footer
    band so the logo reads as a footer element no matter what size the
    user selected.
    """
    pos = (position or "").lower()
    if pos.startswith("footer") or pos.startswith("bottom"):
        if kind == "agency":
            return min(max_w, Inches(2.5)), min(max_h, Inches(0.6))
        return min(max_w, Inches(3.0)), min(max_h, Inches(0.8))
    return max_w, max_h


def _logo_corner_xy(
    *, position: str, slide_w: int, slide_h: int, logo_w: int, logo_h: int,
) -> tuple[int, int]:
    """Map a named position to explicit (left, top) EMU on the cover slide."""
    m = _LOGO_MARGIN
    if position == "top-left":
        result = (m, m)
    elif position in ("top-right", "header"):
        result = (slide_w - logo_w - m, m)
    elif position == "top-center":
        result = ((slide_w - logo_w) // 2, m)
    elif position in ("footer", "footer-left", "bottom-left"):
        result = (m, slide_h - logo_h - m)
    elif position in ("footer-right", "bottom-right"):
        result = (slide_w - logo_w - m, slide_h - logo_h - m)
    elif position in ("footer-center", "bottom-center"):
        result = ((slide_w - logo_w) // 2, slide_h - logo_h - m)
    elif position == "center":
        result = ((slide_w - logo_w) // 2, (slide_h - logo_h) // 2)
    else:
        # Unknown value — fall back to top-right so logos remain visible,
        # but log the surprise so we notice unknown values in production.
        logger.warning("_logo_corner_xy: unknown position %r → top-right fallback", position)
        result = (slide_w - logo_w - m, m)

    # Phase-3 fix v4 — step 4 of 4: log the exact mapping we used. If this
    # disagrees with the expected branch, the position string arriving here
    # doesn't match what the user picked.
    logger.info(
        "_logo_corner_xy position=%r logo_w=%d logo_h=%d → left=%d top=%d",
        position, int(logo_w), int(logo_h), int(result[0]), int(result[1]),
    )
    return result


def _embed_logos(prs: Any, branding: dict | None) -> None:
    """Embed agency and client logos on the cover slide.

    When an image URL is provided: download and embed it at the placeholder
    position, constrained to a max bounding box so wide logos never overflow.
    When no image is provided: DELETE the placeholder shape so no empty box
    appears.

    Logos render directly on the slide with no background pad — agencies
    are expected to upload logos that work against the cover-header color.
    """
    cover = prs.slides[0]
    sw = prs.slide_width

    # Max bounding boxes (EMU). Sized for a 13.33" × 7.5" widescreen slide.
    _AGENCY_MAX_W = Inches(2.5)
    _AGENCY_MAX_H = Inches(0.8)
    _CLIENT_MAX_W = Inches(3.0)
    _CLIENT_MAX_H = Inches(2.0)
    _RIGHT_MARGIN = Inches(0.3)

    # Locate placeholder shapes by their text content
    agency_logo_shape = None
    client_logo_shape = None
    for shape in list(cover.shapes):
        if not shape.has_text_frame:
            continue
        text = shape.text_frame.text.strip()
        if "agency_logo" in text.lower():
            agency_logo_shape = shape
        elif "client_logo" in text.lower():
            client_logo_shape = shape

    br = branding or {}

    # Optional Phase-3 per-client overrides. 'default' preserves the
    # original placeholder-based behaviour.
    agency_pos  = (br.get("agency_logo_position") or "default").strip().lower()
    agency_size = (br.get("agency_logo_size")     or "default").strip().lower()
    client_pos  = (br.get("client_logo_position") or "default").strip().lower()
    client_size = (br.get("client_logo_size")     or "default").strip().lower()

    # Phase-3 fix v3 — log resolved placement so we can diagnose
    # any flow-through issue from the client row / API payload.
    logger.info(
        "_embed_logos: agency_pos=%r size=%r  client_pos=%r size=%r",
        agency_pos, agency_size, client_pos, client_size,
    )

    # ── Agency logo ──────────────────────────────────────────────────────
    agency_img = _download_image(br.get("agency_logo_url", ""))
    if agency_img:
        try:
            # Measure native pixel dimensions so we can preserve aspect ratio.
            img_w_px, img_h_px = _measure_image(agency_img)

            # Size selector overrides default max box.
            max_w, max_h = _logo_max_box(agency_size, kind="agency")
            # Clamp for footer slots so tall logos don't extend into the body.
            max_w, max_h = _clamp_logo_for_position(agency_pos, max_w, max_h, kind="agency")

            if agency_pos != "default":
                # User-picked placement — derive corner coordinates explicitly.
                fit_w, fit_h = _fit_image_to_box(img_w_px, img_h_px, max_w, max_h)
                pic_left, pic_top = _logo_corner_xy(
                    position=agency_pos, slide_w=sw, slide_h=prs.slide_height,
                    logo_w=fit_w, logo_h=fit_h,
                )
                logger.info(
                    "agency logo explicit placement: pos=%s → left=%d top=%d (EMU)",
                    agency_pos, int(pic_left), int(pic_top),
                )
                # We're placing by explicit coords — clear the placeholder if present.
                if agency_logo_shape:
                    _delete_shape(cover, agency_logo_shape)
                    agency_logo_shape = None
            else:
                # Legacy placeholder-based placement.
                if agency_logo_shape:
                    max_w = min(int(agency_logo_shape.width or max_w), max_w)
                    max_h = min(int(agency_logo_shape.height or max_h), max_h)
                    base_top = int(agency_logo_shape.top)
                else:
                    base_top = Inches(0.3)

                fit_w, fit_h = _fit_image_to_box(img_w_px, img_h_px, max_w, max_h)
                pic_left = sw - fit_w - _RIGHT_MARGIN
                pic_top = base_top

            # Logo renders directly on the indigo header with no background
            # pad — agencies are expected to upload a logo that works on
            # dark backgrounds (most already do).
            agency_img.seek(0)
            cover.shapes.add_picture(
                agency_img, pic_left, pic_top, width=fit_w, height=fit_h,
            )
            if agency_logo_shape:
                _delete_shape(cover, agency_logo_shape)
        except Exception as e:
            logger.debug("Could not embed agency logo: %s", e)
    else:
        # No logo provided — delete the empty placeholder so it doesn't clutter the slide
        if agency_logo_shape:
            _delete_shape(cover, agency_logo_shape)

    # ── Client logo ──────────────────────────────────────────────────────
    client_img = _download_image(br.get("client_logo_url", ""))
    if client_img:
        try:
            img_w_px, img_h_px = _measure_image(client_img)

            max_w, max_h = _logo_max_box(client_size, kind="client")
            max_w, max_h = _clamp_logo_for_position(client_pos, max_w, max_h, kind="client")

            if client_pos != "default":
                fit_w, fit_h = _fit_image_to_box(img_w_px, img_h_px, max_w, max_h)
                pic_left, pic_top = _logo_corner_xy(
                    position=client_pos, slide_w=sw, slide_h=prs.slide_height,
                    logo_w=fit_w, logo_h=fit_h,
                )
                logger.info(
                    "client logo explicit placement: pos=%s → left=%d top=%d (EMU)",
                    client_pos, int(pic_left), int(pic_top),
                )
                if client_logo_shape:
                    _delete_shape(cover, client_logo_shape)
                    client_logo_shape = None
            else:
                if client_logo_shape:
                    max_w = min(int(client_logo_shape.width or max_w), max_w)
                    max_h = min(int(client_logo_shape.height or max_h), max_h)
                    base_top = int(client_logo_shape.top)
                    center_x = int(client_logo_shape.left) + max_w // 2
                else:
                    base_top = Inches(5.5)
                    center_x = sw // 2

                fit_w, fit_h = _fit_image_to_box(img_w_px, img_h_px, max_w, max_h)

                # Horizontally centered around the placeholder (or slide center).
                pic_left = center_x - fit_w // 2
                pic_top = base_top

            # Client logos typically sit on the light body area; skip the pad
            # there. They only need contrast help when the slide background
            # is dark, which the cover-header agency path already handles.
            client_img.seek(0)
            cover.shapes.add_picture(
                client_img, pic_left, pic_top, width=fit_w, height=fit_h,
            )
            if client_logo_shape:
                _delete_shape(cover, client_logo_shape)
        except Exception as e:
            logger.debug("Could not embed client logo: %s", e)
    else:
        if client_logo_shape:
            _delete_shape(cover, client_logo_shape)


def _embed_custom_section_image(prs: Any, custom_section: dict | None, slide_index: int) -> None:
    """
    Embed an uploaded image on the custom section slide.
    Adjusts text area width if image is present.
    """
    if not custom_section:
        return
    image_url = custom_section.get("image_url", "")
    if not image_url:
        return

    # Get the custom section slide
    if slide_index >= len(prs.slides):
        return
    slide = prs.slides[slide_index]

    img_data = _download_image(image_url)
    if not img_data:
        return

    try:
        img_data.seek(0)
        # Place image on right side (35% of slide width)
        slide_w = prs.slide_width
        slide_h = prs.slide_height
        img_w = int(slide_w * 0.35)
        img_x = slide_w - img_w - Inches(0.3)
        img_y = Inches(1.5)
        img_h = int(slide_h * 0.55)
        slide.shapes.add_picture(img_data, img_x, img_y, width=img_w, height=img_h)

        # Narrow the text area on this slide to avoid overlap
        for shape in slide.shapes:
            if shape.has_text_frame:
                text = shape.text_frame.text
                if "custom_section_text" in text or (shape.left < Inches(1) and shape.top > Inches(1.2)):
                    # Limit width to 60% of slide
                    shape.width = int(slide_w * 0.60)
                    break
        logger.info("Custom section image embedded successfully")
    except Exception as e:
        logger.warning("Could not embed custom section image: %s", e)


def _populate_text_frame_formatted(
    text_frame: Any,
    content: str,
    font_size: Any = None,
    font_color: Any = None,
    bold: bool = False,
) -> None:
    """
    Replace a text frame's content with properly-formatted multi-paragraph text.
    Splits content on double newlines into separate paragraphs.
    Preserves font name and size from the template's first run where possible.
    """
    import re as _re

    if not content or not content.strip():
        text_frame.clear()
        return

    # Capture template formatting from first paragraph / first run
    template_name: str | None = None
    template_size: Any = font_size
    template_color: Any = font_color
    template_bold: bool = bold
    try:
        first_para = text_frame.paragraphs[0]
        if first_para.runs:
            r0 = first_para.runs[0]
            if r0.font.name:
                template_name = r0.font.name
            if r0.font.size:
                template_size = r0.font.size
            if r0.font.color and r0.font.color.type is not None:
                try:
                    template_color = r0.font.color.rgb
                except Exception:
                    pass
            if r0.font.bold is not None:
                template_bold = r0.font.bold
    except Exception:
        pass

    text_frame.clear()

    paragraphs = [p.strip() for p in content.split("\n\n") if p.strip()]
    if not paragraphs:
        paragraphs = [content.strip()]

    for i, para_text in enumerate(paragraphs):
        p = text_frame.paragraphs[0] if i == 0 else text_frame.add_paragraph()
        p.space_after  = Pt(6)
        p.space_before = Pt(0)
        # Explicit 18pt leading for body text (≈1.45 ratio at 12pt body).
        # See docs/REPORT-QUALITY-RESEARCH-2026.md §Phase 2 Typography.
        p.line_spacing = Pt(18)

        run = p.add_run()
        run.text = para_text
        if template_name:
            run.font.name = template_name
        if template_size:
            run.font.size = template_size
        if template_color:
            run.font.color.rgb = template_color
        run.font.bold = template_bold


def _populate_bullet_list(
    text_frame: Any,
    content: str,
    prefix: str = "•",
    font_size: Any = None,
    font_color: Any = None,
    prefix_color: Any = None,
) -> None:
    """
    Replace a text frame's content with a formatted bullet list.
    Splits content on single newlines into individual items.
    Strips any existing prefix chars (✓ ⚠ • numbered) before re-adding prefix.
    """
    import re as _re

    if not content or not content.strip():
        text_frame.clear()
        return

    # Capture template formatting
    template_name: str | None = None
    template_size: Any = font_size
    template_color: Any = font_color
    try:
        first_para = text_frame.paragraphs[0]
        if first_para.runs:
            r0 = first_para.runs[0]
            if r0.font.name:
                template_name = r0.font.name
            if r0.font.size:
                template_size = r0.font.size
            if r0.font.color and r0.font.color.type is not None:
                try:
                    template_color = r0.font.color.rgb
                except Exception:
                    pass
    except Exception:
        pass

    text_frame.clear()

    items = [item.strip() for item in content.split("\n") if item.strip()]

    for i, item_text in enumerate(items):
        p = text_frame.paragraphs[0] if i == 0 else text_frame.add_paragraph()
        p.space_after  = Pt(8)
        p.space_before = Pt(2)
        # Explicit 18pt leading (~1.45 ratio) — matches body paragraph style.
        p.line_spacing = Pt(18)

        # Strip any existing prefix that AI or previous code may have added
        clean = item_text
        for strip_pfx in ("✓ ", "⚠ ", "• ", "→ ", "- "):
            if clean.startswith(strip_pfx):
                clean = clean[len(strip_pfx):]
                break
        # Strip numbered prefixes like "1. " or "2) "
        clean = _re.sub(r"^\d+[.)]\s*", "", clean)

        # Prefix run (bold, coloured)
        if prefix:
            pr = p.add_run()
            pr.text = f"{prefix}  "
            if template_name:
                pr.font.name = template_name
            if template_size:
                pr.font.size = template_size
            pr.font.bold = True
            if prefix_color:
                pr.font.color.rgb = prefix_color

        # Content run
        cr = p.add_run()
        cr.text = clean
        if template_name:
            cr.font.name = template_name
        if template_size:
            cr.font.size = template_size
        if template_color:
            cr.font.color.rgb = template_color


# ── CSV slide helpers ────────────────────────────────────────────────────────

# Shared unit-alias lookup used by both _populate_csv_slide and select_kpis
_CSV_UNIT_ALIASES: dict[str, str] = {
    "%": "percent",
    "$": "currency", "₹": "currency", "€": "currency",
    "£": "currency", "¥": "currency", "rs": "currency",
}


def _fmt_csv_value(raw: Any, unit: str, cur_sym: str) -> str:
    """
    Format a single CSV metric value with its unit.

    unit must already be normalised to lowercase semantic form
    ("currency" | "percent" | "number").
    Uses 2 decimal places for currency values < 10 (e.g. ₹0.24 not ₹0).
    """
    try:
        num = float(str(raw)) if raw != "" else None
    except (ValueError, TypeError):
        num = None
    if num is None:
        return str(raw) if raw != "" else ""
    if unit == "currency":
        return f"{cur_sym}{num:,.2f}" if num < 10 else f"{cur_sym}{num:,.0f}"
    if unit == "percent":
        return f"{num:.2f}%"
    return f"{int(num):,}" if num == int(num) else f"{num:,.2f}"


def _duplicate_slide(prs: Any, source_slide_idx: int) -> int:
    """
    Duplicate the slide at *source_slide_idx* and append the copy at the end.

    Copies the entire shape tree (text boxes, solid background shapes).
    Does NOT copy embedded images — suitable for template slides whose
    images are inserted programmatically after duplication.

    Returns the 0-based index of the newly created slide.
    """
    import copy as _copy
    source = prs.slides[source_slide_idx]
    new_slide = prs.slides.add_slide(source.slide_layout)

    src_tree = source.shapes._spTree
    dst_tree = new_slide.shapes._spTree

    # Replace the new slide's default shapes with a deep copy of the source
    for elem in list(dst_tree):
        dst_tree.remove(elem)
    for elem in src_tree:
        dst_tree.append(_copy.deepcopy(elem))

    return len(prs.slides) - 1


def _populate_csv_slide(
    slide: Any,
    csv_src: dict,
    cur_sym: str,
    charts: Dict[str, str],
) -> None:
    """
    Populate ONE CSV template slide with data from *csv_src*.

    - Replaces all {{csv_source_name}}, {{csv_kpi_N_label/value/change}} tokens.
    - Embeds the matching CSV chart image, respecting the KPI-card clearance
      rule (chart starts at ≥ 4.30" to avoid overlapping KPI values).
    - When no chart image exists the placeholder shape is silently deleted.
    """
    source_name = csv_src.get("source_name", csv_src.get("name", "Custom Data"))

    # ── Build per-source replacement dict ────────────────────────────────────
    src_repl: dict[str, str] = {"{{csv_source_name}}": source_name}

    csv_metrics = list(csv_src.get("metrics", []))[:6]
    while len(csv_metrics) < 6:
        csv_metrics.append({})

    for i, metric in enumerate(csv_metrics):
        unit = _CSV_UNIT_ALIASES.get(
            (metric.get("unit") or "number").lower().strip(),
            (metric.get("unit") or "number").lower().strip(),
        )
        raw_curr = metric.get("current_value", "")
        raw_prev = metric.get("previous_value")
        val_str  = _fmt_csv_value(raw_curr, unit, cur_sym)

        try:
            curr_num = float(str(raw_curr)) if raw_curr != "" else None
            prev_num = float(str(raw_prev)) if raw_prev is not None else None
        except (ValueError, TypeError):
            curr_num = prev_num = None

        if curr_num is not None and prev_num is not None and prev_num > 0:
            chg = round((curr_num - prev_num) / prev_num * 100, 1)
            change_str = f"+{chg}%" if chg >= 0 else f"{chg}%"
        else:
            change_str = ""

        src_repl[f"{{{{csv_kpi_{i}_label}}}}"]  = (metric.get("name") or "").upper()
        src_repl[f"{{{{csv_kpi_{i}_value}}}}"]  = val_str
        src_repl[f"{{{{csv_kpi_{i}_change}}}}"] = change_str

    _replace_placeholders_in_slide(slide, src_repl)

    # ── Embed CSV chart image ─────────────────────────────────────────────────
    # Chart key in the charts dict is "csv_{safe_source_name}"
    safe_name  = source_name.lower().replace(" ", "_").replace("/", "_")
    chart_path = charts.get(f"csv_{safe_name}")

    # Fallback: when there is exactly ONE csv_* chart available and the
    # safe_name didn't match exactly (e.g. chart generator used a slightly
    # different normalisation), use that one chart.  With multiple sources
    # we intentionally skip the fallback to avoid showing the wrong chart.
    if not chart_path:
        csv_chart_paths = [
            _p for _k, _p in charts.items()
            if _k.startswith("csv_") and _p and os.path.exists(_p)
        ]
        if len(csv_chart_paths) == 1:
            chart_path = csv_chart_paths[0]

    _SLIDE_WIDTH_EMU = int(13.33 * 914400)   # standard widescreen width
    _MAX_CHART_W_EMU = int(8.0 * 914400)    # max chart width for bar charts

    for shape in list(slide.shapes):
        if not shape.has_text_frame:
            continue
        text = shape.text_frame.text.strip()
        if "{{chart_csv_data}}" in text or text.startswith("{{chart_csv_"):
            if chart_path and os.path.exists(chart_path):
                # Template placeholder gives us the correct top and height
                # (Phase 2 positioned it 0.20" below the last KPI card).
                # Override width to max 8.0" and center horizontally so the
                # comparison bar chart maintains a readable ~2:1 aspect ratio
                # instead of being squeezed across the full 11.70" placeholder.
                chart_w = min(shape.width, _MAX_CHART_W_EMU)
                chart_l = (_SLIDE_WIDTH_EMU - chart_w) // 2
                slide.shapes.add_picture(
                    chart_path,
                    chart_l,
                    shape.top,
                    chart_w,
                    shape.height,
                )
                for para in shape.text_frame.paragraphs:
                    for run in para.runs:
                        run.text = ""
            else:
                _delete_shape(slide, shape)
            break


def _reorder_slides(prs: Any, desired_index_order: list[int]) -> None:
    """
    Physically reorder slides in *prs* to match *desired_index_order*.

    *desired_index_order* is a permutation of range(len(prs.slides)) that
    lists current slide indices in the desired output order.  Indices that
    are omitted are appended at the end in their current relative order.

    Manipulates prs.slides._sldIdLst (the XML slide-reference list) directly —
    the standard python-pptx approach for slide reordering.
    """
    n = len(prs.slides)
    if not desired_index_order or n == 0:
        return

    # Append any indices not explicitly listed (preserve relative order)
    listed_set = set(desired_index_order)
    full_order = desired_index_order + [i for i in range(n) if i not in listed_set]

    # Validate
    if sorted(full_order) != list(range(n)):
        logger.warning(
            "_reorder_slides: invalid index list (len=%d, expected %d) — skipping reorder",
            len(full_order), n,
        )
        return

    slide_id_list = prs.slides._sldIdLst
    all_sld_ids = list(slide_id_list)   # snapshot of current sldId elements

    # Clear the list, then re-append in desired order
    for elem in all_sld_ids:
        slide_id_list.remove(elem)
    for i in full_order:
        slide_id_list.append(all_sld_ids[i])

    logger.debug("Slides reordered: %s → %s", list(range(n)), full_order)


# Keys that get formatted population instead of plain text replacement.
# These are popped from replacements before the generic pass so the
# generic _replace_placeholders_in_slide() doesn't flatten them first.
_NARRATIVE_KEYS = {
    "{{executive_summary}}",
    "{{website_narrative}}",
    "{{ads_narrative}}",
}
_LIST_KEYS = {
    "{{key_wins}}",
    "{{concerns}}",
    "{{next_steps}}",
}
_FORMATTED_KEYS = _NARRATIVE_KEYS | _LIST_KEYS


def generate_pptx_report(
    data: Dict[str, Any],
    narrative: Dict[str, str],
    charts: Dict[str, str],
    client_info: Dict[str, Any],
    enabled_sections: dict | None = None,
    template: str = "full",
    custom_section: dict | None = None,
    branding: dict | None = None,
    visual_template: str = "modern_clean",
    language: str = "en",
    cover_customization: dict | None = None,
) -> bytes:
    """
    Generate a report by populating a pre-designed PPTX template.

    Uses adaptive slide selection: only includes slides with actual data.
    Never shows empty slides or N/A KPIs.

    Args:
        data:             Combined data dict (ga4, meta_ads, google_ads, search_console, csv_sources).
        narrative:        AI-generated narrative sections dict.
        charts:           {chart_name: file_path} for PNG chart images.
        client_info:      Dict with "name", "agency_name".
        enabled_sections: Per-client section toggles (all True by default).
        template:         "full" | "summary" | "brief" — detail level.
        custom_section:   {"title": str, "text": str} for the optional custom section.
        branding:         Agency branding dict.
        visual_template:  "modern_clean" | "dark_executive" | "colorful_agency" | etc.

    Returns raw bytes suitable for writing to disk or streaming as a download.
    """
    # Load the visual template
    template_path = VISUAL_TEMPLATES.get(visual_template, VISUAL_TEMPLATES["modern_clean"])
    if not os.path.exists(template_path):
        logger.warning("Visual template %s not found, falling back to modern_clean", visual_template)
        template_path = VISUAL_TEMPLATES["modern_clean"]

    prs = Presentation(template_path)
    num_slides = len(prs.slides)

    # Determine whether this is a legacy 9-slide template or new 19-slide template
    use_legacy = num_slides <= 10
    logger.info("Template %s has %d slides — using %s slide mapping",
                visual_template, num_slides, "legacy" if use_legacy else "v2")

    # Build the full replacement map, then split out the keys that get
    # rich formatting (paragraphs / bullet lists) so the generic pass
    # doesn't flatten them to a single unstyled run first.
    replacements = _build_replacements(data, narrative, client_info, branding, custom_section, template, language)
    formatted_content: dict[str, str] = {}
    for key in _FORMATTED_KEYS:
        if key in replacements:
            formatted_content[key] = replacements.pop(key)

    # Also add new narrative keys for expanded sections
    extra_narratives = {
        "{{google_ads_narrative}}": _narrative_to_text(narrative.get("google_ads_performance", "")),
        "{{seo_narrative}}": _narrative_to_text(narrative.get("seo_performance", "")),
    }
    for k, v in extra_narratives.items():
        if v:
            formatted_content[k] = v
        else:
            replacements[k] = ""  # Clear unused placeholders

    # ── KPI sparkline embedding (pre-replacement) ──────────────────────────
    # Must run BEFORE _replace_placeholders_in_slide so the {{kpi_i_value}}
    # and {{kpi_i_change}} tokens are still present for shape lookup.
    try:
        cur_sym_spark = _currency_symbol(data)
        _selected = select_kpis(data, currency_symbol=cur_sym_spark)
        _kpi_labels = [k["label"] for k in _selected]
        _embed_kpi_sparklines(prs, charts, _kpi_labels)
    except Exception as exc:
        logger.warning("KPI sparkline embedding skipped: %s", exc)

    # ── Cover preset (Phase 3) — MUST run BEFORE placeholder substitution
    # so apply_cover_preset can identify shapes by their placeholder tokens
    # ({{client_name}}, {{report_period}}, etc.). Recolouring here persists
    # through the later text-replacement pass.
    if cover_customization:
        try:
            from services.cover_presets import apply_cover_preset  # noqa: PLC0415
            apply_cover_preset(
                prs,
                preset=cover_customization.get("preset") or "default",
                headline=cover_customization.get("headline"),
                subtitle=cover_customization.get("subtitle"),
                hero_image_url=cover_customization.get("hero_image_url"),
                brand_color=(branding or {}).get("brand_color"),
                accent_color=cover_customization.get("accent_color"),
            )
        except Exception as exc:
            logger.warning("Cover preset application failed: %s", exc)

    # ── Generic pass: single-value placeholders (names, KPIs, dates, logos) ─
    for slide in prs.slides:
        _replace_placeholders_in_slide(slide, replacements)

    # ── Formatted pass: narrative paragraphs ──────────────────────────────────
    if use_legacy:
        # Legacy 9-slide mapping
        _NARRATIVE_SLIDES = {
            1: "{{executive_summary}}",
            3: "{{website_narrative}}",
            4: "{{ads_narrative}}",
        }
        _LIST_SLIDES = {
            5: ("{{key_wins}}",    "\u2713", RGBColor(0x04, 0x78, 0x57)),
            6: ("{{concerns}}",   "\u26A0", RGBColor(0xD9, 0x77, 0x06)),
            7: ("{{next_steps}}", "\u2192", _hex_to_rgb((branding or {}).get("brand_color") or "#4338CA")),
        }
    else:
        # New 19-slide mapping.
        # Intentionally only maps ONE slide per narrative type — each slide that
        # shares the same placeholder key would otherwise show identical text.
        # Slides not listed here (e.g. website_audience, bounce_rate_analysis)
        # are chart-heavy; their leftover placeholder tokens are cleared by the
        # cleanup pass below.
        _NARRATIVE_SLIDES = {
            SLIDE_INDEX["executive_summary"]:   "{{executive_summary}}",
            SLIDE_INDEX["website_traffic"]:     "{{website_narrative}}",
            SLIDE_INDEX["website_engagement"]:  "{{engagement_narrative}}",
            SLIDE_INDEX["meta_ads_overview"]:   "{{ads_narrative}}",
            SLIDE_INDEX["google_ads_overview"]: "{{google_ads_narrative}}",
            SLIDE_INDEX["seo_overview"]:        "{{seo_narrative}}",
            SLIDE_INDEX.get("conversion_funnel", 14): "{{website_narrative}}",
        }
        _LIST_SLIDES = {
            SLIDE_INDEX["key_wins"]:   ("{{key_wins}}",    "\u2713", RGBColor(0x04, 0x78, 0x57)),
            SLIDE_INDEX["concerns"]:   ("{{concerns}}",   "\u26A0", RGBColor(0xD9, 0x77, 0x06)),
            SLIDE_INDEX["next_steps"]: ("{{next_steps}}", "\u2192", _hex_to_rgb((branding or {}).get("brand_color") or "#4338CA")),
        }

    for slide_idx, placeholder_key in _NARRATIVE_SLIDES.items():
        if slide_idx >= len(prs.slides):
            continue
        content = formatted_content.get(placeholder_key, "")
        if not content:
            continue
        slide = prs.slides[slide_idx]
        for shape in slide.shapes:
            if shape.has_text_frame and placeholder_key in shape.text_frame.text:
                _populate_text_frame_formatted(shape.text_frame, content)
                break

    for slide_idx, (placeholder_key, prefix, prefix_color) in _LIST_SLIDES.items():
        if slide_idx >= len(prs.slides):
            continue
        content = formatted_content.get(placeholder_key, "")
        if not content:
            continue
        slide = prs.slides[slide_idx]
        for shape in slide.shapes:
            if shape.has_text_frame and placeholder_key in shape.text_frame.text:
                _populate_bullet_list(shape.text_frame, content,
                                      prefix=prefix, prefix_color=prefix_color)
                break

    # ── Custom section rich formatting ─────────────────────────────────────────
    cs = custom_section or {}
    cs_text = cs.get("text", "")
    if cs_text:
        cs_slide_idx = SLIDE_INDEX.get("custom_section", 18) if not use_legacy else 8
        if cs_slide_idx < len(prs.slides):
            cs_slide = prs.slides[cs_slide_idx]
            for shape in cs_slide.shapes:
                if shape.has_text_frame and "{{custom_section_text}}" in shape.text_frame.text:
                    # Apply rich formatting using text_formatter
                    from services.text_formatter import parse_structured_text as _parse
                    blocks = _parse(cs_text)
                    tf = shape.text_frame
                    tf.clear()
                    brand_color = _hex_to_rgb((branding or {}).get("brand_color") or "#4338CA")
                    for block in blocks:
                        p = tf.add_paragraph()
                        p.space_after = Pt(4)
                        p.line_spacing = Pt(18)   # consistent with body populators
                        if block["type"] == "header":
                            run = p.add_run()
                            run.text = block["text"]
                            run.font.bold = True
                            run.font.size = Pt(14)
                            run.font.color.rgb = brand_color
                        elif block["type"] == "bullet":
                            pr = p.add_run()
                            pr.text = "•  "
                            pr.font.bold = True
                            pr.font.color.rgb = brand_color
                            cr = p.add_run()
                            cr.text = block["text"]
                            cr.font.size = Pt(11)
                        elif block["type"] == "numbered":
                            pr = p.add_run()
                            pr.text = f"{block['number']}.  "
                            pr.font.bold = True
                            pr.font.color.rgb = brand_color
                            cr = p.add_run()
                            cr.text = block["text"]
                            cr.font.size = Pt(11)
                        else:
                            run = p.add_run()
                            run.text = block["text"]
                            run.font.size = Pt(11)
                    break

    # ── Non-CSV charts ────────────────────────────────────────────────────────
    # CSV chart placeholders are intentionally skipped here; they are embedded
    # per-source by _populate_csv_slide() below, after slide duplication.
    _replace_charts(prs, charts)

    # ── KPI colour coding ─────────────────────────────────────────────────────
    _colorize_kpi_changes(prs, data)

    # ── Logos ─────────────────────────────────────────────────────────────────
    # NOTE: cover preset was already applied earlier in the pipeline (before
    # text substitution) so it could identify shapes by placeholder tokens.
    # Logos are embedded LAST so optional placement/size overrides don't
    # fight the preset's header recolour.
    _embed_logos(prs, branding)

    # ── Custom section image ─────────────────────────────────────────────────
    if not use_legacy:
        custom_slide_idx = SLIDE_INDEX.get("custom_section", 18)
        _embed_custom_section_image(prs, custom_section, custom_slide_idx)
    else:
        _embed_custom_section_image(prs, custom_section, 8)  # legacy index

    # ── Smart slide deletion ──────────────────────────────────────────────────
    if use_legacy:
        # Legacy behavior for existing 9-slide templates
        slides_to_delete = _get_slides_to_delete(enabled_sections, template, narrative, custom_section)
    else:
        # New adaptive selection for 19-slide templates
        selected = select_slides(data, template, custom_section, narrative)
        slides_to_delete = get_slides_to_delete(
            selected, len(data.get("csv_sources", []))
        )

    # Delete in reverse order to preserve indices
    for idx in sorted(slides_to_delete, reverse=True):
        if idx < len(prs.slides):
            try:
                rId = prs.slides._sldIdLst[idx].rId
                prs.part.drop_rel(rId)
                del prs.slides._sldIdLst[idx]
            except Exception as e:
                logger.warning("Could not delete slide %d: %s", idx, e)

    # ── Purge orphaned slide parts ────────────────────────────────────────────
    # drop_rel() removes the relationship but the Part object stays in the
    # package's internal part collection.  When add_slide() later creates
    # duplicates it can pick a /ppt/slides/slideN.xml name that collides
    # with an orphan, causing the orphan to overwrite the new slide during
    # save (python-pptx #689).  A save→reload cycle writes only reachable
    # parts to the byte stream, then reloads — cleanly eliminating orphans.
    if slides_to_delete:
        _buf = io.BytesIO()
        prs.save(_buf)
        _buf.seek(0)
        prs = Presentation(_buf)
        logger.debug("Presentation reloaded after deleting %d slide(s) — orphaned parts purged",
                     len(slides_to_delete))

    # ── CSV multi-slide population ────────────────────────────────────────────
    # After unused slides are deleted, find the csv_data template slide
    # (identified by its {{csv_source_name}} placeholder).  For N CSV sources,
    # duplicate it N-1 times so every source gets its own slide, then populate
    # each with per-source data and the matching chart image.
    #
    # These variables are declared outside the if-block so the slide-reorder
    # step below can reference them regardless of whether CSV data was present.
    _csv_tpl_idx: int | None = None
    _csv_slide_indices: list[int] = []

    if not use_legacy:
        csv_sources_list = data.get("csv_sources", [])
        if csv_sources_list:
            cur_sym = _currency_symbol(data)

            # ── Deduplication guard: drop sources with repeated names ─────────
            _seen_csv_names: set[str] = set()
            _unique_csv: list[dict] = []
            for _src in csv_sources_list:
                _sname = (_src.get("source_name") or _src.get("name") or "Custom").strip()
                if _sname not in _seen_csv_names:
                    _seen_csv_names.add(_sname)
                    _unique_csv.append(_src)
                else:
                    logger.warning(
                        "Duplicate CSV source '%s' skipped — each source name must be unique",
                        _sname,
                    )
            csv_sources_list = _unique_csv

            # Locate the csv_data template slide by its placeholder text
            for _idx, _sl in enumerate(prs.slides):
                for _sh in _sl.shapes:
                    if _sh.has_text_frame and "{{csv_source_name}}" in _sh.text_frame.text:
                        _csv_tpl_idx = _idx
                        break
                if _csv_tpl_idx is not None:
                    break

            if _csv_tpl_idx is not None:
                # First entry = existing template slide; subsequent = freshly duplicated copies.
                # The first source reuses the template slide directly (no duplication needed).
                _csv_slide_indices = [_csv_tpl_idx]
                for _ in range(len(csv_sources_list) - 1):
                    new_idx = _duplicate_slide(prs, _csv_tpl_idx)
                    _csv_slide_indices.append(new_idx)

                # Populate each slide with its corresponding source
                for _slide_idx, _csv_src in zip(_csv_slide_indices, csv_sources_list):
                    _populate_csv_slide(prs.slides[_slide_idx], _csv_src, cur_sym, charts)
                logger.info(
                    "CSV slides populated: %d source(s) → slide indices %s",
                    len(csv_sources_list), _csv_slide_indices,
                )
            else:
                logger.warning("csv_data template slide not found after deletion — CSV data not shown")

    # ── Leftover placeholder cleanup ──────────────────────────────────────────
    # Clear any remaining {{...}} tokens that weren't substituted above.
    # This handles {{website_narrative}} on chart-heavy detail slides (4/5/6),
    # unreferenced optional sections, and any {{csv_*}} tokens on non-CSV slides.
    _leftover_re = re.compile(r"\{\{[^}]+\}\}")
    for _slide in prs.slides:
        for _shape in _slide.shapes:
            if not _shape.has_text_frame:
                continue
            for _para in _shape.text_frame.paragraphs:
                _full = "".join(r.text for r in _para.runs)
                if _leftover_re.search(_full):
                    _cleaned = _leftover_re.sub("", _full).strip()
                    if _para.runs:
                        _para.runs[0].text = _cleaned
                        for _run in _para.runs[1:]:
                            _run.text = ""

    # ── Slide reordering ──────────────────────────────────────────────────────
    # Duplicated CSV slides are appended at the end by _duplicate_slide().
    # Move them back to their correct logical position: immediately after the
    # original csv_data template slide and before the key_wins/concerns/next_steps
    # conclusion slides.
    # Logical order: Cover → Exec → KPIs → GA4 → Ads → SEO → Funnel →
    #                ALL CSV slides → Key Wins → Concerns → Next Steps → Custom
    if _csv_tpl_idx is not None and len(_csv_slide_indices) > 1:
        _n = len(prs.slides)
        _extra_csv = _csv_slide_indices[1:]    # the slides that were appended
        _extra_set  = set(_extra_csv)
        _non_extra  = [i for i in range(_n) if i not in _extra_set]
        try:
            _orig_pos = _non_extra.index(_csv_slide_indices[0])
            _desired  = (
                _non_extra[:_orig_pos + 1]   # everything up to and including original
                + _extra_csv                  # the duplicates — placed right after
                + _non_extra[_orig_pos + 1:]  # conclusion slides and anything after
            )
            _reorder_slides(prs, _desired)
            logger.info(
                "CSV slides reordered: %d extra slide(s) moved to position %d+",
                len(_extra_csv), _orig_pos + 1,
            )
        except (ValueError, Exception) as _reorder_err:
            logger.warning("Could not reorder CSV slides: %s", _reorder_err)

    # ── Post-deletion clean-up ────────────────────────────────────────────────
    # Renumber footers AFTER all slide operations (deletion, duplication,
    # reordering) so "Page N" reflects the true final slide position.
    _renumber_slide_footers(prs, language)
    # Remove email-only CTA shapes ("Questions? Reply to this email …")
    _remove_static_email_cta(prs)

    # Save to bytes
    output = io.BytesIO()
    prs.save(output)
    logger.info("PPTX report generated for %s (detail=%s, visual=%s, slides=%d)",
                client_info.get("name"), template, visual_template, len(prs.slides))
    return output.getvalue()


# ═══════════════════════════════════════════════════════════════════════════════
# ── PDF report — PPTX→PDF via LibreOffice with ReportLab fallback ──────────
# ═══════════════════════════════════════════════════════════════════════════════

# Languages whose scripts ReportLab cannot render without special fonts
_NON_LATIN_LANGUAGES: frozenset[str] = frozenset({
    "hi",  # Hindi (Devanagari)
    "ar",  # Arabic
    "ja",  # Japanese
    "zh",  # Chinese
    "ko",  # Korean
    "th",  # Thai
    "he",  # Hebrew
    "bn",  # Bengali
    "ta",  # Tamil
    "te",  # Telugu
    "mr",  # Marathi (Devanagari)
})

# Noto font family names per language (used by matplotlib font manager lookup)
_NOTO_FONT_FAMILIES: dict[str, list[str]] = {
    "hi": ["Noto Sans Devanagari", "Noto Sans"],
    "mr": ["Noto Sans Devanagari", "Noto Sans"],
    "ar": ["Noto Sans Arabic", "Noto Sans"],
    "ja": ["Noto Sans CJK JP", "Noto Sans JP", "Noto Sans"],
    "zh": ["Noto Sans CJK SC", "Noto Sans SC", "Noto Sans"],
    "ko": ["Noto Sans CJK KR", "Noto Sans KR", "Noto Sans"],
    "th": ["Noto Sans Thai", "Noto Sans"],
    "he": ["Noto Sans Hebrew", "Noto Sans"],
    "bn": ["Noto Sans Bengali", "Noto Sans"],
    "ta": ["Noto Sans Tamil", "Noto Sans"],
    "te": ["Noto Sans Telugu", "Noto Sans"],
}

_noto_registered: dict[str, str] = {}  # language → registered font name


def _try_register_noto_font(language: str) -> str | None:
    """
    Try to find and register a Noto font for *language* with ReportLab.
    Returns the registered font name on success, None if unavailable.
    Caches results so registration only happens once per process lifetime.
    """
    if language in _noto_registered:
        return _noto_registered[language]

    families = _NOTO_FONT_FAMILIES.get(language, ["Noto Sans"])
    try:
        import matplotlib.font_manager as _fm
        from reportlab.pdfbase import pdfmetrics as _pm
        from reportlab.pdfbase.ttfonts import TTFont as _TTF

        for family in families:
            try:
                path = _fm.findfont(_fm.FontProperties(family=family))
                # findfont never raises; check the path contains a recognisable name
                if "Noto" in path or "noto" in path:
                    font_name = f"Noto-{language}"
                    _pm.registerFont(_TTF(font_name, path))
                    _noto_registered[language] = font_name
                    logger.info("Registered Noto font '%s' for language '%s'", path, language)
                    return font_name
            except Exception:
                continue
    except Exception as _e:
        logger.debug("Noto font registration failed for '%s': %s", language, _e)

    _noto_registered[language] = ""  # cache miss so we don't retry
    return None


def generate_pdf_report(
    data: Dict[str, Any],
    narrative: Dict[str, str],
    charts: Dict[str, str],
    client_info: Dict[str, Any],
    enabled_sections: dict | None = None,
    template: str = "full",
    custom_section: dict | None = None,
    branding: dict | None = None,
    visual_template: str = "modern_clean",
    language: str = "en",
) -> bytes | None:
    """
    Generate PDF by first creating PPTX then converting via LibreOffice.
    Falls back to the ReportLab-based PDF generator for Latin languages when
    LibreOffice is unavailable.

    For non-Latin languages (Hindi, Arabic, CJK …) LibreOffice is the only path
    that renders the script correctly.  When LibreOffice is unavailable AND the
    language is non-Latin, returns None so the caller can skip the PDF and show
    "Download PPTX" instead — much better UX than a PDF full of boxes.
    """
    logger.info(
        "PDF generation started: client=%s, language=%s, visual=%s",
        client_info.get("name"), language, visual_template,
    )

    # First generate the PPTX (always correct — python-pptx is Unicode-safe)
    try:
        pptx_bytes = generate_pptx_report(
            data, narrative, charts, client_info,
            enabled_sections, template, custom_section, branding, visual_template,
            language,
        )
        logger.info("PPTX generated (%d bytes) — attempting LibreOffice conversion", len(pptx_bytes))
    except Exception as _pptx_err:
        logger.error("PPTX generation failed inside generate_pdf_report: %s", _pptx_err, exc_info=True)
        raise

    # ── Path 1: LibreOffice (ALL languages — best quality) ───────────────────
    # Tries ALL common install locations across Windows, macOS, Linux, Docker.
    with tempfile.TemporaryDirectory() as tmpdir:
        pptx_path = os.path.join(tmpdir, "report.pptx")
        with open(pptx_path, "wb") as f:
            f.write(pptx_bytes)

        for soffice_cmd in (
            "soffice",              # Windows / macOS (brew), many Linux distros
            "libreoffice",          # Some Linux packages use this name
            r"C:\Program Files\LibreOffice\program\soffice.exe",  # Windows default install
            "/usr/bin/soffice",     # Debian/Ubuntu apt full path
            "/usr/bin/libreoffice", # Alternative full path
            "/Applications/LibreOffice.app/Contents/MacOS/soffice",  # macOS manual install
        ):
            try:
                logger.debug("Trying LibreOffice: %s", soffice_cmd)
                result = subprocess.run(
                    [soffice_cmd, "--headless", "--convert-to", "pdf",
                     "--outdir", tmpdir, pptx_path],
                    timeout=60, check=True, capture_output=True,
                )
                pdf_out = os.path.join(tmpdir, "report.pdf")
                if os.path.exists(pdf_out):
                    with open(pdf_out, "rb") as f:
                        pdf_data = f.read()
                    logger.info(
                        "PDF generated via LibreOffice (%s) for %s — %d bytes",
                        soffice_cmd, client_info.get("name"), len(pdf_data),
                    )
                    return pdf_data
                else:
                    logger.warning(
                        "LibreOffice (%s) exited 0 but no PDF found at %s. stdout=%s",
                        soffice_cmd, pdf_out,
                        (result.stdout or b"").decode("utf-8", errors="replace")[:300],
                    )
            except FileNotFoundError:
                logger.debug("LibreOffice binary not found: %s", soffice_cmd)
            except subprocess.TimeoutExpired:
                logger.warning("LibreOffice conversion timed out (60 s) — %s", soffice_cmd)
            except subprocess.CalledProcessError as _cpe:
                logger.warning(
                    "LibreOffice (%s) exited with code %d. stderr: %s",
                    soffice_cmd, _cpe.returncode,
                    (_cpe.stderr or b"").decode("utf-8", errors="replace")[:500],
                )
            except Exception as _lo_err:
                logger.warning("LibreOffice unexpected error (%s): %s", soffice_cmd, _lo_err)

    logger.info("LibreOffice not available or conversion failed for client=%s", client_info.get("name"))

    # ── Path 2: non-Latin language + no LibreOffice → PPTX-only ─────────────
    if language in _NON_LATIN_LANGUAGES:
        logger.warning(
            "PDF skipped: language '%s' requires LibreOffice for correct script rendering "
            "(Devanagari/CJK/Arabic). Client will be offered PPTX download instead.",
            language,
        )
        return None  # Caller saves None → frontend shows "Download PPTX"

    # ── Path 3: ReportLab fallback (Latin languages only) ────────────────────
    logger.info("Falling back to ReportLab PDF for language=%s", language)
    try:
        return _generate_pdf_reportlab(
            data, narrative, charts, client_info,
            enabled_sections, template, custom_section, branding,
            language=language,
        )
    except Exception as _rl_err:
        logger.error("ReportLab PDF generation failed: %s", _rl_err, exc_info=True)
        return None  # Do not crash — caller handles None gracefully


# ═══════════════════════════════════════════════════════════════════════════════
# ── ReportLab PDF fallback ─────────────────────────────────────────────────
# ═══════════════════════════════════════════════════════════════════════════════

def _generate_pdf_reportlab(
    data: Dict[str, Any],
    narrative: Dict[str, str],
    charts: Dict[str, str],
    client_info: Dict[str, Any],
    enabled_sections: dict | None = None,
    template: str = "full",
    custom_section: dict | None = None,
    branding: dict | None = None,
    language: str = "en",
) -> bytes:
    """
    Generate a clean PDF report using ReportLab (used when LibreOffice is unavailable).
    For non-Latin languages, attempts to use a Noto font; adds a rendering notice
    if the font is not available on this server.
    """
    _br         = branding or {}
    agency_name = _br.get("agency_name") or client_info.get("agency_name") or "Your Agency"
    _brand_rl   = _hex_to_rl(_br.get("brand_color") or "#4338CA")

    # ── Non-Latin language font handling ─────────────────────────────────────
    _pdf_body = _PDF_BODY_FONT
    _pdf_bold = _PDF_BOLD_FONT
    _needs_unicode_font = language in _NON_LATIN_LANGUAGES
    _unicode_font_missing = False
    if _needs_unicode_font:
        noto_name = _try_register_noto_font(language)
        if noto_name:
            _pdf_body = noto_name
            _pdf_bold = noto_name  # Noto Bold may not always be separate; use regular as fallback
            logger.info("Using Noto font '%s' for language '%s' in PDF", noto_name, language)
        else:
            _unicode_font_missing = True
            logger.warning(
                "No suitable font found for language '%s' — PDF may show boxes for non-ASCII text. "
                "Install LibreOffice or Noto fonts on the server for correct rendering.",
                language,
            )

    agency_logo_img = _download_image(_br.get("agency_logo_url", ""))
    client_logo_img = _download_image(_br.get("client_logo_url", ""))

    ga4         = data.get("ga4", {}).get("summary", {})
    meta        = data.get("meta_ads", {}).get("summary", {})
    p_start     = data.get("period_start", "")
    p_end       = data.get("period_end", "")
    client_name = client_info.get("name", "Client")
    report_date = datetime.now().strftime("%B %d, %Y")
    cur_sym     = _currency_symbol(data)
    _page_w, _page_h = A4
    _lm = 0.75 * inch

    # ── Canvas callbacks ──────────────────────────────────────────────────────
    def _draw_cover_canvas(canvas: Any, doc: Any) -> None:
        canvas.saveState()
        band_h = 1.6 * inch
        canvas.setFillColor(_brand_rl)
        canvas.rect(0, _page_h - band_h, _page_w, band_h, fill=1, stroke=0)
        canvas.setFillColor(rl_colors.white)
        canvas.setFont(_pdf_bold, 12)
        canvas.drawString(_lm, _page_h - 0.9 * inch, agency_name.upper())
        canvas.setFont(_pdf_body, 9)
        canvas.drawString(_lm, _page_h - 1.18 * inch, "CLIENT PERFORMANCE REPORT")
        if agency_logo_img:
            try:
                agency_logo_img.seek(0)
                canvas.drawImage(
                    ImageReader(agency_logo_img),
                    _page_w - 2.3 * inch, _page_h - 1.42 * inch,
                    width=1.6 * inch, height=0.8 * inch,
                    preserveAspectRatio=True, mask="auto",
                )
            except Exception:
                pass
        _draw_pdf_footer(canvas, doc)
        canvas.restoreState()

    def _draw_content_canvas(canvas: Any, doc: Any) -> None:
        canvas.saveState()
        y_line = _page_h - 0.55 * inch
        canvas.setStrokeColor(_brand_rl)
        canvas.setLineWidth(1.5)
        canvas.line(_lm, y_line, _page_w - _lm, y_line)
        canvas.setFont(_pdf_body, 8)
        canvas.setFillColor(_RL_SLATE_500)
        canvas.drawString(_lm, _page_h - 0.42 * inch, agency_name)
        canvas.drawRightString(_page_w - _lm, _page_h - 0.42 * inch, f"{t(language, 'page')} {doc.page}")
        _draw_pdf_footer(canvas, doc)
        canvas.restoreState()

    def _draw_pdf_footer(canvas: Any, doc: Any) -> None:
        canvas.setStrokeColor(_RL_SLATE_200)
        canvas.setLineWidth(0.5)
        canvas.line(_lm, 0.52 * inch, _page_w - _lm, 0.52 * inch)
        canvas.setFont(_pdf_body, 8)
        canvas.setFillColor(_RL_SLATE_400)
        canvas.drawString(_lm, 0.34 * inch,
                          f"{t(language, 'prepared_by')} {agency_name}  \u2022  {t(language, 'confidential')}  \u2022  {report_date}")
        canvas.drawRightString(_page_w - _lm, 0.34 * inch, f"{t(language, 'page')} {doc.page}")

    # ── Doc setup ─────────────────────────────────────────────────────────────
    buf = io.BytesIO()
    doc = SimpleDocTemplate(
        buf, pagesize=A4,
        rightMargin=_lm, leftMargin=_lm,
        topMargin=1.75 * inch, bottomMargin=0.75 * inch,
    )
    ss = getSampleStyleSheet()

    h1 = ParagraphStyle("RPH1", parent=ss["Heading1"],
                        fontSize=18, textColor=_brand_rl, spaceBefore=18, spaceAfter=4,
                        leading=22, fontName=_pdf_bold)
    body = ParagraphStyle("RPBody", parent=ss["Normal"],
                          fontSize=11, textColor=_RL_SLATE_700, spaceBefore=4, spaceAfter=4,
                          leading=18, fontName=_pdf_body)
    label_s = ParagraphStyle("RPLabel", parent=ss["Normal"],
                              fontSize=9, textColor=_RL_SLATE_400, spaceAfter=2, leading=11,
                              fontName=_pdf_bold)
    value_s = ParagraphStyle("RPValue", parent=ss["Normal"],
                              fontSize=20, textColor=_RL_SLATE_900, spaceAfter=2, leading=24,
                              fontName=_pdf_bold)
    chg_pos = ParagraphStyle("RPChgPos", parent=ss["Normal"],
                              fontSize=10, textColor=_RL_EMERALD, leading=13, fontName=_pdf_bold)
    chg_neg = ParagraphStyle("RPChgNeg", parent=ss["Normal"],
                              fontSize=10, textColor=_RL_ROSE,    leading=13, fontName=_pdf_bold)
    chg_na  = ParagraphStyle("RPChgNa",  parent=ss["Normal"],
                              fontSize=10, textColor=_RL_SLATE_400, leading=13, fontName=_pdf_body)

    def _chg_style(val: float | None) -> ParagraphStyle:
        if val is None: return chg_na
        return chg_pos if val >= 0 else chg_neg

    def _section_heading(title: str) -> None:
        story.append(Paragraph(title, h1))
        story.append(HRFlowable(width="100%", thickness=1.0, color=_brand_rl, spaceAfter=8))

    # KPI data
    s_chg    = ga4.get("sessions_change")
    u_chg    = ga4.get("users_change")
    c_chg    = ga4.get("conversions_change")
    sp_chg   = meta.get("spend_change")
    roas     = meta.get("roas") or 0.0
    proas    = meta.get("prev_roas") or 0.0
    roas_chg = round((roas - proas) / max(proas, 0.01) * 100, 1) if proas else None
    cpa      = meta.get("cost_per_conversion") or 0.0
    pcpa     = meta.get("prev_cost_per_conversion") or 0.0
    cpa_chg  = round((cpa - pcpa) / max(pcpa, 0.01) * 100, 1) if pcpa else None

    template_label = {"summary": "Summary Report", "brief": "Executive Brief"}.get(template, "Full Performance Report")

    # ── Story ──────────────────────────────────────────────────────────────────
    story: list[Any] = []

    # Cover page
    # ── Unicode font warning banner (non-Latin languages without Noto) ───────
    if _unicode_font_missing:
        story.append(Paragraph(
            "⚠ PDF rendering note: This report is in a non-Latin language but a compatible "
            "font (Noto) was not found on the server. Text may appear as boxes. "
            "Install LibreOffice on the server, or install Noto fonts, for correct rendering. "
            "The PPTX download renders correctly.",
            ParagraphStyle("RPFontWarning", parent=ss["Normal"],
                           fontSize=9, textColor=rl_colors.HexColor("#92400E"),
                           backColor=rl_colors.HexColor("#FEF3C7"),
                           borderColor=rl_colors.HexColor("#F59E0B"),
                           borderWidth=1, borderPadding=6,
                           leading=14, spaceBefore=0, spaceAfter=10),
        ))

    story.append(Paragraph(
        client_name,
        ParagraphStyle("RPClientName", parent=ss["Normal"],
                       fontSize=30, textColor=_RL_SLATE_900, fontName=_pdf_bold,
                       leading=36, spaceBefore=6, spaceAfter=10),
    ))
    if client_logo_img:
        try:
            client_logo_img.seek(0)
            _clogo = RLImage(client_logo_img, width=1.4 * inch, height=1.4 * inch)
            _clogo.hAlign = "LEFT"
            story.append(_clogo)
            story.append(Spacer(1, 0.1 * inch))
        except Exception:
            pass
    story.append(Paragraph(
        f"{template_label}  \u2022  {p_start} to {p_end}",
        ParagraphStyle("RPSub", parent=ss["Normal"], fontSize=13,
                       textColor=_RL_SLATE_500, fontName=_pdf_body, leading=17, spaceAfter=6),
    ))
    story.append(Paragraph(
        (
            f"Prepared by {agency_name}  \u2022  Powered by GoReportPilot  \u2022  {report_date}"
            if _br.get("powered_by_badge", True)
            else f"Prepared by {agency_name}  \u2022  {report_date}"
        ),
        ParagraphStyle("RPPrep", parent=ss["Normal"], fontSize=10,
                       textColor=_RL_SLATE_400, fontName=_pdf_body, leading=14, spaceAfter=16),
    ))
    story.append(HRFlowable(width="100%", thickness=1.5, color=_brand_rl, spaceAfter=0))
    story.append(PageBreak())

    # Executive Summary
    if _section_enabled(enabled_sections, "executive_summary"):
        _exec = _to_lines(narrative.get("executive_summary", ""))
        if _exec:
            _section_heading(t(language, "executive_summary"))
            for line in _exec:
                story.append(Paragraph(line, body))
            story.append(Spacer(1, 0.15 * inch))

    # KPI Scorecard
    if _section_enabled(enabled_sections, "kpi_scorecard"):
        _section_heading(t(language, "key_performance_indicators"))
        col_w = (_page_w - 2 * _lm) / 3
        _neg_cpa = -cpa_chg if cpa_chg is not None else None

        def _kpi_cell(lbl, val, chg):
            return [Paragraph(lbl, label_s), Paragraph(val, value_s),
                    Paragraph(_fmt_change(chg), _chg_style(chg))]

        kpi_data = [
            [_kpi_cell(t(language, "sessions"), f"{ga4.get('sessions',0):,}", s_chg),
             _kpi_cell(t(language, "users"), f"{ga4.get('users',0):,}", u_chg),
             _kpi_cell(t(language, "conversions"), f"{ga4.get('conversions',0):,}", c_chg)],
            [_kpi_cell(t(language, "ad_spend"), f"{cur_sym}{meta.get('spend',0):,.0f}", sp_chg),
             _kpi_cell(t(language, "roas"), f"{roas:.1f}x", roas_chg),
             _kpi_cell(t(language, "cost_per_conv"), f"{cur_sym}{cpa:.2f}", _neg_cpa)],
        ]
        kpi_table = Table(kpi_data, colWidths=[col_w] * 3, rowHeights=[1.05 * inch] * 2)
        kpi_table.setStyle(TableStyle([
            ("BACKGROUND", (0, 0), (2, 0), _RL_SLATE_100),
            ("BACKGROUND", (0, 1), (2, 1), _RL_SLATE_50),
            ("TOPPADDING", (0, 0), (-1, -1), 10),
            ("BOTTOMPADDING", (0, 0), (-1, -1), 10),
            ("LEFTPADDING", (0, 0), (-1, -1), 14),
            ("RIGHTPADDING", (0, 0), (-1, -1), 14),
            ("VALIGN", (0, 0), (-1, -1), "TOP"),
            ("GRID", (0, 0), (-1, -1), 1.5, rl_colors.white),
        ]))
        story.append(kpi_table)
        story.append(Spacer(1, 0.2 * inch))

    # Website Performance
    if _section_enabled(enabled_sections, "website_traffic"):
        s_chart = charts.get("sessions")
        t_chart = charts.get("traffic_sources")
        _web = _to_lines(narrative.get("website_performance", ""))
        if (s_chart and os.path.exists(s_chart)) or (t_chart and os.path.exists(t_chart)) or _web:
            _section_heading(t(language, "website_performance"))
            for c in [s_chart, t_chart]:
                if c and os.path.exists(c):
                    img = RLImage(c, width=6.5 * inch, height=2.9 * inch)
                    img.hAlign = "CENTER"
                    story.append(img)
                    story.append(Spacer(1, 0.1 * inch))
            for line in _web:
                story.append(Paragraph(line, body))
            story.append(Spacer(1, 0.15 * inch))

    # Paid Advertising
    if _section_enabled(enabled_sections, "meta_ads"):
        sc = charts.get("spend_conversions")
        cp = charts.get("campaigns")
        _ads = _to_lines(narrative.get("paid_advertising", ""))
        if (sc and os.path.exists(sc)) or (cp and os.path.exists(cp)) or _ads:
            _section_heading(t(language, "paid_advertising_meta"))
            for c in [sc, cp]:
                if c and os.path.exists(c):
                    img = RLImage(c, width=6.5 * inch, height=2.9 * inch)
                    img.hAlign = "CENTER"
                    story.append(img)
                    story.append(Spacer(1, 0.1 * inch))
            for line in _ads:
                story.append(Paragraph(line, body))
            story.append(Spacer(1, 0.15 * inch))

    # Key Wins
    if _section_enabled(enabled_sections, "key_wins"):
        _wins = _to_lines(narrative.get("key_wins", ""))
        if _wins:
            _section_heading(t(language, "key_wins"))
            win_s = ParagraphStyle("RPWin", parent=body, leftIndent=16, firstLineIndent=-16, spaceAfter=6)
            for line in _wins:
                clean = line.lstrip("\u2022\u2713\u2714\u26A0-\u2013\u2014 ").strip()
                if clean:
                    story.append(Paragraph(f'<font color="#047857"><b>\u2713</b></font>\u2002{clean}', win_s))
            story.append(Spacer(1, 0.15 * inch))

    # Concerns
    if _section_enabled(enabled_sections, "concerns"):
        _concerns = _to_lines(narrative.get("concerns", ""))
        if _concerns:
            _section_heading(t(language, "concerns"))
            con_s = ParagraphStyle("RPCon", parent=body, leftIndent=16, firstLineIndent=-16, spaceAfter=6)
            for line in _concerns:
                clean = line.lstrip("\u2022\u2713\u2714\u26A0-\u2013\u2014 ").strip()
                if clean:
                    story.append(Paragraph(f'<font color="#D97706"><b>\u26A0</b></font>\u2002{clean}', con_s))
            story.append(Spacer(1, 0.15 * inch))

    # Next Steps
    if _section_enabled(enabled_sections, "next_steps"):
        _steps = _to_lines(narrative.get("next_steps", ""))
        if _steps:
            _section_heading(t(language, "next_steps"))
            step_s = ParagraphStyle("RPStep", parent=body, leftIndent=20, firstLineIndent=-20, spaceAfter=8)
            bc = _br.get("brand_color", "#4338CA")
            for i, line in enumerate(_steps, 1):
                clean = line.lstrip("0123456789.)- ").strip()
                if clean:
                    story.append(Paragraph(f'<font color="{bc}"><b>{i}.</b></font>\u2002{clean}', step_s))
            story.append(Spacer(1, 0.1 * inch))

    # CSV Data Sources
    csv_sources_list = data.get("csv_sources", [])
    if csv_sources_list:
        for csv_src in csv_sources_list:
            src_name = csv_src.get("source_name", csv_src.get("name", "Custom Data"))
            metrics  = csv_src.get("metrics", [])
            if not metrics:
                continue
            _section_heading(f"Supplementary Data — {src_name}")
            # Build a 3-column KPI table for the CSV metrics
            csv_col_w = (_page_w - 2 * _lm) / 3
            csv_label_s = ParagraphStyle("RPCsvLabel", parent=label_s)
            csv_value_s = ParagraphStyle("RPCsvVal",   parent=body,
                                          fontSize=14, textColor=_RL_SLATE_900,
                                          fontName=_pdf_bold, leading=18)
            csv_chg_s   = ParagraphStyle("RPCsvChg",  parent=body,
                                          fontSize=9,  textColor=_RL_SLATE_400,
                                          fontName=_pdf_body, leading=11)

            def _csv_kpi_cell(lbl: str, val: str, chg: str) -> list:
                return [
                    Paragraph(lbl.upper(), csv_label_s),
                    Paragraph(val or "—", csv_value_s),
                    Paragraph(chg, csv_chg_s),
                ]

            # Group metrics into rows of 3
            rows = []
            row: list = []
            for metric in metrics:
                m_unit = metric.get("unit", "number")
                m_curr = metric.get("current_value")
                m_name = metric.get("name", "")
                # Format value
                try:
                    m_num = float(str(m_curr)) if m_curr is not None else None
                except (ValueError, TypeError):
                    m_num = None
                if m_num is not None:
                    if m_unit == "currency":
                        m_val = f"{cur_sym}{m_num:,.0f}"
                    elif m_unit == "percent":
                        m_val = f"{m_num:.2f}%"
                    else:
                        m_val = f"{int(m_num):,}" if m_num == int(m_num) else f"{m_num:,.2f}"
                else:
                    m_val = str(m_curr) if m_curr is not None else "—"
                # Format change
                m_prev = metric.get("previous_value")
                try:
                    p_num = float(str(m_prev)) if m_prev is not None else None
                except (ValueError, TypeError):
                    p_num = None
                if m_num is not None and p_num is not None and p_num > 0:
                    m_chg_pct = round((m_num - p_num) / p_num * 100, 1)
                    m_chg = f"+{m_chg_pct}%" if m_chg_pct >= 0 else f"{m_chg_pct}%"
                else:
                    m_chg = ""
                row.append(_csv_kpi_cell(m_name, m_val, m_chg))
                if len(row) == 3:
                    rows.append(row)
                    row = []
            if row:
                # Pad incomplete last row
                while len(row) < 3:
                    row.append(_csv_kpi_cell("", "", ""))
                rows.append(row)

            if rows:
                csv_table = Table(rows, colWidths=[csv_col_w] * 3,
                                  rowHeights=[0.9 * inch] * len(rows))
                csv_table.setStyle(TableStyle([
                    ("BACKGROUND", (0, 0), (-1, -1), _RL_SLATE_50),
                    ("TOPPADDING",    (0, 0), (-1, -1), 8),
                    ("BOTTOMPADDING", (0, 0), (-1, -1), 8),
                    ("LEFTPADDING",   (0, 0), (-1, -1), 12),
                    ("RIGHTPADDING",  (0, 0), (-1, -1), 12),
                    ("VALIGN",        (0, 0), (-1, -1), "TOP"),
                    ("GRID",          (0, 0), (-1, -1), 1.5, rl_colors.white),
                ]))
                story.append(csv_table)
                story.append(Spacer(1, 0.15 * inch))

            # Embed CSV chart if available
            safe_name = src_name.lower().replace(" ", "_").replace("/", "_")
            csv_chart_path = charts.get(f"csv_{safe_name}")
            if csv_chart_path and os.path.exists(csv_chart_path):
                try:
                    csv_img = RLImage(csv_chart_path, width=6.5 * inch, height=2.9 * inch)
                    csv_img.hAlign = "CENTER"
                    story.append(csv_img)
                    story.append(Spacer(1, 0.15 * inch))
                except Exception as _csv_img_err:
                    logger.debug("Could not embed CSV chart in PDF: %s", _csv_img_err)

    # Custom Section
    if (_section_enabled(enabled_sections, "custom_section")
            and custom_section and custom_section.get("title") and custom_section.get("text")):
        _cl = _to_lines(custom_section["text"])
        if _cl:
            _section_heading(custom_section["title"])
            for line in _cl:
                story.append(Paragraph(line, body))

    # Build
    doc.build(story, onFirstPage=_draw_cover_canvas, onLaterPages=_draw_content_canvas)
    buf.seek(0)
    logger.info("PDF report (ReportLab) generated for %s", client_info.get("name"))
    return buf.read()
