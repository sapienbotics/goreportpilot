"""
Cover Customisation — Option F v1 (post-chrome-strip architecture).

Philosophy
----------
The 6 visual-template masters carry only design CHROME on the cover
slide — header bands, gradient strips, divider lines, static labels
("PERFORMANCE REPORT"). Every placeholder text shape has been stripped
by ``scripts/regenerate_cover_thumbnails.py``, so the cover starts
text-free and logo-free at report-generate time.

This module is now the SOLE writer of cover content. It:

  1. Recolours the theme's header band using the user's brand primary
     (only for themes with ``brand_tint_strategy='header_band'``).
     **C-fix:** strips ``<a:lumMod>`` / ``<a:lumOff>`` / ``<a:tint>`` /
     ``<a:shade>`` child elements from the srgbClr so the rendered
     colour exactly matches the user's hex — python-pptx's
     ``fore_color.rgb`` setter does not clear these inherited theme
     modifiers, which caused ``#3AB2CB`` → ``~#4CC5D5`` drift.

  2. Draws a thin (0.10") accent bar at the band/body boundary on
     the body side using the user's accent colour (brand-tinted themes
     only). **v2 fix:** repositioned from ``band.y + band.h - 0.10``
     (inside band — camouflaged against similar-hue tints) to
     ``band.y + band.h`` (just below the band in the body area, high
     contrast across all themes). Colour scrub applied so hex parity
     holds.

  3. Adds a headline text box at ``theme_layout.client_name_box``
     coordinates, using the theme's native font spec. The headline
     content is the user's ``cover_headline`` override, or the
     client's name.

  4. Adds a subtitle / period text box at ``theme_layout.report_period_box``
     coordinates. Contains the period label ("April 2026") and, when
     the user provided a subtitle, appends " — <subtitle>".

  5. **v2 fix (D-D Option D-1):** Adds a "Prepared by <agency_name>"
     attribution text box at ``theme_layout.agency_attribution`` coords.
     Restores the composition density stripped out by the chrome-only
     template pass — modern_clean and gradient_modern had attribution
     INSIDE the tinted band (so the band isn't empty), the other
     themes had it in the footer. Drawn only when the theme has the
     slot configured and a meaningful ``agency_name`` is provided.

  6. Logos are NOT drawn here. They're added by
     ``report_generator._embed_logos`` after this function runs, using
     the client's per-report position + size overrides. When the
     generator finds no placeholder shape on the chrome-only cover, it
     falls back to the theme's ``agency_logo_placeholder`` /
     ``client_logo_placeholder`` coordinates for the "default"
     position.

What this module does NOT do
----------------------------
  * Does not touch content slides (2-19) — they retain placeholder
    shapes and are populated by the generator's generic substitution
    pass.
  * Does not load its own fonts — uses whatever LibreOffice / the user's
    PowerPoint renderer picks for the font name declared in theme_layout.
  * Does not scale text to fit — the configured font sizes are
    template-calibrated; extremely long headlines (>80 chars) may
    overflow their box.
"""
from __future__ import annotations

import logging
from typing import Any, Optional

from pptx.dml.color import RGBColor
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.oxml.ns import qn
from pptx.util import Emu, Inches, Pt

from services.theme_layout import (
    get as get_theme_layout,
    supports_brand_tint,
)

logger = logging.getLogger(__name__)


# ── Public entry point ────────────────────────────────────────────────────────


def apply_cover_customization(
    prs: Any,
    *,
    theme: str,
    headline: str,
    period_label: str,
    subtitle: Optional[str] = None,
    brand_primary_color: Optional[str] = None,
    accent_color: Optional[str] = None,
    agency_name: Optional[str] = None,
) -> None:
    """
    Draw the cover's text + colour overrides on slide[0].

    Arguments
    ---------
    prs
        The ``pptx.Presentation`` instance. Its ``slides[0]`` is the
        cover (chrome-only).
    theme
        Theme id — one of the 6 keys in ``theme_layout.THEME_LAYOUT``.
    headline
        The client name OR the user's custom headline. Required. Falls
        back to "Report" if empty.
    period_label
        Rendered period string, e.g. "April 2026". Required.
    subtitle
        Optional appended subtitle, e.g. "Executive Brief". When set,
        the period line becomes ``"{period_label} — {subtitle}"``.
    brand_primary_color
        Hex (#rrggbb or rrggbb). Applied to the header band. Skipped
        for themes without a band.
    accent_color
        Hex. Drawn as a thin bar at the band/body boundary.
    agency_name
        Agency name used for the "Prepared by …" attribution text.
        Drawn only when the theme has an ``agency_attribution`` slot
        and the name is non-empty. v2 fix (D-D Option D-1).
    """
    try:
        cover = prs.slides[0]
    except Exception as exc:  # noqa: BLE001
        logger.warning("Cover customisation: no cover slide: %s", exc)
        return

    brand_hex = _normalise_hex(brand_primary_color)
    accent_hex = _normalise_hex(accent_color)

    # 1. Header-band tint (only for themes that support it).
    if brand_hex and supports_brand_tint(theme):
        try:
            _recolor_header_band(prs, cover, theme, brand_hex)
        except Exception as exc:  # noqa: BLE001
            logger.debug("Cover: header-band tint failed: %s", exc)

    # 2. Accent bar (only for themes with a band).
    if accent_hex and supports_brand_tint(theme):
        try:
            _draw_accent_bar(prs, cover, theme, accent_hex)
        except Exception as exc:  # noqa: BLE001
            logger.debug("Cover: accent bar draw failed: %s", exc)

    # 3. Headline + 4. period/subtitle text boxes — always drawn.
    try:
        _draw_headline(cover, theme, headline or "Report")
    except Exception as exc:  # noqa: BLE001
        logger.warning("Cover: headline draw failed: %s", exc)

    period_line = period_label if not subtitle else f"{period_label} \u2014 {subtitle}"
    try:
        _draw_period_line(cover, theme, period_line)
    except Exception as exc:  # noqa: BLE001
        logger.warning("Cover: period line draw failed: %s", exc)

    # 5. Agency attribution — v2 fix (D-D Option D-1). Restores
    # composition density for themes whose pre-strip design relied on
    # "Prepared by …" text. Skipped if the theme has no attribution
    # slot or the agency name is empty.
    clean_agency = (agency_name or "").strip()
    if clean_agency:
        try:
            _draw_agency_attribution(cover, theme, clean_agency)
        except Exception as exc:  # noqa: BLE001
            logger.warning("Cover: agency attribution draw failed: %s", exc)


# ── Step 1: header-band recolour (+ C-fix colour scrub) ───────────────────────


def _recolor_header_band(prs: Any, cover: Any, theme: str, hex_str: str) -> None:
    """
    Set the header band's fill to ``hex_str`` (exact, no luminance drift).

    Finds the band via a 2-phase heuristic (see ``_find_header_band``).
    After setting the fill colour, scrubs theme-luminance child
    elements from the resulting ``<a:srgbClr>`` so the rendered colour
    matches the hex byte-for-byte.
    """
    rgb = _hex_to_rgb(hex_str)
    band = _find_header_band(cover, prs, theme)
    if band is None:
        logger.debug("Cover: header band not found for theme %s; skipping tint", theme)
        return
    try:
        band.fill.solid()
        band.fill.fore_color.rgb = RGBColor(*rgb)
        _scrub_color_modifiers_on_shape(band)
        try:
            band.line.fill.background()
        except Exception:  # noqa: BLE001
            pass
        logger.info("Cover: header band tinted with brand %s (theme=%s)", hex_str, theme)
    except Exception as exc:  # noqa: BLE001
        logger.debug("Cover: failed to set header-band fill: %s", exc)


def _find_header_band(cover: Any, prs: Any, theme: str) -> Optional[Any]:
    """
    Heuristic: topmost non-text full-width shape in the top third of the
    slide. For chrome-only covers this lands on the header band shape
    directly. Falls back to coordinate match against the theme layout
    if the heuristic returns no candidate.
    """
    slide_w = prs.slide_width
    slide_h = prs.slide_height
    min_w   = int(slide_w * 0.6)
    top_cap = int(slide_h * 0.33)

    # Phase 1 — heuristic.
    candidates: list[Any] = []
    for shape in cover.shapes:
        if getattr(shape, "shape_type", None) == 13:  # PICTURE
            continue
        if shape.has_text_frame and (shape.text_frame.text or "").strip():
            continue
        try:
            if (shape.width and shape.width >= min_w
                and shape.top is not None and shape.top <= top_cap):
                candidates.append(shape)
        except Exception:  # noqa: BLE001
            continue
    if candidates:
        candidates.sort(key=lambda s: (int(s.top or 0), -int(s.width or 0)))
        return candidates[0]

    # Phase 2 — coordinate fallback.
    band_spec = get_theme_layout(theme).get("header_band")
    if band_spec is None:
        return None
    EMU = 914400
    target_x = int(band_spec["x"] * EMU)
    target_y = int(band_spec["y"] * EMU)
    target_w = int(band_spec["w"] * EMU)
    for shape in cover.shapes:
        try:
            if (shape.left is not None and shape.top is not None
                and abs(int(shape.left) - target_x) < EMU * 0.05
                and abs(int(shape.top)  - target_y) < EMU * 0.1
                and shape.width and abs(int(shape.width) - target_w) < EMU * 0.2):
                return shape
        except Exception:  # noqa: BLE001
            continue
    return None


# ── Step 2: accent bar ────────────────────────────────────────────────────────


def _draw_accent_bar(prs: Any, cover: Any, theme: str, hex_str: str) -> None:
    """
    Thin (0.10") accent bar at the band-body boundary, on the body side.

    v2 fix (D-A Option A-1): previously positioned at
    ``band.y + band.h - 0.10`` which put the bar INSIDE thick bands,
    camouflaging it against similar-hue brand tints (e.g. modern_clean
    with teal primary + teal accent). Moved to ``band.y + band.h`` so
    the bar sits just below the band, in the body area, with consistent
    visibility across all brand-tinted themes regardless of band height.

    Colour-scrubbed for exact hex rendering.
    """
    band_spec = get_theme_layout(theme).get("header_band")
    if band_spec is None:
        return

    rgb = _hex_to_rgb(hex_str)
    EMU = 914400
    bar_h_emu = int(Inches(0.10))
    bar_top_emu = int((band_spec["y"] + band_spec["h"]) * EMU)
    bar_left_emu = int(band_spec["x"] * EMU)
    bar_width_emu = int(band_spec["w"] * EMU)

    bar = cover.shapes.add_shape(
        1,  # MSO_SHAPE.RECTANGLE
        Emu(bar_left_emu), Emu(bar_top_emu),
        Emu(bar_width_emu), Emu(bar_h_emu),
    )
    try:
        bar.fill.solid()
        bar.fill.fore_color.rgb = RGBColor(*rgb)
        _scrub_color_modifiers_on_shape(bar)
        bar.line.fill.background()
    except Exception as exc:  # noqa: BLE001
        logger.debug("Cover: accent bar fill failed: %s", exc)


# ── Step 3: headline text box ─────────────────────────────────────────────────


def _draw_headline(cover: Any, theme: str, headline: str) -> None:
    """Add a text box for the headline at the theme's ``client_name_box``."""
    layout = get_theme_layout(theme)
    box  = layout["client_name_box"]
    font = layout["client_name_font"]
    _add_text_box(
        cover,
        text=headline,
        box=box,
        font_name=font["name"],
        size_pt=int(font["size_pt"]),
        color_hex=font["color_hex"],
        bold=bool(font.get("bold")),
        align=PP_ALIGN.CENTER,
        anchor=MSO_ANCHOR.MIDDLE,
    )


# ── Step 4: period / subtitle line ────────────────────────────────────────────


def _draw_period_line(cover: Any, theme: str, text: str) -> None:
    """Add a text box for the period + optional subtitle."""
    layout = get_theme_layout(theme)
    box  = layout["report_period_box"]
    font = layout["report_period_font"]
    _add_text_box(
        cover,
        text=text,
        box=box,
        font_name=font["name"],
        size_pt=int(font["size_pt"]),
        color_hex=font["color_hex"],
        bold=bool(font.get("bold")),
        align=PP_ALIGN.CENTER,
        anchor=MSO_ANCHOR.TOP,
    )


# ── Step 5: agency attribution ("Prepared by …") ──────────────────────────────


def _draw_agency_attribution(cover: Any, theme: str, agency_name: str) -> None:
    """
    Add a "Prepared by <agency_name>" text box at the theme's
    ``agency_attribution`` coordinates.

    v2 fix (D-D Option D-1). Four of the six themes placed this text
    IN the tinted band or just below it (modern_clean, gradient_modern,
    colorful_agency, bold_geometric), the other two used it as a
    footer signature (dark_executive, minimal_elegant). In every case
    we restore the designer's original position + font + colour so the
    composition matches what the template designer intended when they
    drew the layout.
    """
    spec = get_theme_layout(theme).get("agency_attribution")
    if not spec:
        return
    box  = spec["box"]
    font = spec["font"]
    _add_text_box(
        cover,
        text=f"Prepared by {agency_name}",
        box=box,
        font_name=font["name"],
        size_pt=int(font["size_pt"]),
        color_hex=font["color_hex"],
        bold=bool(font.get("bold")),
        align=_resolve_alignment(spec.get("align")),
        anchor=MSO_ANCHOR.MIDDLE,
    )


def _resolve_alignment(key: Optional[str]) -> Any:
    """Translate a theme_layout alignment string to a PP_ALIGN enum."""
    mapping = {
        "left":   PP_ALIGN.LEFT,
        "center": PP_ALIGN.CENTER,
        "right":  PP_ALIGN.RIGHT,
    }
    return mapping.get((key or "left").lower(), PP_ALIGN.LEFT)


# ── Shared text-box helper ────────────────────────────────────────────────────


def _add_text_box(
    cover: Any,
    *,
    text: str,
    box: dict,
    font_name: str,
    size_pt: int,
    color_hex: str,
    bold: bool,
    align: Any,
    anchor: Any,
) -> None:
    """Add a single-paragraph text box with the given formatting."""
    shape = cover.shapes.add_textbox(
        Inches(float(box["x"])),
        Inches(float(box["y"])),
        Inches(float(box["w"])),
        Inches(float(box["h"])),
    )
    tf = shape.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    # No internal margin — text_box starts with 0.1" default which would
    # visually shift the content off its spec position.
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0

    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text or ""
    run.font.name = font_name
    run.font.size = Pt(size_pt)
    run.font.bold = bold
    rgb = _hex_to_rgb(color_hex)
    run.font.color.rgb = RGBColor(*rgb)


# ── Colour XML scrubber (C-fix) ───────────────────────────────────────────────


_MODIFIER_TAGS = (
    "lumMod", "lumOff", "tint", "shade", "alpha",
    "gamma", "invGamma", "satMod", "satOff", "hueMod", "hueOff",
    "comp", "inv", "gray", "red", "redMod", "redOff",
    "green", "greenMod", "greenOff", "blue", "blueMod", "blueOff",
)


def _scrub_color_modifiers_on_shape(shape: Any) -> None:
    """
    Remove theme-luminance modifier children from every ``<a:srgbClr>``
    descendant of ``<p:spPr>/<a:solidFill>``.

    This is the C-fix for hex drift: python-pptx's ``fore_color.rgb``
    setter writes a new ``val`` attribute on the existing ``<a:srgbClr>``
    but does not strip its existing modifier children (``<a:lumMod>``,
    ``<a:lumOff>``, etc.) when the template had them. Those modifiers
    then re-apply over the user's exact hex at render time, shifting
    the colour. We remove them here.
    """
    try:
        spPr = shape._element.spPr
    except AttributeError:
        return
    for solid_fill in spPr.findall(qn("a:solidFill")):
        for srgb in solid_fill.findall(qn("a:srgbClr")):
            for tag in _MODIFIER_TAGS:
                for child in srgb.findall(qn(f"a:{tag}")):
                    srgb.remove(child)


# ── Small helpers ─────────────────────────────────────────────────────────────


def _normalise_hex(value: Optional[str]) -> Optional[str]:
    if not value:
        return None
    v = value.strip()
    if v.startswith("#"):
        v = v[1:]
    if len(v) == 6:
        try:
            int(v, 16)
            return v.upper()
        except ValueError:
            return None
    return None


def _hex_to_rgb(hex_str: str) -> tuple[int, int, int]:
    h = _normalise_hex(hex_str) or "4338CA"
    return int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16)
