"""
Report generation and management endpoints.
Trigger report generation, list reports, get report details, download files.
"""
import asyncio
import logging
import os
import uuid
from datetime import datetime

from fastapi import APIRouter, Depends, HTTPException, Request, status
from fastapi.responses import FileResponse

from middleware.auth import get_current_user_id
from middleware.plan_enforcement import get_user_subscription, can_use_feature
from services.plans import get_plan
from models.schemas import (
    CoverPreviewRequest,
    ReportGenerateRequest,
    ReportListItem,
    ReportListResponse,
    ReportResponse,
    ReportSectionRegenerateRequest,
    ReportSendRequest,
    ReportUpdateRequest,
)
from services.supabase_client import get_supabase_admin
from slowapi import Limiter
from slowapi.util import get_remote_address

limiter = Limiter(key_func=get_remote_address)
logger = logging.getLogger(__name__)
router = APIRouter()


def _sanitize_data_for_ai(data: dict) -> dict:
    """
    Return a deep copy of *data* with all float values rounded to 2 decimal places.

    Prevents raw floats (e.g. 2.9600000001, 0.030000000000000002) from appearing
    verbatim in AI-generated report text.  NaN and Infinity are replaced with 0.
    """
    import copy
    import math

    def _clean(obj):
        if isinstance(obj, float):
            if not math.isfinite(obj):
                return 0.0
            return round(obj, 2)
        if isinstance(obj, dict):
            return {k: _clean(v) for k, v in obj.items()}
        if isinstance(obj, list):
            return [_clean(item) for item in obj]
        return obj

    return _clean(copy.deepcopy(data))


def _map_db_row(row: dict, client_name: str | None = None) -> dict:
    """
    Translate raw Supabase report row to the field names expected by
    ReportResponse / ReportListItem.

    DB column  →  API field
    ai_narrative  →  narrative
    sections      →  data_summary  (extracted from sections JSONB)
    pptx_file_url →  pptx_url
    pdf_file_url  →  pdf_url
    """
    sections = row.get("sections") or {}
    return {
        "id":            row["id"],
        "user_id":       row["user_id"],
        "client_id":     row["client_id"],
        "client_name":   client_name,
        "title":         row["title"],
        "status":        row["status"],
        "period_start":  str(row["period_start"]),
        "period_end":    str(row["period_end"]),
        "narrative":     row.get("ai_narrative"),
        "data_summary":  sections.get("data_summary") if isinstance(sections, dict) else None,
        "meta_currency": sections.get("meta_currency", "USD") if isinstance(sections, dict) else "USD",
        "user_edits":    row.get("user_edits"),
        "pptx_url":      row.get("pptx_file_url"),
        "pdf_url":       row.get("pdf_file_url"),
        "created_at":    row["created_at"],
        "updated_at":    row["updated_at"],
    }

# Local storage directory — lives inside backend/generated_reports/
# Gitignored; move to Supabase Storage in a later phase.
_HERE = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))  # backend/
REPORTS_BASE_DIR = os.path.join(_HERE, "generated_reports")


# ---------------------------------------------------------------------------
# POST /api/reports/preview-cover  (Phase 3)
# ---------------------------------------------------------------------------


@router.post("/preview-cover")
async def preview_cover(
    body: CoverPreviewRequest,
    user_id: str = Depends(get_current_user_id),
):
    """
    Generate a single-slide PPTX preview of a cover design (Option F v1).

    Loads the client's theme (or the per-request override), opens the
    matching visual-template PPTX, keeps only slide[0], applies the
    minimal cover customisations (text, brand tint, accent bar),
    substitutes the remaining placeholder tokens with sample values,
    embeds logos at the requested positions, returns PPTX bytes.
    """
    from fastapi.responses import Response  # noqa: PLC0415

    supabase = get_supabase_admin()

    client_result = (
        supabase.table("clients")
        .select("*")
        .eq("id", body.client_id)
        .eq("user_id", user_id)
        .eq("is_active", True)
        .single()
        .execute()
    )
    if not client_result.data:
        raise HTTPException(status_code=status.HTTP_404_NOT_FOUND, detail="Client not found")
    client = client_result.data

    # Resolve theme + override fields — explicit overrides > stored values.
    theme          = body.theme           or client.get("theme")           or "modern_clean"
    headline       = body.headline        if body.headline        is not None else client.get("cover_headline")
    subtitle       = body.subtitle        if body.subtitle        is not None else client.get("cover_subtitle")
    primary_color  = body.primary_color   if body.primary_color   is not None else client.get("cover_brand_primary_color")
    accent_color   = body.accent_color    if body.accent_color    is not None else client.get("cover_brand_accent_color")
    agency_pos     = body.agency_logo_position or client.get("cover_agency_logo_position") or "default"
    agency_sz      = body.agency_logo_size     or client.get("cover_agency_logo_size")     or "default"
    client_pos     = body.client_logo_position or client.get("cover_client_logo_position") or "default"
    client_sz      = body.client_logo_size     or client.get("cover_client_logo_size")     or "default"

    logger.info(
        "preview_cover[%s] theme=%r headline=%r subtitle=%r primary=%r accent=%r",
        client["id"], theme, headline, subtitle, primary_color, accent_color,
    )

    # Fetch agency branding for logo + default brand fallback.
    profile_result = (
        supabase.table("profiles")
        .select("agency_name,agency_logo_url,brand_color,agency_email")
        .eq("id", user_id)
        .maybe_single()
        .execute()
    )
    _profile = profile_result.data or {}
    branding = {
        "agency_name":          (_profile.get("agency_name") or "Your Agency").strip() or "Your Agency",
        "agency_logo_url":      _profile.get("agency_logo_url") or "",
        "brand_color":          primary_color or _profile.get("brand_color") or "#4338CA",
        "accent_color":         accent_color or "",
        "agency_logo_position": agency_pos,
        "agency_logo_size":     agency_sz,
        "client_logo_position": client_pos,
        "client_logo_size":     client_sz,
        "client_logo_url":      client.get("logo_url") or "",
        "powered_by_badge":     False,
        # Theme hint consumed by _embed_logos when no placeholder shape
        # exists on the chrome-only cover (A-fix plumbing).
        "_cover_theme":         theme,
    }

    def _render_preview() -> bytes:
        import io as _io  # noqa: PLC0415
        from pptx import Presentation  # noqa: PLC0415
        from services.report_generator import (  # noqa: PLC0415
            VISUAL_TEMPLATES,
            _embed_logos,
            _replace_placeholders_in_slide,
        )
        from services.cover_customization import apply_cover_customization  # noqa: PLC0415

        tpl_path = VISUAL_TEMPLATES.get(theme, VISUAL_TEMPLATES["modern_clean"])
        prs = Presentation(tpl_path)

        # Keep ONLY slide[0]; delete everything else.
        for idx in range(len(prs.slides) - 1, 0, -1):
            rId = prs.slides._sldIdLst[idx].rId
            prs.part.drop_rel(rId)
            del prs.slides._sldIdLst[idx]

        # 1) Cover customisation is the SOLE writer of cover text + colours
        #    on the chrome-only cover (no placeholders remain to substitute).
        import datetime as _dt  # noqa: PLC0415
        today = _dt.date.today()
        _effective_headline = (headline or "").strip() or (client.get("name") or "Acme Corp")
        apply_cover_customization(
            prs,
            theme=theme,
            headline=_effective_headline,
            period_label=today.strftime("%B %Y"),
            subtitle=subtitle,
            brand_primary_color=branding["brand_color"],
            accent_color=branding["accent_color"] or None,
            agency_name=branding["agency_name"],
        )

        # 2) Stray-token substitution on slide[0] — a safety net for any
        #    leftover placeholders in the chrome (e.g. a stale template
        #    that hasn't been re-stripped). For the 6 design templates
        #    this pass is a no-op because the cover is chrome-only.
        sample = {
            "{{agency_name}}":   branding["agency_name"],
            "{{agency_email}}":  _profile.get("agency_email") or "",
            "{{agency_logo}}":   "",
            "{{client_logo}}":   "",
            "{{footer_text}}":   f"Prepared by {branding['agency_name']}",
        }
        for slide in prs.slides:
            _replace_placeholders_in_slide(slide, sample)

        # 3) Logos LAST so colour changes don't touch them.
        _embed_logos(prs, branding)

        buf = _io.BytesIO()
        prs.save(buf)
        return buf.getvalue()

    pptx_bytes = await asyncio.to_thread(_render_preview)

    return Response(
        content=pptx_bytes,
        media_type="application/vnd.openxmlformats-officedocument.presentationml.presentation",
        headers={
            "Content-Disposition": 'attachment; filename="cover-preview.pptx"',
        },
    )


# ---------------------------------------------------------------------------
# _generate_report_internal — shared by the API endpoint and the scheduler
# ---------------------------------------------------------------------------

async def _generate_report_internal(
    *,
    client_id: str,
    user_id: str,
    period_start: str,
    period_end: str,
    template: str = "full",
    visual_template: str = "modern_clean",
    csv_sources: list[dict] | None = None,
    supabase=None,
    report_id: str | None = None,
) -> tuple[dict, str]:
    """
    Core report generation pipeline.  Returns ``(raw DB row, client_name)``.
    Raises HTTPException (or any exception) on failure.

    Shared between:
      • POST /api/reports/generate       (API endpoint)
      • POST /api/reports/{id}/regenerate (rebuild after file loss)
      • services/scheduler.py            (automated scheduled reports)

    Pass ``report_id`` to rebuild an existing report in place: the row is
    UPDATEd rather than INSERTed, stale user edits are cleared, and the trial
    report-count check is skipped because no new report is being created.
    Regeneration used to be a separate 480-line copy of this function; the copy
    had drifted (dropped CSV sources, hardcoded chart theme, skipped white-label
    enforcement), which is why it is gone.
    """
    if supabase is None:
        supabase = get_supabase_admin()

    is_regeneration = report_id is not None
    report_id = report_id or str(uuid.uuid4())

    # 0 — Subscription status check: block expired/cancelled users
    sub = get_user_subscription(user_id)
    if sub.get("status") in ("expired", "cancelled"):
        raise HTTPException(
            status_code=status.HTTP_403_FORBIDDEN,
            detail="Your subscription has expired. Please upgrade to continue generating reports.",
        )

    # 0a — Trial report limit (5 reports during free trial).
    # Regeneration rebuilds an existing report, so it does not consume a slot.
    if sub.get("status") == "trialing" and not is_regeneration:
        report_count_resp = (
            supabase.table("reports")
            .select("id", count="exact")
            .eq("user_id", user_id)
            .execute()
        )
        trial_report_count = report_count_resp.count if report_count_resp.count is not None else 0
        if trial_report_count >= 5:
            raise HTTPException(
                status_code=status.HTTP_403_FORBIDDEN,
                detail="You've used all 5 trial reports. Upgrade to a paid plan for unlimited reports.",
            )

    # 1 — Verify client ownership
    client_result = (
        supabase.table("clients")
        .select("*")
        .eq("id", client_id)
        .eq("user_id", user_id)
        .eq("is_active", True)
        .single()
        .execute()
    )
    if not client_result.data:
        raise HTTPException(status_code=status.HTTP_404_NOT_FOUND, detail="Client not found")
    client = client_result.data

    # 1a — Plan enforcement: clamp AI tone and visual template to allowed values
    # (sub was fetched above in step 0)
    user_plan = sub.get("plan", "trial")
    plan_cfg = get_plan(user_plan)
    plan_features = plan_cfg.get("features", {})

    allowed_tones = plan_features.get("ai_tones", ["professional"])
    client_tone = client.get("ai_tone", "professional")
    if client_tone not in allowed_tones:
        client["ai_tone"] = allowed_tones[0]  # override to first allowed tone

    allowed_templates = plan_features.get("visual_templates", ["modern_clean"])
    if visual_template not in allowed_templates:
        visual_template = allowed_templates[0]  # override to first allowed template

    show_powered_by = plan_features.get("powered_by_badge", True)

    # 1b — Read report_config (section toggles, KPI selection, template, custom section)
    report_config: dict = client.get("report_config") or {}
    cfg_sections    = report_config.get("sections", {})     # section toggle dict
    cfg_template    = template or report_config.get("template", "full")
    cfg_custom      = {
        "title": report_config.get("custom_section_title", ""),
        "text":  report_config.get("custom_section_text",  ""),
    }

    # 1c — Connection health pre-generation gate (Phase 2).
    # Block generation when any of the client's connections is broken or
    # expiring_soon so we never ship an empty scheduled report to a client.
    # Bypassed for CSV-only generations (csv_sources supplied, no native
    # connections required).
    unhealthy = (
        supabase.table("connections")
        .select("id,platform,health_status,account_name")
        .eq("client_id", client_id)
        .in_("health_status", ["broken", "expiring_soon"])
        .execute()
    )
    unhealthy_rows = unhealthy.data or []
    if unhealthy_rows and not csv_sources:
        detail_lines = [
            f"{row.get('platform')} ({row.get('health_status')})"
            for row in unhealthy_rows
        ]
        raise HTTPException(
            status_code=status.HTTP_422_UNPROCESSABLE_ENTITY,
            detail=(
                "One or more platform connections need attention before this "
                "report can be "
                + ("regenerated" if is_regeneration else "generated")
                + ": "
                + ", ".join(detail_lines)
                + ". Open the Integrations page and reconnect, then try again."
            ),
        )

    # 2 — Pull every connected source through the adapter registry.
    # This single call replaces the per-platform boilerplate that used to be
    # copy-pasted here and again in regenerate_report. A source that fails is
    # recorded in pull.failures and skipped — never fatal on its own.
    from services.sources import pull_all_sources  # noqa: PLC0415

    pull = await pull_all_sources(
        supabase=supabase,
        client_id=client_id,
        period_start=period_start,
        period_end=period_end,
    )
    meta_currency = pull.currency

    # Build raw_data from ONLY real/connected data — no mock data.
    raw_data: dict = {
        "client_name":  client["name"],
        "period_start": period_start,
        "period_end":   period_end,
        **pull.data,
    }

    # Ad-hoc CSV sources supplied with the request, or replayed on regeneration.
    if csv_sources:
        raw_data["csv_sources"] = csv_sources

    # Require at least one data source.
    if not (pull.has_data or csv_sources):
        if pull.failures:
            detail = (
                "Every connected data source failed to return data: "
                + "; ".join(f"{src} ({msg})" for src, msg in pull.failures.items())
                + "."
            )
            if pull.reauth_required:
                detail += (
                    " Reconnect "
                    + ", ".join(pull.reauth_required)
                    + " on the Integrations page."
                )
        else:
            detail = (
                "No data sources connected for this client. Connect GA4, Meta Ads, "
                "or upload a CSV before generating a report."
            )
        raise HTTPException(status_code=status.HTTP_422_UNPROCESSABLE_ENTITY, detail=detail)

    # 3 — AI narrative (async I/O)
    # Round all floats before passing to GPT-4o so the AI doesn't copy raw
    # floating-point noise (e.g. 2.9600000001) into generated report text.
    from services.ai_narrative import generate_narrative  # noqa: PLC0415
    narrative = await generate_narrative(
        data=_sanitize_data_for_ai(raw_data),
        client_name=client["name"],
        client_goals=client.get("goals_context"),
        tone=client.get("ai_tone", "professional"),
        template=cfg_template,
        language=client.get("report_language", "en") or "en",
    )

    # 4 — Resolve branding and theme BEFORE charts are drawn.
    #
    # Ordering matters and used to be wrong: charts were generated from the
    # request's visual_template and the agency-profile brand colour, while the
    # deck was rendered from the client's theme and the client's brand-colour
    # override. The two disagreed on every client that had either override set,
    # so charts looked pasted in. Everything the renderer needs is now resolved
    # once, up front, and both charts and deck read the same values.
    profile_result = (
        supabase.table("profiles")
        .select("agency_name,agency_logo_url,brand_color,sender_name,agency_email")
        .eq("id", user_id)
        .maybe_single()
        .execute()
    )
    _profile = profile_result.data or {}

    # White-label enforcement: Starter plan gets no custom branding.
    has_white_label = plan_features.get("white_label", False)
    branding = {
        "agency_name":     ((_profile.get("agency_name") or "").strip() or "Your Agency") if has_white_label else "Your Agency",
        "agency_logo_url": (_profile.get("agency_logo_url") or "") if has_white_label else "",
        "brand_color":     (_profile.get("brand_color") or "#4338CA") if has_white_label else "#4338CA",
        "client_logo_url": (client.get("logo_url") or "") if has_white_label else "",
        "powered_by_badge": show_powered_by,
    }

    # Design System (Option F v1) — the client's theme governs the whole deck.
    # Falls back to modern_clean for clients created before migration 017.
    # visual_template is a legacy request field kept for API compatibility.
    client_theme = client.get("theme") or "modern_clean"
    if visual_template and visual_template != client_theme:
        logger.warning(
            "Ignoring legacy visual_template=%r for client %s — using theme=%r",
            visual_template, client_id, client_theme,
        )
    visual_template = client_theme

    # Per-client cover overrides take precedence over the agency defaults.
    # Setting brand_color here (rather than after chart generation, as it was)
    # is what makes the chart palette match the deck.
    #
    # Gated on white_label for the same reason the agency logo and brand colour
    # are: on Starter these per-client colour overrides were the one branding
    # channel that bypassed plan enforcement entirely. Cover and charts read the
    # same resolved values below, so they can no longer disagree.
    _client_primary = (client.get("cover_brand_primary_color") or "") if has_white_label else ""
    _client_accent  = (client.get("cover_brand_accent_color")  or "") if has_white_label else ""
    if _client_primary:
        branding["brand_color"] = _client_primary
    branding["accent_color"]           = _client_accent
    branding["agency_logo_position"]   = client.get("cover_agency_logo_position") or "default"
    branding["agency_logo_size"]       = client.get("cover_agency_logo_size")     or "default"
    branding["client_logo_position"]   = client.get("cover_client_logo_position") or "default"
    branding["client_logo_size"]       = client.get("cover_client_logo_size")     or "default"
    # Theme hint consumed by _embed_logos when no placeholder shape exists on
    # the chrome-only cover — the A-fix strip removed placeholder shapes, so
    # "default" logo position must fall back to theme_layout coordinates.
    branding["_cover_theme"]           = client_theme

    cover_customization = {
        "theme":                client_theme,
        "headline":             client.get("cover_headline"),
        "subtitle":             client.get("cover_subtitle"),
        # Resolved (plan-gated) values, so the cover cannot diverge from the
        # charts and the rest of the deck.
        "brand_primary_color":  _client_primary or None,
        "accent_color":         _client_accent or None,
    }

    client_info = {
        "name":        client["name"],
        "agency_name": branding["agency_name"],
    }

    logger.info(
        "%s[%s] theme=%r brand=%s agency_pos=%r client_pos=%r sources=%s",
        "regenerate" if is_regeneration else "generate",
        client_id, client_theme, branding.get("brand_color"),
        branding.get("agency_logo_position"), branding.get("client_logo_position"),
        list(pull.data) + (["csv_sources"] if csv_sources else []),
    )

    # 5 — Generate charts (sync, run in thread pool), themed to match the deck.
    charts_dir = os.path.join(REPORTS_BASE_DIR, report_id, "charts")

    # Pass AI-generated per-chart action titles (from the narrative engine)
    # so charts render with takeaway headlines instead of generic labels.
    _chart_insights = (narrative or {}).get("chart_insights") or {}
    _report_language = client.get("report_language", "en") or "en"

    from services.chart_generator import generate_all_charts  # noqa: PLC0415
    charts = await asyncio.to_thread(
        generate_all_charts, raw_data, charts_dir, branding["brand_color"], client_theme,
        _chart_insights, _report_language,
    )

    # 6 — Build report files (sync, run in thread pool)
    from services.report_generator import generate_pdf_report, generate_pptx_report  # noqa: PLC0415
    pptx_bytes, pdf_bytes = await asyncio.gather(
        asyncio.to_thread(
            generate_pptx_report, raw_data, narrative, charts, client_info,
            cfg_sections if cfg_sections else None,
            cfg_template,
            cfg_custom if cfg_custom.get("title") else None,
            branding,
            visual_template,
            _report_language,
            cover_customization,
        ),
        asyncio.to_thread(
            generate_pdf_report, raw_data, narrative, charts, client_info,
            cfg_sections if cfg_sections else None,
            cfg_template,
            cfg_custom if cfg_custom.get("title") else None,
            branding,
            visual_template,
            _report_language,
        ),
    )

    # 6 — Save to disk
    report_dir = os.path.join(REPORTS_BASE_DIR, report_id)
    os.makedirs(report_dir, exist_ok=True)

    pptx_path = os.path.join(report_dir, "report.pptx")
    pdf_path  = os.path.join(report_dir, "report.pdf")

    with open(pptx_path, "wb") as f:
        f.write(pptx_bytes)

    if pdf_bytes is not None:
        if not os.path.exists(pdf_path):
            # PDF not yet written (non-trial path, or watermark regen skipped)
            with open(pdf_path, "wb") as f:
                f.write(pdf_bytes)
        db_pdf_path: str | None = pdf_path
    else:
        db_pdf_path = None  # Non-Latin language, LibreOffice unavailable — PPTX only
        logger.info(
            "PDF not saved for report %s — non-Latin language with no LibreOffice. "
            "PPTX download will be offered instead.",
            report_id,
        )

    logger.info("Report files saved to %s", report_dir)

    # 7 — Build data_summary for the DB / frontend preview
    ga4_s  = raw_data.get("ga4",      {}).get("summary", {})
    meta_s = raw_data.get("meta_ads", {}).get("summary", {})
    data_summary = {
        "sessions":            ga4_s.get("sessions"),
        "sessions_change":     ga4_s.get("sessions_change"),
        "users":               ga4_s.get("users"),
        "users_change":        ga4_s.get("users_change"),
        "conversions":         ga4_s.get("conversions"),
        "conversions_change":  ga4_s.get("conversions_change"),
        "pageviews":           ga4_s.get("pageviews"),
        "bounce_rate":         ga4_s.get("bounce_rate"),
        "avg_session_duration": ga4_s.get("avg_session_duration"),
        "spend":               meta_s.get("spend"),
        "spend_change":        meta_s.get("spend_change"),
        "impressions":         meta_s.get("impressions"),
        "clicks":              meta_s.get("clicks"),
        "ctr":                 meta_s.get("ctr"),
        "cpc":                 meta_s.get("cpc"),
        "roas":                meta_s.get("roas"),
        "cost_per_conversion": meta_s.get("cost_per_conversion"),
    }

    # 8 — Human-readable title — use the END month of the reporting period
    # so e.g. "Acme — March 2026 Performance Report" rather than the start.
    # Reports are named by the most recent month they cover.
    try:
        month_year = datetime.strptime(period_end, "%Y-%m-%d").strftime("%B %Y")
    except ValueError:
        month_year = period_end
    title = f"{client['name']} — {month_year} Performance Report"

    # 9 — Persist report record in Supabase (use actual DB column names)
    sections_payload = {
        "data_summary":   data_summary,
        "meta_currency":  meta_currency,
        "ai_model":       "gpt-4.1",
        # Everything needed to rebuild this exact report later. Report files are
        # ephemeral on Railway, so regeneration is a normal-path operation — and
        # it used to lose the CSV sources and the detail level because they were
        # never recorded anywhere. Stored on the existing `sections` JSONB, so
        # no migration is required.
        "generation_settings": {
            "template":        cfg_template,
            "visual_template": client_theme,
            "language":        _report_language,
            "csv_sources":     csv_sources or [],
        },
        # Which sources contributed, and which were connected but failed. Lets
        # the UI explain a thin report instead of leaving the user guessing.
        "source_status": {
            "succeeded":       list(pull.data),
            "failed":          pull.failures,
            "reauth_required": pull.reauth_required,
            "csv_count":       len(csv_sources or []),
            # Set when the narrative engine was unavailable and placeholder text
            # was used. The deck still renders with correct figures, but the
            # agency needs to know to regenerate before sending it on.
            "narrative_error": (narrative or {}).get("_narrative_error"),
        },
        # Compact raw_data for section regeneration (daily arrays omitted to save space)
        "narrative_data": {
            "ga4": {k: v for k, v in raw_data.get("ga4", {}).items() if k != "daily"},
            "meta_ads": {k: v for k, v in raw_data.get("meta_ads", {}).items() if k != "daily"},
            "period_start": raw_data.get("period_start"),
            "period_end":   raw_data.get("period_end"),
        },
    }

    if is_regeneration:
        update_payload = {
            "title":         title,
            "status":        "draft",
            "pptx_file_url": pptx_path,
            "pdf_file_url":  db_pdf_path,
            "ai_narrative":  narrative,
            "user_edits":    None,   # narrative was rebuilt — stale edits no longer apply
            "sections":      sections_payload,
            "updated_at":    datetime.utcnow().isoformat(),
        }
        result = (
            supabase.table("reports")
            .update(update_payload)
            .eq("id", report_id)
            .eq("user_id", user_id)
            .execute()
        )
        failure_detail = "Report regenerated but failed to update database"
    else:
        insert_payload = {
            "id":            report_id,
            "user_id":       user_id,
            "client_id":     client_id,
            "title":         title,
            "status":        "draft",       # CHECK: generating|draft|approved|sent|failed
            "period_start":  period_start,
            "period_end":    period_end,
            "pptx_file_url": pptx_path,     # actual DB column name
            "pdf_file_url":  db_pdf_path,   # None when non-Latin + no LibreOffice
            "ai_narrative":  narrative,     # actual DB column name
            "sections":      sections_payload,
        }
        result = supabase.table("reports").insert(insert_payload).execute()
        failure_detail = "Report generated but failed to save to database"

    if not result.data:
        raise HTTPException(
            status_code=status.HTTP_500_INTERNAL_SERVER_ERROR,
            detail=failure_detail,
        )

    # Return raw DB row so callers can map it however they need
    return result.data[0], client["name"]


# ---------------------------------------------------------------------------
# POST /generate  — thin wrapper around _generate_report_internal
# ---------------------------------------------------------------------------

@router.post("/generate", response_model=ReportResponse, status_code=status.HTTP_201_CREATED)
@limiter.limit("10/hour")
async def generate_report(
    request: Request,
    payload: ReportGenerateRequest,
    user_id: str = Depends(get_current_user_id),
) -> ReportResponse:
    """
    Full report generation pipeline:
    1. Fetch client  2. Data pull (real or mock)  3. AI narrative  4. Charts
    5. PPTX + PDF    6. Save to disk  7. Store in Supabase
    """
    row, client_name = await _generate_report_internal(
        client_id=payload.client_id,
        user_id=user_id,
        period_start=payload.period_start,
        period_end=payload.period_end,
        template=payload.template,
        visual_template=payload.visual_template,
        csv_sources=payload.csv_sources,
    )
    return ReportResponse(**_map_db_row(row, client_name=client_name))


# ---------------------------------------------------------------------------
# GET /client/{client_id}  — list reports for one client
# ---------------------------------------------------------------------------

@router.get("/client/{client_id}", response_model=ReportListResponse)
async def list_client_reports(
    client_id: str,
    user_id: str = Depends(get_current_user_id),
) -> ReportListResponse:
    """List all reports for a specific client (owned by the authenticated user)."""
    supabase = get_supabase_admin()

    # Verify client ownership
    client_result = (
        supabase.table("clients")
        .select("id,name")
        .eq("id", client_id)
        .eq("user_id", user_id)
        .single()
        .execute()
    )
    if not client_result.data:
        raise HTTPException(status_code=status.HTTP_404_NOT_FOUND, detail="Client not found")
    client_name = client_result.data["name"]

    reports_result = (
        supabase.table("reports")
        .select("*")
        .eq("client_id", client_id)
        .eq("user_id", user_id)
        .order("created_at", desc=True)
        .execute()
    )

    items = [
        ReportListItem(**_map_db_row(row, client_name=client_name))
        for row in reports_result.data
    ]
    return ReportListResponse(reports=items, total=len(items))


# ---------------------------------------------------------------------------
# GET /  — list all reports for the authenticated user
# ---------------------------------------------------------------------------

@router.get("", response_model=ReportListResponse)
async def list_all_reports(
    user_id: str = Depends(get_current_user_id),
) -> ReportListResponse:
    """List all reports generated by the authenticated user, newest first."""
    supabase = get_supabase_admin()

    reports_result = (
        supabase.table("reports")
        .select("*")
        .eq("user_id", user_id)
        .order("created_at", desc=True)
        .execute()
    )
    if not reports_result.data:
        return ReportListResponse(reports=[], total=0)

    # Fetch client names in one query
    client_ids = list({r["client_id"] for r in reports_result.data})
    clients_result = (
        supabase.table("clients")
        .select("id,name")
        .in_("id", client_ids)
        .execute()
    )
    client_map = {c["id"]: c["name"] for c in (clients_result.data or [])}

    items = [
        ReportListItem(**_map_db_row(row, client_name=client_map.get(row["client_id"])))
        for row in reports_result.data
    ]
    return ReportListResponse(reports=items, total=len(items))


# ---------------------------------------------------------------------------
# GET /{report_id}
# ---------------------------------------------------------------------------

@router.get("/{report_id}", response_model=ReportResponse)
async def get_report(
    report_id: str,
    user_id: str = Depends(get_current_user_id),
) -> ReportResponse:
    """Fetch a single report with full narrative and data summary."""
    supabase = get_supabase_admin()

    result = (
        supabase.table("reports")
        .select("*")
        .eq("id", report_id)
        .eq("user_id", user_id)
        .single()
        .execute()
    )
    if not result.data:
        raise HTTPException(status_code=status.HTTP_404_NOT_FOUND, detail="Report not found")

    row = result.data
    client_result = (
        supabase.table("clients").select("name").eq("id", row["client_id"]).single().execute()
    )
    client_name = client_result.data["name"] if client_result.data else None

    return ReportResponse(**_map_db_row(row, client_name=client_name))


# ---------------------------------------------------------------------------
# GET /{report_id}/download/pptx
# ---------------------------------------------------------------------------

@router.get("/{report_id}/download/pptx")
async def download_pptx(
    report_id: str,
    user_id: str = Depends(get_current_user_id),
) -> FileResponse:
    """Download the PowerPoint (.pptx) file for a report."""
    # Plan check: PPTX export is Pro+ only
    allowed, reason = can_use_feature(user_id, "pptx_export")
    if not allowed:
        raise HTTPException(
            status_code=status.HTTP_403_FORBIDDEN,
            detail="PPTX export is available on Pro and Agency plans. Upgrade to download PowerPoint files.",
        )

    supabase = get_supabase_admin()
    result = (
        supabase.table("reports")
        .select("user_id,title")
        .eq("id", report_id)
        .single()
        .execute()
    )
    if not result.data or result.data["user_id"] != user_id:
        raise HTTPException(status_code=status.HTTP_404_NOT_FOUND, detail="Report not found")

    pptx_path = os.path.join(REPORTS_BASE_DIR, report_id, "report.pptx")
    if not os.path.exists(pptx_path):
        raise HTTPException(
            status_code=status.HTTP_410_GONE,
            detail={
                "error": "Report files are no longer available. Please regenerate the report.",
                "code": "FILES_EXPIRED",
            },
        )

    safe_title = result.data.get("title", "report").replace(" — ", " - ").replace(" ", "_")[:80]
    return FileResponse(
        pptx_path,
        media_type="application/vnd.openxmlformats-officedocument.presentationml.presentation",
        filename=f"{safe_title}.pptx",
    )


# ---------------------------------------------------------------------------
# GET /{report_id}/download/pdf
# ---------------------------------------------------------------------------

@router.get("/{report_id}/download/pdf")
async def download_pdf(
    report_id: str,
    user_id: str = Depends(get_current_user_id),
) -> FileResponse:
    """Download the PDF file for a report."""
    # Subscription check: block expired/cancelled users
    pdf_sub = get_user_subscription(user_id)
    if pdf_sub.get("status") in ("expired", "cancelled"):
        raise HTTPException(
            status_code=status.HTTP_403_FORBIDDEN,
            detail="Your subscription has expired. Please upgrade to continue downloading reports.",
        )

    supabase = get_supabase_admin()
    result = (
        supabase.table("reports")
        .select("user_id,title")
        .eq("id", report_id)
        .single()
        .execute()
    )
    if not result.data or result.data["user_id"] != user_id:
        raise HTTPException(status_code=status.HTTP_404_NOT_FOUND, detail="Report not found")

    pdf_path = os.path.join(REPORTS_BASE_DIR, report_id, "report.pdf")
    if not os.path.exists(pdf_path):
        # Distinguish between "never had a PDF" (non-Latin) vs "files expired"
        pptx_also_missing = not os.path.exists(
            os.path.join(REPORTS_BASE_DIR, report_id, "report.pptx")
        )
        if pptx_also_missing:
            # Both files missing → ephemeral filesystem wiped after redeployment
            raise HTTPException(
                status_code=status.HTTP_410_GONE,
                detail={
                    "error": "Report files are no longer available. Please regenerate the report.",
                    "code": "FILES_EXPIRED",
                },
            )
        # PPTX exists but PDF doesn't → non-Latin language, LibreOffice wasn't available
        raise HTTPException(
            status_code=status.HTTP_404_NOT_FOUND,
            detail="PDF not available for this report. The report language may require LibreOffice for PDF rendering — please download the PPTX instead.",
        )

    safe_title = result.data.get("title", "report").replace(" — ", " - ").replace(" ", "_")[:80]
    return FileResponse(pdf_path, media_type="application/pdf", filename=f"{safe_title}.pdf")


# ---------------------------------------------------------------------------
# PATCH /{report_id}  — save manual user_edits
# ---------------------------------------------------------------------------

@router.patch("/{report_id}", response_model=ReportResponse)
async def update_report_edits(
    report_id: str,
    payload: ReportUpdateRequest,
    user_id: str = Depends(get_current_user_id),
) -> ReportResponse:
    """
    Persist manual text edits for one or more narrative sections.
    Merges the incoming user_edits dict with any existing edits in the DB,
    so editing section A doesn't wipe a previous edit to section B.
    """
    supabase = get_supabase_admin()

    # Fetch existing report (ownership check)
    existing = (
        supabase.table("reports")
        .select("*")
        .eq("id", report_id)
        .eq("user_id", user_id)
        .single()
        .execute()
    )
    if not existing.data:
        raise HTTPException(status_code=status.HTTP_404_NOT_FOUND, detail="Report not found")

    row = existing.data
    merged_edits = dict(row.get("user_edits") or {})
    merged_edits.update(payload.user_edits)

    result = (
        supabase.table("reports")
        .update({"user_edits": merged_edits})
        .eq("id", report_id)
        .eq("user_id", user_id)
        .execute()
    )
    if not result.data:
        raise HTTPException(
            status_code=status.HTTP_500_INTERNAL_SERVER_ERROR,
            detail="Failed to save edits",
        )

    client_result = (
        supabase.table("clients").select("name").eq("id", row["client_id"]).single().execute()
    )
    client_name = client_result.data["name"] if client_result.data else None
    return ReportResponse(**_map_db_row(result.data[0], client_name=client_name))


# ---------------------------------------------------------------------------
# POST /{report_id}/regenerate-section  — re-run AI for one section
# ---------------------------------------------------------------------------

@router.post("/{report_id}/regenerate-section", response_model=ReportResponse)
async def regenerate_section(
    report_id: str,
    payload: ReportSectionRegenerateRequest,
    user_id: str = Depends(get_current_user_id),
) -> ReportResponse:
    """
    Re-run GPT-4o for a single narrative section.
    Uses the compact narrative_data stored in sections JSONB at generation time.
    """
    valid_sections = {
        "executive_summary", "website_performance", "paid_advertising",
        "key_wins", "concerns", "next_steps",
    }
    if payload.section not in valid_sections:
        raise HTTPException(
            status_code=status.HTTP_422_UNPROCESSABLE_ENTITY,
            detail=f"Invalid section '{payload.section}'. Must be one of: {valid_sections}",
        )

    supabase = get_supabase_admin()

    # Fetch report + client for context
    report_result = (
        supabase.table("reports")
        .select("*")
        .eq("id", report_id)
        .eq("user_id", user_id)
        .single()
        .execute()
    )
    if not report_result.data:
        raise HTTPException(status_code=status.HTTP_404_NOT_FOUND, detail="Report not found")

    row = report_result.data
    sections_json = row.get("sections") or {}
    narrative_data = sections_json.get("narrative_data")

    if not narrative_data:
        raise HTTPException(
            status_code=status.HTTP_422_UNPROCESSABLE_ENTITY,
            detail="This report does not have stored narrative data for regeneration. "
                   "Please generate a new report.",
        )

    client_result = (
        supabase.table("clients")
        .select("name,goals_context,ai_tone,report_config")
        .eq("id", row["client_id"])
        .single()
        .execute()
    )
    if not client_result.data:
        raise HTTPException(status_code=status.HTTP_404_NOT_FOUND, detail="Client not found")

    client = client_result.data
    report_config = client.get("report_config") or {}
    cfg_template  = report_config.get("template", "monthly")

    # Re-run AI for just this one section
    from services.ai_narrative import generate_narrative  # noqa: PLC0415
    new_section_narrative = await generate_narrative(
        data=narrative_data,
        client_name=client["name"],
        client_goals=client.get("goals_context"),
        tone=client.get("ai_tone", "professional"),
        template=cfg_template,
        sections=[payload.section],
    )

    # Merge new section into existing ai_narrative
    existing_narrative = dict(row.get("ai_narrative") or {})
    existing_narrative[payload.section] = new_section_narrative.get(payload.section, "")

    # Clear any user_edit for this section (user requested a fresh AI version)
    existing_user_edits = dict(row.get("user_edits") or {})
    existing_user_edits.pop(payload.section, None)

    result = (
        supabase.table("reports")
        .update({
            "ai_narrative": existing_narrative,
            "user_edits":   existing_user_edits,
        })
        .eq("id", report_id)
        .eq("user_id", user_id)
        .execute()
    )
    if not result.data:
        raise HTTPException(
            status_code=status.HTTP_500_INTERNAL_SERVER_ERROR,
            detail="Failed to save regenerated section",
        )

    client_name = client.get("name")
    return ReportResponse(**_map_db_row(result.data[0], client_name=client_name))


# ---------------------------------------------------------------------------
# POST /{report_id}/send  — deliver report by email
# ---------------------------------------------------------------------------

@router.post("/{report_id}/send", status_code=status.HTTP_200_OK)
@limiter.limit("20/hour")
async def send_report(
    request: Request,
    report_id: str,
    payload: ReportSendRequest,
    user_id: str = Depends(get_current_user_id),
) -> dict:
    """
    Send the report to one or more email addresses via Resend.
    Attaches PDF and/or PPTX depending on payload.attachment.
    Logs the delivery attempt in the report_deliveries table.
    """
    # Subscription check: block expired/cancelled users from sending reports
    send_sub = get_user_subscription(user_id)
    if send_sub.get("status") in ("expired", "cancelled"):
        raise HTTPException(
            status_code=status.HTTP_403_FORBIDDEN,
            detail="Your subscription has expired. Please upgrade to continue sending reports.",
        )

    supabase = get_supabase_admin()

    # Fetch report
    report_result = (
        supabase.table("reports")
        .select("*")
        .eq("id", report_id)
        .eq("user_id", user_id)
        .single()
        .execute()
    )
    if not report_result.data:
        raise HTTPException(status_code=status.HTTP_404_NOT_FOUND, detail="Report not found")

    row = report_result.data

    # Fetch client for contact email fallback
    client_result = (
        supabase.table("clients")
        .select("name,primary_contact_email")
        .eq("id", row["client_id"])
        .single()
        .execute()
    )
    client = client_result.data or {}
    client_name = client.get("name", "Client")

    # Fetch agency/profile settings for sender customisation
    profile_result = (
        supabase.table("profiles")
        .select("agency_name,agency_email,sender_name,reply_to_email,email_footer")
        .eq("id", user_id)
        .single()
        .execute()
    )
    profile = profile_result.data or {}
    agency_name  = profile.get("agency_name")  or "Your Agency"
    agency_email = profile.get("agency_email") or ""
    sender_name  = payload.sender_name or profile.get("sender_name") or agency_name
    reply_to     = payload.reply_to   or profile.get("reply_to_email") or None
    email_footer = profile.get("email_footer") or ""

    # Resolve files — check if ephemeral filesystem still has them
    attach = payload.attachment.lower()
    pptx_path_resolved = os.path.join(REPORTS_BASE_DIR, report_id, "report.pptx")
    pdf_path_resolved  = os.path.join(REPORTS_BASE_DIR, report_id, "report.pdf")

    pptx_exists = os.path.exists(pptx_path_resolved)
    pdf_exists  = os.path.exists(pdf_path_resolved)

    if not pptx_exists and not pdf_exists:
        raise HTTPException(
            status_code=status.HTTP_410_GONE,
            detail={
                "error": "Report files are no longer available. Please regenerate the report.",
                "code": "FILES_EXPIRED",
            },
        )

    send_pptx = attach in ("pptx", "both") and pptx_exists
    send_pdf  = attach in ("pdf",  "both") and pdf_exists

    if not send_pptx and not send_pdf:
        raise HTTPException(
            status_code=status.HTTP_422_UNPROCESSABLE_ENTITY,
            detail="Requested file format not available for this report.",
        )

    # Build email
    sections_json = row.get("sections") or {}
    meta_currency = sections_json.get("meta_currency", "USD") if isinstance(sections_json, dict) else "USD"

    # Use user_edits-merged narrative for executive_summary snippet
    narrative = dict(row.get("ai_narrative") or {})
    user_edits = row.get("user_edits") or {}
    merged_narrative = {**narrative, **{k: v for k, v in user_edits.items() if v}}
    exec_summary = merged_narrative.get("executive_summary", "")

    title = row.get("title", f"{client_name} Performance Report")
    subject = payload.subject or title

    from services.email_service import build_report_email_html, send_report_email  # noqa: PLC0415
    html_body = build_report_email_html(
        client_name=client_name,
        period_start=str(row.get("period_start", "")),
        period_end=str(row.get("period_end", "")),
        report_title=title,
        executive_summary=exec_summary,
        agency_name=agency_name,
        agency_email=agency_email,
        email_footer=email_footer,
    )

    try:
        resend_result = await send_report_email(
            to_emails=payload.to_emails,
            subject=subject,
            html_body=html_body,
            sender_name=sender_name,
            reply_to=reply_to,
            pptx_path=pptx_path_resolved if send_pptx else None,
            pdf_path=pdf_path_resolved  if send_pdf  else None,
        )
    except Exception as exc:
        logger.error("Email send failed for report %s: %s", report_id, exc)
        # Log failed delivery
        supabase.table("report_deliveries").insert({
            "report_id":       report_id,
            "user_id":         user_id,
            "delivery_method": "email",
            "recipient_emails": payload.to_emails,
            "status":          "failed",
            "error_message":   str(exc),
            "email_subject":   subject,
        }).execute()
        raise HTTPException(
            status_code=status.HTTP_502_BAD_GATEWAY,
            detail=f"Email delivery failed: {exc}",
        )

    # Log successful delivery and update report status to "sent"
    supabase.table("report_deliveries").insert({
        "report_id":        report_id,
        "user_id":          user_id,
        "delivery_method":  "email",
        "recipient_emails": payload.to_emails,
        "status":           "sent",
        "resend_id":        resend_result.get("id"),
        "email_subject":    subject,
        "attachment_type":  payload.attachment,
        "sent_at":          datetime.utcnow().isoformat(),
    }).execute()

    supabase.table("reports").update({"status": "sent"}).eq("id", report_id).execute()

    logger.info(
        "Report %s sent to %s (Resend ID: %s)",
        report_id, payload.to_emails, resend_result.get("id"),
    )
    return {
        "success":   True,
        "resend_id": resend_result.get("id"),
        "to":        payload.to_emails,
        "subject":   subject,
    }


# ---------------------------------------------------------------------------
# POST /{report_id}/regenerate  — re-run full pipeline, reuse same report ID
# ---------------------------------------------------------------------------

@router.post("/{report_id}/regenerate", response_model=ReportResponse)
@limiter.limit("10/hour")
async def regenerate_report(
    request: Request,
    report_id: str,
    user_id: str = Depends(get_current_user_id),
) -> ReportResponse:
    """
    Re-run the full pipeline for an existing report, reusing the same report ID.

    Use case: report files were lost after a container redeployment (Railway
    storage is ephemeral) and the user clicks "Regenerate Report".

    This is now a thin wrapper. It used to be a ~480-line copy of
    _generate_report_internal, and the copy had drifted in three ways that all
    reached production:
      * CSV sources were silently dropped, so a regenerated report lost every
        uploaded data source;
      * the chart theme was hardcoded to "modern_clean" while the deck used the
        client's theme, so charts and slides disagreed;
      * white-label plan enforcement was missing, so a Starter user could
        regenerate their way to a fully branded, badge-free deck.
    All three are fixed by deleting the copy.
    """
    supabase = get_supabase_admin()

    existing = (
        supabase.table("reports")
        .select("client_id,period_start,period_end,sections")
        .eq("id", report_id)
        .eq("user_id", user_id)
        .single()
        .execute()
    )
    if not existing.data:
        raise HTTPException(status_code=status.HTTP_404_NOT_FOUND, detail="Report not found")

    row = existing.data
    client_id = str(row["client_id"])

    # Reports generated before generation_settings existed fall back to the
    # defaults - the same behaviour they had before, so no regression.
    gen_settings = (row.get("sections") or {}).get("generation_settings") or {}

    row_dict, client_name = await _generate_report_internal(
        client_id=client_id,
        user_id=user_id,
        period_start=str(row["period_start"]),
        period_end=str(row["period_end"]),
        template=gen_settings.get("template") or "full",
        visual_template=gen_settings.get("visual_template") or "modern_clean",
        csv_sources=gen_settings.get("csv_sources") or None,
        supabase=supabase,
        report_id=report_id,
    )

    logger.info("Report %s regenerated for client %s", report_id, client_id)
    return ReportResponse(**_map_db_row(row_dict, client_name=client_name))
