"""
Verify universal CSV/XLSX ingestion end to end — no network, no OpenAI calls.

Everything checked here is the deterministic half of the pipeline: profiling,
locale resolution, date disambiguation, aggregation, and the legacy-KPI path.
The AI mapping step is stubbed with a hand-written mapping, because the whole
design point is that the model only chooses column *names* — every number is
produced by the code exercised below.

    cd backend && python scripts/verify_csv_ingest.py
"""
import io
import os
import sys
import tempfile

sys.path.insert(0, os.path.dirname(os.path.dirname(os.path.abspath(__file__))))

import matplotlib  # noqa: E402

matplotlib.use("Agg")

from services.csv_ingest.normalizer import normalize, preview_rows  # noqa: E402
from services.csv_ingest.profiler import detect_date_format, profile_file  # noqa: E402
from services.csv_ingest.schema import (  # noqa: E402
    ColumnMapping,
    DateColumn,
    MappingProposal,
)
from services.csv_ingest.templates import (  # noqa: E402
    LEGACY_KPI_FINGERPRINT,
    system_mapping,
)
from services.csv_parser import parse_kpi_csv  # noqa: E402

PASSED: list[str] = []
FAILED: list[str] = []


def check(name: str, condition: bool, detail: str = "") -> None:
    if condition:
        PASSED.append(name)
        print(f"  PASS  {name}")
    else:
        FAILED.append(f"{name} — {detail}")
        print(f"  FAIL  {name} — {detail}")


# ---------------------------------------------------------------------------

def test_messy_export() -> None:
    """Title rows, a blank line, quoted thousands, %, currency, and a totals row."""
    print("\nMessy platform export")
    raw = (
        b"Campaign Performance Report\n"
        b"Date range: 01/07/2026 - 05/07/2026\n"
        b"\n"
        b"Day,Impressions,Clicks,Click Through Rate,Total Spent\n"
        b'01/07/2026,"12,400",310,2.50%,"$1,240.00"\n'
        b'02/07/2026,"11,980",298,2.49%,"$1,198.50"\n'
        b'15/07/2026,"13,220",356,2.69%,"$1,322.00"\n'
        b'30/07/2026,"12,050",301,2.50%,"$1,205.00"\n'
        b'31/07/2026,"12,900",340,2.64%,"$1,290.00"\n'
        b'Total,"62,550","1,605",2.57%,"$6,255.50"\n'
    )
    p = profile_file(raw, "linkedin.csv")[0]
    check("header row found below the preamble", p.header_row_index == 3, str(p.header_row_index))
    check("totals row detected", p.totals_row_index is not None)
    check("totals row excluded from the row count", p.data_row_count == 5, str(p.data_row_count))
    check("date column detected as DMY", p.columns[0].date_format == "%d/%m/%Y", str(p.columns[0].date_format))
    check("totals row excluded from column max", p.columns[1].max_value == 13220.0, str(p.columns[1].max_value))
    check("percent sign detected", p.columns[3].has_percent_sign)
    check("currency symbol detected", p.columns[4].has_currency_symbol)


def test_locales() -> None:
    """Locale is a property of the file, not of one column."""
    print("\nLocale handling")
    eu = (
        "Datum;Impressionen;Kosten;CTR\n"
        "01.02.2026;1.234;1.234,56;2,45\n"
        "02.02.2026;2.500;2.100,00;3,10\n"
        "15.02.2026;3.100;900,50;1,90\n"
    ).encode("utf-8")
    p = profile_file(eu, "de.csv")[0]
    check("semicolon delimiter sniffed", len(p.columns) == 4, str(len(p.columns)))
    check("European date format", p.columns[0].date_format == "%d.%m.%Y", str(p.columns[0].date_format))
    check(
        "dot-thousands column read as thousands, not decimals",
        p.columns[1].max_value == 3100.0,
        f"got {p.columns[1].max_value}, expected 3100.0",
    )
    check("comma-decimal column parsed", p.columns[2].max_value == 2100.0, str(p.columns[2].max_value))

    us = (
        b"Day,Impressions,Cost,CTR\n"
        b'2026-02-01,"1,234","1,234.56",2.45\n'
        b'2026-02-02,"2,500","2,100.00",3.10\n'
    )
    p2 = profile_file(us, "us.csv")[0]
    check("US thousands parsed", p2.columns[1].max_value == 2500.0, str(p2.columns[1].max_value))
    check("US decimals parsed", p2.columns[2].max_value == 2100.0, str(p2.columns[2].max_value))

    ratio = b"Name,Ratio\nA,1.25\nB,2.50\nC,0.75\n"
    p3 = profile_file(ratio, "r.csv")[0]
    check("a genuine decimal column is not flipped", p3.columns[1].max_value == 2.50, str(p3.columns[1].max_value))


def test_date_disambiguation() -> None:
    print("\nDate disambiguation")
    check("DMY settled by a day > 12", detect_date_format(["31/07/2026", "01/07/2026"]) == "%d/%m/%Y")
    check("MDY settled by a day > 12", detect_date_format(["12/25/2026", "01/15/2026"]) == "%m/%d/%Y")
    check("text is not a date", detect_date_format(["not a date", "also not"]) is None)
    check("numbers are not dates", detect_date_format(["1,234", "5,678"]) is None)
    check("ISO dates", detect_date_format(["2026-07-01", "2026-07-02"]) == "%Y-%m-%d")


def test_aggregation() -> None:
    """
    Counts sum over the whole period; rates recompute from their components.

    The fixture is built so the two wrong answers and the right one are all
    different numbers. Days 1-3 run at 4% CTR on light volume, days 4-6 at 6%
    on double the volume:

        summed CTR            = 30.0    (nonsense, the classic bug)
        unweighted mean CTR   =  5.0    (plausible, still wrong)
        Sigma clicks / Sigma impressions = 5.3333  (correct)

    The previous version of this test could not tell those apart: its CTR
    column said 6% on days whose own clicks and impressions worked out to 5%,
    so any method looked defensible. A self-contradictory fixture cannot fail.
    """
    print("\nAggregation")
    data = b"Day,Impressions,Clicks,CTR,Spend\n" + b"".join(
        f"2026-07-{d:02d},{1000 if d <= 3 else 2000},{40 if d <= 3 else 120},"
        f"{4.0 if d <= 3 else 6.0},{100 if d <= 3 else 200}\n".encode()
        for d in range(1, 7)
    )
    p = profile_file(data, "ads.csv")[0]
    m = MappingProposal(
        table_shape="wide_timeseries",
        source_label="Ads",
        date_column=DateColumn(name="Day", format="%Y-%m-%d", confidence=0.99),
        columns=[
            ColumnMapping(source_column="Impressions", target_metric="impressions", label="Impressions", confidence=0.98),
            ColumnMapping(source_column="Clicks", target_metric="clicks", label="Clicks", confidence=0.98),
            ColumnMapping(source_column="CTR", target_metric="ctr", label="CTR", unit="percent", confidence=0.97),
            ColumnMapping(source_column="Spend", target_metric="spend", label="Spend", unit="currency", confidence=0.99),
        ],
    )
    out = normalize(p, m)
    by = {r["metric_key"]: r for r in out["metrics"]}

    # The headline is the whole uploaded period: 3x1000 + 3x2000.
    check("counts sum over the whole period",
          by["impressions"]["current_value"] == 9000.0, str(by["impressions"]))
    check("halves retained for the trend",
          by["impressions"]["first_half_value"] == 3000.0
          and by["impressions"]["second_half_value"] == 6000.0,
          str(by["impressions"]))
    # One upload is one period; there is no prior period to compare against.
    check("no invented previous period",
          by["impressions"]["previous_value"] is None, str(by["impressions"]))
    check("change is the within-period trend",
          by["impressions"]["change"] == 100.0
          and by["impressions"]["change_basis"] == "within_period",
          str(by["impressions"]))

    ctr = by["ctr"]["current_value"]
    check("rate recomputed from components, not summed",
          abs(ctr - 5.3333) < 0.001, f"{ctr} (summed would be 30.0)")
    check("rate is volume-weighted, not an unweighted mean",
          abs(ctr - 5.0) > 0.001, f"{ctr} == the unweighted mean of the CTR column")
    check("rate reports how it was derived",
          out["derivations"]["ctr"]["method"] == "recomputed",
          str(out["derivations"]["ctr"]))

    check("daily series emitted", len(out.get("daily", [])) == 6, str(len(out.get("daily", []))))

    undated = b"Channel,Sessions\nOrganic,500\nPaid,300\n"
    p2 = profile_file(undated, "x.csv")[0]
    m2 = MappingProposal(
        source_label="Analytics",
        columns=[ColumnMapping(source_column="Sessions", target_metric="sessions", label="Sessions", confidence=0.95)],
    )
    o2 = normalize(p2, m2)
    check("undated file totals without comparison", o2["metrics"][0]["current_value"] == 800.0, str(o2["metrics"][0]))
    check("undated file emits no series", "daily" not in o2)


def test_legacy_parity() -> None:
    """The five bundled templates must produce identical values on the new path."""
    print("\nLegacy KPI template parity")
    legacy = (
        b"metric_name,current_value,previous_value,unit\n"
        b"Impressions,45200,38900,number\n"
        b"Clicks,1340,1100,number\n"
        b"Click-Through Rate,2.96,2.83,percent\n"
        b"Spend,1850.00,1600.00,currency\n"
    )
    p = profile_file(legacy, "linkedin_ads_template.csv")[0]
    check("legacy layout matches the seeded system fingerprint",
          p.column_fingerprint == LEGACY_KPI_FINGERPRINT)

    new = normalize(p, system_mapping("linkedin_ads"), source_name="LinkedIn Ads")
    old = parse_kpi_csv(legacy, "linkedin_ads_template.csv")
    new_by = {r["name"]: r for r in new["metrics"]}

    check("same metric count", len(new["metrics"]) == len(old["metrics"]),
          f"{len(new['metrics'])} vs {len(old['metrics'])}")
    mismatches = []
    for row in old["metrics"]:
        fresh = new_by.get(row["name"])
        if not fresh:
            mismatches.append(f"missing {row['name']}")
            continue
        if fresh["current_value"] != row["current_value"]:
            mismatches.append(f"{row['name']} current {fresh['current_value']} != {row['current_value']}")
        if fresh["previous_value"] != row["previous_value"]:
            mismatches.append(f"{row['name']} previous {fresh['previous_value']} != {row['previous_value']}")
        if fresh["unit"] != row["unit"]:
            mismatches.append(f"{row['name']} unit {fresh['unit']} != {row['unit']}")
    check("every value matches the legacy parser", not mismatches, "; ".join(mismatches))

    eu = (
        "metric_name;current_value;previous_value;unit\n"
        "Umsatz;1.234,56;1.100,00;currency\n"
        "Klicks;2.500;2.100;number\n"
    ).encode("utf-8")
    p2 = profile_file(eu, "de_kpis.csv")[0]
    o2 = normalize(p2, system_mapping("generic"), source_name="DE KPIs")
    check("European long-KPI currency", o2["metrics"][0]["current_value"] == 1234.56, str(o2["metrics"][0]))
    check("European long-KPI thousands", o2["metrics"][1]["current_value"] == 2500.0, str(o2["metrics"][1]))


def test_xlsx() -> None:
    print("\nExcel workbooks")
    try:
        import openpyxl
    except ImportError:
        check("openpyxl installed", False, "pip install openpyxl")
        return

    from datetime import date

    wb = openpyxl.Workbook()
    readme = wb.active
    readme.title = "Read Me"
    readme.append(["Exported from AdPlatform"])
    detail = wb.create_sheet("Campaign Detail")
    detail.append(["Campaign", "Impressions", "Clicks", "Spend"])
    for i in range(1, 26):
        detail.append([f"Campaign {i}", 1000 + i * 10, 40 + i, 100 + i])
    buf = io.BytesIO()
    wb.save(buf)

    profiles = profile_file(buf.getvalue(), "export.xlsx")
    check("largest sheet ranked first", profiles[0].sheet_name == "Campaign Detail", profiles[0].sheet_name)
    check("all non-empty sheets profiled", len(profiles) == 2, str(len(profiles)))

    wb2 = openpyxl.Workbook()
    ws = wb2.active
    ws.append(["Day", "Clicks"])
    for d in range(1, 11):
        ws.append([date(2026, 7, d), 10 + d])
    buf2 = io.BytesIO()
    wb2.save(buf2)
    p2 = profile_file(buf2.getvalue(), "dated.xlsx")[0]
    check("native Excel dates recognised", p2.columns[0].inferred_type == "date", p2.columns[0].inferred_type)


def test_render() -> None:
    """A mapped, dated upload must reach a rendered deck with no blank tokens."""
    print("\nRender")
    from services.chart_generator import generate_all_charts
    from services.report_generator import generate_pptx_report

    data_csv = b"Day,Impressions,Clicks,Spend\n" + b"".join(
        f"2026-07-{d:02d},{1000 + d * 20},{50 + d},{100 + d}\n".encode() for d in range(1, 29)
    )
    p = profile_file(data_csv, "linkedin.csv")[0]
    m = MappingProposal(
        table_shape="wide_timeseries",
        source_label="LinkedIn Ads",
        date_column=DateColumn(name="Day", format="%Y-%m-%d", confidence=0.99),
        columns=[
            ColumnMapping(source_column="Impressions", target_metric="impressions", label="Impressions", confidence=0.98),
            ColumnMapping(source_column="Clicks", target_metric="clicks", label="Clicks", confidence=0.98),
            ColumnMapping(source_column="Spend", target_metric="spend", label="Spend", unit="currency", confidence=0.99),
        ],
    )
    source = normalize(p, m, source_name="LinkedIn Ads")
    data = {
        "client_name": "Acme Ltd",
        "period_start": "2026-07-01",
        "period_end": "2026-07-28",
        "csv_sources": [source],
    }

    out_dir = tempfile.mkdtemp()
    charts = generate_all_charts(data, out_dir, "#4338CA", "modern_clean")
    check("dated upload produces a trend chart", "csv_linkedin_ads" in charts, str(list(charts)))

    pptx = generate_pptx_report(
        data,
        {"executive_summary": "Summary.", "csv_performance": "LinkedIn delivered.",
         "key_wins": ["Leads up"], "concerns": ["CPC rising"], "next_steps": ["Scale winners"]},
        charts,
        {"name": "Acme Ltd", "agency_name": "SapienBotics"},
        None, "full", None,
        {"agency_name": "SapienBotics", "brand_color": "#4338CA", "agency_logo_url": "",
         "client_logo_url": "", "powered_by_badge": False, "_cover_theme": "modern_clean"},
        "modern_clean", "en",
        {"theme": "modern_clean", "headline": "Acme Ltd", "subtitle": None,
         "brand_primary_color": None, "accent_color": None},
    )
    path = os.path.join(out_dir, "report.pptx")
    with open(path, "wb") as handle:
        handle.write(pptx)

    from pptx import Presentation

    prs = Presentation(path)
    texts = [sh.text_frame.text for s in prs.slides for sh in s.shapes if sh.has_text_frame]
    check("no unreplaced tokens", not [t for t in texts if "{{" in t],
          str([t for t in texts if "{{" in t][:3]))
    check("source name reaches the deck", any("LinkedIn" in t for t in texts))
    check("no blank slides", len(prs.slides) > 0 and len(prs.slides) < 19, str(len(prs.slides)))


def test_failure_messages() -> None:
    """Every failure must say what to do next, never 'parse error'."""
    print("\nFailure messages")
    from services.csv_ingest.profiler import IngestError

    cases = [
        (b"", "empty.csv"),
        (b"%PDF-1.4 junk", "report.pdf"),
        (b"\xd0\xcf\x11\xe0junk", "old.xls"),
    ]
    for content, name in cases:
        try:
            profile_file(content, name)
            check(f"{name} rejected", False, "no error raised")
        except IngestError as exc:
            message = str(exc)
            lowered = message.lower()
            check(
                f"{name} rejected with actionable guidance",
                len(message) > 30 and any(
                    verb in lowered
                    for verb in ("try", "export", "save", "upload", "convert", "open it")
                ),
                message,
            )
        except Exception as exc:  # noqa: BLE001
            check(f"{name} rejected cleanly", False, f"{type(exc).__name__}: {exc}")


def main() -> int:
    print("Universal CSV ingestion verification")
    print("=" * 60)
    test_messy_export()
    test_locales()
    test_date_disambiguation()
    test_aggregation()
    test_legacy_parity()
    test_xlsx()
    test_render()
    test_failure_messages()

    print("\n" + "=" * 60)
    print(f"{len(PASSED)} passed, {len(FAILED)} failed")
    for failure in FAILED:
        print(f"  FAILED: {failure}")
    return 1 if FAILED else 0


if __name__ == "__main__":
    sys.exit(main())
