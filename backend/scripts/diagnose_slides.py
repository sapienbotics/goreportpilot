"""Detailed slide content dump of a generated report."""
import os, sys, tempfile, io
sys.stdout = io.TextIOWrapper(sys.stdout.buffer, encoding='utf-8')
sys.path.insert(0, os.path.join(os.path.dirname(os.path.dirname(os.path.abspath(__file__)))))

from pptx import Presentation
from services.report_generator import generate_pptx_report
from services.chart_generator import generate_all_charts
from services.mock_data import generate_all_mock_data

raw_data = generate_all_mock_data("Test Client", "2026-03-01", "2026-03-30")
raw_data["csv_sources"] = [
    {"source_name": "TikTok Ads", "metrics": [
        {"name": "Impressions", "current_value": 320000, "previous_value": 275000, "unit": "number", "change": 16.36},
        {"name": "Clicks", "current_value": 9800, "previous_value": 8100, "unit": "number", "change": 20.99},
    ]},
    {"source_name": "LinkedIn Ads", "metrics": [
        {"name": "Impressions", "current_value": 45200, "previous_value": 38900, "unit": "number", "change": 16.2},
        {"name": "Spend", "current_value": 1850, "previous_value": 1600, "unit": "currency", "change": 15.63},
    ]},
    {"source_name": "Mailchimp", "metrics": [
        {"name": "Emails Sent", "current_value": 12500, "previous_value": 11800, "unit": "number", "change": 5.93},
        {"name": "Open Rate", "current_value": 28.4, "previous_value": 26.1, "unit": "percent", "change": 8.81},
    ]},
]

with tempfile.TemporaryDirectory() as tmpdir:
    charts = generate_all_charts(raw_data, os.path.join(tmpdir, "charts"), "#4338CA", "modern_clean")
    narrative = {
        "executive_summary": "Test summary.",
        "website_performance": "Website OK.",
        "paid_advertising": "Ads OK.",
        "key_wins": "Win 1\nWin 2",
        "concerns": "Issue 1\nIssue 2",
        "next_steps": "Step 1\nStep 2",
    }
    pptx_bytes = generate_pptx_report(
        data=raw_data, narrative=narrative, charts=charts,
        client_info={"name": "Test Client", "agency_name": "Test Agency"},
        visual_template="modern_clean",
    )
    path = os.path.join(tmpdir, "test.pptx")
    with open(path, "wb") as f:
        f.write(pptx_bytes)

    prs = Presentation(path)
    print(f"Total slides: {len(prs.slides)}\n")
    for idx, slide in enumerate(prs.slides):
        texts = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                t = shape.text_frame.text.strip()
                if t:
                    texts.append(t[:60])
        # Show first 3 text blocks
        summary = " | ".join(texts[:3]) if texts else "(no text)"
        print(f"  {idx:2d}: {summary}")
