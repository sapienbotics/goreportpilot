"""
Regenerate chrome-only cover thumbnails for the 6 design-system templates.

Phase 3 parity fix — D1-B + D2-B combined.

What it does
------------
For each template under ``backend/templates/pptx/``:

  1. Loads the PPTX and walks the cover slide (``slides[0]``).
  2. Deletes every text shape whose text contains a ``{{placeholder}}``
     token. That removes, in one pass:

        * ``{{client_name}}``, ``{{report_period}}``, ``{{report_type}}``
        * ``{{agency_name}}`` and ``Prepared by {{agency_name}}``
        * ``{{agency_logo}}`` and ``{{client_logo}}`` placeholder boxes

     Decorative chrome is kept — header bands, divider lines, gradient
     strips, and any static label text (e.g. "PERFORMANCE REPORT").
  3. Saves the stripped template back in place. The master is now
     **chrome-only** on the cover. Content slides 2-19 are untouched.
  4. Builds a one-slide scratch copy of the template (cover only) and
     converts it to PNG via LibreOffice CLI. The PNG replaces the old
     thumbnail under ``backend/static/cover_thumbnails/``.

After this script runs, text + logos on the cover come exclusively
from ``cover_customization.py`` at report-generate time, which draws
them at ``theme_layout`` coordinates. The thumbnail now shows only
the design language (band, gradient, geometric block, whitespace) so
the Design tab's CSS overlay renders cleanly with no double-text.

Invocation (Windows)
--------------------

    python backend/scripts/regenerate_cover_thumbnails.py

Set ``SOFFICE_PATH`` if your LibreOffice install lives somewhere else.
"""
from __future__ import annotations

import os
import shutil
import subprocess
import sys
import tempfile
from pathlib import Path

from pptx import Presentation

# ─── Config ────────────────────────────────────────────────────────────────────
BACKEND       = Path(__file__).resolve().parents[1]
TEMPLATES_DIR = BACKEND / "templates" / "pptx"
THUMBS_DIR    = BACKEND / "static"    / "cover_thumbnails"

SOFFICE = os.environ.get("SOFFICE_PATH") or (
    r"C:/Program Files/LibreOffice/program/soffice.com"  # Windows default
    if sys.platform == "win32"
    else "soffice"
)

THEMES = (
    "modern_clean",
    "dark_executive",
    "colorful_agency",
    "bold_geometric",
    "minimal_elegant",
    "gradient_modern",
)


# ─── Strip step ───────────────────────────────────────────────────────────────


def _is_placeholder_text(txt: str) -> bool:
    """Return True iff the shape's text carries a ``{{...}}`` token."""
    return bool(txt) and "{{" in txt and "}}" in txt


def strip_cover_placeholders(pptx_path: Path) -> list[str]:
    """
    Delete placeholder text shapes from the cover slide of ``pptx_path``
    (overwriting the file). Returns the list of removed shape names for
    logging.
    """
    prs = Presentation(str(pptx_path))
    cover = prs.slides[0]

    removed: list[str] = []
    to_remove = []
    for shape in list(cover.shapes):
        if not shape.has_text_frame:
            continue
        txt = shape.text_frame.text or ""
        if _is_placeholder_text(txt):
            to_remove.append(shape)
            removed.append(
                f"{shape.name!r} text={txt.strip().replace(chr(10), ' | ')[:50]!r}"
            )

    for shape in to_remove:
        sp = shape._element
        sp.getparent().remove(sp)

    prs.save(str(pptx_path))
    return removed


# ─── Thumbnail step ───────────────────────────────────────────────────────────


def _extract_cover_only(src_pptx: Path, dst_pptx: Path) -> None:
    """Copy ``src`` → ``dst`` then delete every slide except index 0."""
    shutil.copyfile(src_pptx, dst_pptx)
    prs = Presentation(str(dst_pptx))

    # sldIdLst holds the slide references in presentation order.
    xml_slides = prs.slides._sldIdLst  # pylint: disable=protected-access
    slides_to_remove = list(xml_slides)[1:]
    for s in slides_to_remove:
        xml_slides.remove(s)

    prs.save(str(dst_pptx))


def _convert_to_png(pptx_path: Path, out_dir: Path) -> Path:
    """Run ``soffice --convert-to png`` and return the produced PNG path."""
    cmd = [
        SOFFICE, "--headless",
        "--convert-to", "png",
        "--outdir", str(out_dir),
        str(pptx_path),
    ]
    result = subprocess.run(cmd, check=False, capture_output=True, text=True)
    if result.returncode != 0:
        raise RuntimeError(
            f"LibreOffice convert failed (rc={result.returncode}).\n"
            f"stdout: {result.stdout}\nstderr: {result.stderr}"
        )

    produced = out_dir / (pptx_path.stem + ".png")
    if not produced.exists():
        raise RuntimeError(
            f"LibreOffice did not produce expected PNG: {produced}.\n"
            f"stdout: {result.stdout}\nstderr: {result.stderr}"
        )
    return produced


def regenerate_thumbnail(theme: str, tmp_dir: Path) -> Path:
    """Build a cover-only scratch PPTX for ``theme`` and write its PNG."""
    src = TEMPLATES_DIR / f"{theme}.pptx"
    scratch = tmp_dir / f"{theme}_cover_only.pptx"
    _extract_cover_only(src, scratch)
    produced_png = _convert_to_png(scratch, tmp_dir)
    dst_png = THUMBS_DIR / f"{theme}.png"
    shutil.move(str(produced_png), str(dst_png))
    return dst_png


# ─── Driver ───────────────────────────────────────────────────────────────────


def main() -> int:
    if not TEMPLATES_DIR.exists():
        print(f"[ERR] Templates dir not found: {TEMPLATES_DIR}", file=sys.stderr)
        return 2
    THUMBS_DIR.mkdir(parents=True, exist_ok=True)

    # Probe LibreOffice once before touching files.
    probe = subprocess.run([SOFFICE, "--version"], check=False,
                           capture_output=True, text=True)
    if probe.returncode != 0:
        print(f"[ERR] soffice not runnable at {SOFFICE!r}\n"
              f"stderr: {probe.stderr}", file=sys.stderr)
        return 2
    print(f"LibreOffice: {probe.stdout.strip()}")
    print(f"Templates  : {TEMPLATES_DIR}")
    print(f"Thumbs     : {THUMBS_DIR}")
    print()

    with tempfile.TemporaryDirectory() as tmp:
        tmp_path = Path(tmp)
        for theme in THEMES:
            src = TEMPLATES_DIR / f"{theme}.pptx"
            print(f"▸ {theme}")
            if not src.exists():
                print(f"  [skip] missing {src}")
                continue

            # 1. Strip placeholder shapes (overwrites the master).
            removed = strip_cover_placeholders(src)
            if removed:
                print(f"  stripped {len(removed)} placeholder shape(s):")
                for line in removed:
                    print(f"    · {line}")
            else:
                print("  (no placeholder shapes found — template already clean)")

            # 2. Regenerate the PNG thumbnail from the cleaned master.
            png = regenerate_thumbnail(theme, tmp_path)
            size_kb = png.stat().st_size / 1024
            print(f"  wrote {png.relative_to(BACKEND)} ({size_kb:.1f} KB)")
            print()

    print("Done. Review the templates + thumbnails and commit.")
    return 0


if __name__ == "__main__":
    sys.exit(main())
