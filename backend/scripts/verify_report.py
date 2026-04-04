"""
End-to-end verification: generate a test report with 3 CSV sources,
then audit the output PPTX to verify no shape overlaps on CSV slides.
"""
import os
import sys
import io
import tempfile

sys.stdout = io.TextIOWrapper(sys.stdout.buffer, encoding='utf-8')
# Add backend to path
sys.path.insert(0, os.path.join(os.path.dirname(os.path.dirname(os.path.abspath(__file__)))))

from pptx import Presentation
from pptx.util import Inches

EMU = 914400


def audit_csv_slides(pptx_path: str) -> list[dict]:
    """Audit all CSV-like slides in a generated report."""
    prs = Presentation(pptx_path)
    results = []

    for idx, slide in enumerate(prs.slides):
        # Detect CSV slides: look for source-name-like headers and csv_kpi patterns
        has_csv_kpi = False
        has_chart_img = False
        kpi_bottoms = []
        chart_tops = []
        footer_tops = []
        slide_title = ""

        for shape in slide.shapes:
            t = shape.top / EMU if shape.top else 0
            b = (shape.top + shape.height) / EMU if shape.top and shape.height else 0

            if shape.has_text_frame:
                text = shape.text_frame.text.strip()
                if not slide_title and text and "{{" not in text:
                    slide_title = text[:40]

            # Detect KPI values: shapes with text in the ~1.5–5.0" range
            if shape.has_text_frame:
                text = shape.text_frame.text.strip()
                if text and 1.0 < t < 5.0 and not any(kw in text.lower() for kw in ["agency", "confidential", "page"]):
                    kpi_bottoms.append(b)

            # Detect images (chart embeds)
            if shape.shape_type is not None:
                try:
                    if shape.image:
                        has_chart_img = True
                        chart_tops.append(t)
                except Exception:
                    pass

            # Footer detection
            if shape.has_text_frame:
                text = shape.text_frame.text.strip().lower()
                if "confidential" in text or ("page" in text and "agency" in text.lower()):
                    footer_tops.append(t)

        if not kpi_bottoms or not chart_tops:
            continue

        max_kpi_bottom = max(kpi_bottoms)
        min_chart_top = min(chart_tops)
        gap = min_chart_top - max_kpi_bottom

        result = {
            "slide": idx,
            "title": slide_title,
            "max_kpi_bottom": round(max_kpi_bottom, 2),
            "chart_top": round(min_chart_top, 2),
            "gap": round(gap, 2),
            "overlap": gap < 0,
        }
        results.append(result)

    return results


def main():
    from services.report_generator import generate_pptx_report
    from services.chart_generator import generate_all_charts
    from services.mock_data import generate_all_mock_data

    # Build test data with 3 CSV sources
    raw_data = generate_all_mock_data("Test Client", "2026-03-01", "2026-03-30")
    raw_data["csv_sources"] = [
        {
            "source_name": "TikTok Ads",
            "metrics": [
                {"name": "Impressions", "current_value": 320000, "previous_value": 275000, "unit": "number", "change": 16.36},
                {"name": "Clicks", "current_value": 9800, "previous_value": 8100, "unit": "number", "change": 20.99},
                {"name": "CTR", "current_value": 3.06, "previous_value": 2.95, "unit": "percent", "change": 3.73},
                {"name": "Spend", "current_value": 2400, "previous_value": 2100, "unit": "currency", "change": 14.29},
                {"name": "CPC", "current_value": 0.24, "previous_value": 0.26, "unit": "currency", "change": -7.69},
                {"name": "Video Views", "current_value": 185000, "previous_value": 160000, "unit": "number", "change": 15.63},
            ],
        },
        {
            "source_name": "LinkedIn Ads",
            "metrics": [
                {"name": "Impressions", "current_value": 45200, "previous_value": 38900, "unit": "number", "change": 16.2},
                {"name": "Clicks", "current_value": 1340, "previous_value": 1100, "unit": "number", "change": 21.82},
                {"name": "CTR", "current_value": 2.96, "previous_value": 2.83, "unit": "percent", "change": 4.59},
                {"name": "Spend", "current_value": 1850, "previous_value": 1600, "unit": "currency", "change": 15.63},
                {"name": "CPC", "current_value": 1.38, "previous_value": 1.45, "unit": "currency", "change": -4.83},
                {"name": "Leads", "current_value": 87, "previous_value": 71, "unit": "number", "change": 22.54},
            ],
        },
        {
            "source_name": "Mailchimp",
            "metrics": [
                {"name": "Emails Sent", "current_value": 12500, "previous_value": 11800, "unit": "number", "change": 5.93},
                {"name": "Open Rate", "current_value": 28.4, "previous_value": 26.1, "unit": "percent", "change": 8.81},
                {"name": "Click Rate", "current_value": 4.7, "previous_value": 4.2, "unit": "percent", "change": 11.9},
                {"name": "Revenue", "current_value": 3200, "previous_value": 2750, "unit": "currency", "change": 16.36},
                {"name": "New Subscribers", "current_value": 340, "previous_value": 290, "unit": "number", "change": 17.24},
                {"name": "List Size", "current_value": 18450, "previous_value": 18110, "unit": "number", "change": 1.88},
            ],
        },
    ]

    with tempfile.TemporaryDirectory() as tmpdir:
        charts_dir = os.path.join(tmpdir, "charts")
        charts = generate_all_charts(raw_data, charts_dir, "#4338CA", "modern_clean")

        client_info = {"name": "Test Client", "agency_name": "Test Agency"}
        narrative = {
            "executive_summary": "Test summary for verification.",
            "website_performance": "Website metrics looked strong.",
            "paid_advertising": "Paid ads delivered solid ROAS.",
            "key_wins": "Strong engagement\nGood conversions",
            "concerns": "Rising CPC\nBounce rate trending up",
            "next_steps": "Optimize landing pages\nAdjust bid strategy",
        }

        pptx_bytes = generate_pptx_report(
            data=raw_data,
            narrative=narrative,
            charts=charts,
            client_info=client_info,
            visual_template="modern_clean",
        )

        out_path = os.path.join(tmpdir, "test_report.pptx")
        with open(out_path, "wb") as f:
            f.write(pptx_bytes)

        prs = Presentation(out_path)
        print(f"\nGenerated report: {len(prs.slides)} slides, {len(pptx_bytes):,} bytes")

        # Count CSV slides: any slide containing a known CSV source name
        csv_source_names = {"TikTok Ads", "LinkedIn Ads", "Mailchimp"}
        csv_slide_count = 0
        for slide in prs.slides:
            for shape in slide.shapes:
                if shape.has_text_frame:
                    text = shape.text_frame.text.strip()
                    if text in csv_source_names:
                        csv_slide_count += 1
                        break

        print(f"CSV slides found: {csv_slide_count} (expected 3)")

        # Run audit
        results = audit_csv_slides(out_path)
        if results:
            print(f"\nSlide overlap audit ({len(results)} slides with KPIs + charts):")
            any_overlap = False
            for r in results:
                status = "OVERLAP" if r["overlap"] else "OK"
                if r["overlap"]:
                    any_overlap = True
                print(f"  Slide {r['slide']:2d} | {r['title'][:30]:30s} | "
                      f"KPI bottom={r['max_kpi_bottom']:.2f}\" | "
                      f"Chart top={r['chart_top']:.2f}\" | "
                      f"Gap={r['gap']:+.2f}\" | {status}")

            if any_overlap:
                print("\n*** FAIL: overlapping shapes detected ***")
            else:
                print("\n*** PASS: no overlaps ***")
        else:
            print("\nNo slides with both KPIs and charts found to audit.")

        # Verify slide ordering: CSV slides should come before Key Wins/Concerns/Next Steps
        slide_labels = []
        for slide in prs.slides:
            label = ""
            for shape in slide.shapes:
                if shape.has_text_frame and shape.text_frame.text.strip():
                    label = shape.text_frame.text.strip()[:40]
                    break
            slide_labels.append(label)

        print(f"\nSlide order:")
        for i, label in enumerate(slide_labels):
            print(f"  {i:2d}: {label}")


if __name__ == "__main__":
    main()
