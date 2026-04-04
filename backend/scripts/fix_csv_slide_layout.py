"""
Fix CSV slide layout in ALL 6 PPTX templates.

For each template:
1. Find the CSV data slide (contains {{csv_source_name}})
2. Calculate max bottom edge of all csv_kpi_* shapes
3. Move {{chart_csv_data}} placeholder to: max_kpi_bottom + 0.20" clearance
4. Resize chart height to fill space up to footer - 0.15" margin
5. Also move any background shape that matches the chart's original position/size
"""
from pptx import Presentation
from pptx.util import Inches, Emu
import glob
import os

TPL_DIR = os.path.join(os.path.dirname(os.path.dirname(os.path.abspath(__file__))),
                       "templates", "pptx")

CLEARANCE = Inches(0.20)       # gap between KPI bottom and chart top
FOOTER_MARGIN = Inches(0.15)   # gap between chart bottom and footer top
MIN_CHART_H = Inches(1.5)      # minimum readable chart height

EMU_PER_INCH = 914400

for tpl_path in sorted(glob.glob(os.path.join(TPL_DIR, "*.pptx"))):
    basename = os.path.basename(tpl_path)
    prs = Presentation(tpl_path)

    # Find the CSV slide
    csv_slide = None
    csv_slide_idx = None
    for idx, slide in enumerate(prs.slides):
        for shape in slide.shapes:
            if shape.has_text_frame and "{{csv_source_name}}" in shape.text_frame.text:
                csv_slide = slide
                csv_slide_idx = idx
                break
        if csv_slide:
            break

    if csv_slide is None:
        print(f"{basename}: NO CSV SLIDE — skipping")
        continue

    # Collect positions
    kpi_bottoms = []
    chart_shape = None
    chart_orig_left = None
    chart_orig_top = None
    chart_orig_width = None
    chart_orig_height = None
    footer_top = None

    for shape in csv_slide.shapes:
        if not shape.has_text_frame:
            continue
        text = shape.text_frame.text.strip()

        if "csv_kpi_" in text:
            kpi_bottoms.append(shape.top + shape.height)
        elif "chart_csv" in text:
            chart_shape = shape
            chart_orig_left = shape.left
            chart_orig_top = shape.top
            chart_orig_width = shape.width
            chart_orig_height = shape.height
        elif ("agency_name" in text.lower() and
              ("page" in text.lower() or "confidential" in text.lower())):
            footer_top = shape.top

    if not kpi_bottoms or chart_shape is None:
        print(f"{basename}: Missing KPI shapes or chart placeholder — skipping")
        continue

    max_kpi_bottom = max(kpi_bottoms)
    new_top = max_kpi_bottom + CLEARANCE

    if footer_top:
        new_bottom = footer_top - FOOTER_MARGIN
    else:
        new_bottom = prs.slide_height - Inches(0.5)

    new_height = max(new_bottom - new_top, MIN_CHART_H)

    # Report
    old_top_in = round(chart_orig_top / EMU_PER_INCH, 2)
    old_h_in = round(chart_orig_height / EMU_PER_INCH, 2)
    new_top_in = round(new_top / EMU_PER_INCH, 2)
    new_h_in = round(new_height / EMU_PER_INCH, 2)
    kpi_b_in = round(max_kpi_bottom / EMU_PER_INCH, 2)
    gap_in = round((new_top - max_kpi_bottom) / EMU_PER_INCH, 2)

    print(f"{basename}  slide {csv_slide_idx}:")
    print(f"  KPI max bottom:  {kpi_b_in}\"")
    print(f"  Chart old:       top={old_top_in}\" height={old_h_in}\"")
    print(f"  Chart new:       top={new_top_in}\" height={new_h_in}\" (gap={gap_in}\")")

    # Update chart placeholder
    chart_shape.top = int(new_top)
    chart_shape.height = int(new_height)

    # Update any background shape at the exact original chart position/size
    bg_count = 0
    for shape in csv_slide.shapes:
        if shape is chart_shape:
            continue
        if (shape.left == chart_orig_left and
            shape.top == chart_orig_top and
            shape.width == chart_orig_width and
            shape.height == chart_orig_height):
            shape.top = int(new_top)
            shape.height = int(new_height)
            bg_count += 1
            print(f"  Updated background shape id={shape.shape_id}")

    prs.save(tpl_path)
    print(f"  Saved ({bg_count} bg shape{'s' if bg_count != 1 else ''} updated)\n")
