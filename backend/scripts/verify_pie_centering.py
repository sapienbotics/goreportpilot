"""Verify: when device_breakdown exists but top_pages doesn't, pie chart stays original width and centers."""
import os, sys, io, tempfile
sys.stdout = io.TextIOWrapper(sys.stdout.buffer, encoding='utf-8')
sys.path.insert(0, os.path.join(os.path.dirname(os.path.dirname(os.path.abspath(__file__)))))

from pptx import Presentation
EMU = 914400

from services.report_generator import generate_pptx_report
from services.chart_generator import generate_all_charts
from services.mock_data import generate_all_mock_data

raw_data = generate_all_mock_data("Test Client", "2026-03-01", "2026-03-30")
# Keep device_breakdown but remove top_pages so pie chart is solo
raw_data["ga4"]["top_pages"] = []

with tempfile.TemporaryDirectory() as tmpdir:
    charts = generate_all_charts(raw_data, os.path.join(tmpdir, "charts"), "#4338CA", "modern_clean")
    narrative = {
        "executive_summary": "S", "website_performance": "W",
        "paid_advertising": "A", "key_wins": "K", "concerns": "C", "next_steps": "N",
    }
    pptx_bytes = generate_pptx_report(
        data=raw_data, narrative=narrative, charts=charts,
        client_info={"name": "Test", "agency_name": "Agency"},
        visual_template="modern_clean",
    )
    path = os.path.join(tmpdir, "test.pptx")
    with open(path, "wb") as f:
        f.write(pptx_bytes)

    prs = Presentation(path)
    sw = prs.slide_width / EMU

    for idx, slide in enumerate(prs.slides):
        for shape in slide.shapes:
            if shape.has_text_frame and "Engagement" in shape.text_frame.text:
                print(f"Website Engagement slide: {idx}")
                for s2 in slide.shapes:
                    try:
                        if s2.image:
                            w = s2.width / EMU
                            l = s2.left / EMU
                            center_off = abs(l + w / 2 - sw / 2)
                            print(f"  Pie chart: width={w:.2f}\" left={l:.2f}\" "
                                  f"centered={center_off < 0.15} "
                                  f"{'PASS (≤6\")' if w <= 6.0 and center_off < 0.15 else 'FAIL'}")
                    except Exception:
                        pass
                break
