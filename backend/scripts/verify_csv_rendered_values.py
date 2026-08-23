"""
Assert what the DECK SHOWS, not what the code computes.

This script exists because of a specific failure. Pass 4 verified the
LinkedIn fixture's CTR by recomputing it in Python — Sigma clicks over Sigma
impressions, 0.6952% — compared that against the expected 0.70%, and
reported PASS. The deck being verified displayed 1.04%. The assertion was
correct and measured the wrong thing: it re-derived the expected value
alongside the code under test instead of reading the artifact that ships.

So every assertion below opens the generated .pptx, reads the text actually
rendered onto the CSV KPI slide, and compares THAT. If the renderer formats a
number wrongly, drops its decimals, or labels dollars as rupees, these fail —
none of which a recomputation can see.

    cd backend && python scripts/verify_csv_rendered_values.py
    cd backend && python scripts/verify_csv_rendered_values.py --live

--live re-runs the real GPT-4.1 column mapping instead of the recorded one.
The recorded mapping is the live model's actual output for this fixture; the
numbers are deterministic either way, since the model only names columns.
"""
from __future__ import annotations

import argparse
import os
import re
import sys
import tempfile

sys.path.insert(0, os.path.dirname(os.path.dirname(os.path.abspath(__file__))))

import matplotlib  # noqa: E402

matplotlib.use("Agg")

from pptx import Presentation  # noqa: E402

from services.csv_ingest.normalizer import normalize  # noqa: E402
from services.csv_ingest.profiler import profile_file  # noqa: E402
from services.csv_ingest.schema import (  # noqa: E402
    ColumnMapping,
    DateColumn,
    EntityColumn,
    IgnoredColumn,
    MappingProposal,
)

FIXTURE = os.path.join(
    os.path.dirname(os.path.dirname(os.path.abspath(__file__))),
    "tests", "fixtures", "csv_ingest",
    "linkedin_ads_campaign_performance_july2026.csv",
)
FIXTURE_MD5 = "f6618db3ad18071a9264d1f1afcfafef"
TEMPLATE = os.path.join(
    os.path.dirname(os.path.dirname(os.path.abspath(__file__))),
    "templates", "pptx", "modern_clean.pptx",
)

PASSED: list[str] = []
FAILED: list[str] = []

# Failure details quote rendered cell text, which is exactly where a stray '₹'
# or '€' turns up. On a cp1252 console printing one raises UnicodeEncodeError
# and a legible FAIL becomes a traceback — the currency mutation test hit this.
# A verification tool that crashes on the evidence it was built to surface is
# worse than useless, so stdout is made lossy rather than fatal.
if hasattr(sys.stdout, "reconfigure"):
    sys.stdout.reconfigure(encoding="utf-8", errors="replace")


def check(name: str, condition: bool, detail: str = "") -> None:
    if condition:
        PASSED.append(name)
        print(f"  PASS  {name}")
    else:
        FAILED.append(f"{name} - {detail}")
        print(f"  FAIL  {name} - {detail}")


# ---------------------------------------------------------------------------
# The recorded mapping: GPT-4.1's actual classification of this fixture.
# Column names only — every figure below is produced by deterministic code.
# ---------------------------------------------------------------------------

RECORDED_MAPPING = MappingProposal(
    table_shape="wide_entity",
    source_label="linkedin_ads_campaign_performance_july2026.csv",
    date_column=DateColumn(name="Start Date (in UTC)", format="%Y-%m-%d", confidence=0.95),
    entity_column=EntityColumn(name="Campaign Name", confidence=0.98),
    ignored_columns=[
        IgnoredColumn(
            name="End Date (in UTC)",
            reason="Redundant with Start Date for daily granularity; not a metric",
        )
    ],
    columns=[
        ColumnMapping(source_column="Impressions", target_metric="impressions",
                      label="Impressions", unit="number", confidence=0.99),
        ColumnMapping(source_column="Clicks", target_metric="clicks",
                      label="Clicks", unit="number", confidence=0.99),
        ColumnMapping(source_column="Click Through Rate", target_metric="ctr",
                      label="Click Through Rate", unit="percent", confidence=0.98),
        ColumnMapping(source_column="Average CPC", target_metric="cpc",
                      label="Average CPC", unit="currency",
                      direction="lower_is_better", confidence=0.97),
        ColumnMapping(source_column="Average CPM", target_metric="cpm",
                      label="Average CPM", unit="currency",
                      direction="lower_is_better", confidence=0.97),
        ColumnMapping(source_column="Total Spent", target_metric="spend",
                      label="Total Spent", unit="currency",
                      direction="lower_is_better", confidence=0.99),
        ColumnMapping(source_column="Conversions", target_metric="conversions",
                      label="Conversions", unit="number", confidence=0.98),
        ColumnMapping(source_column="Cost Per Conversion",
                      target_metric="cost_per_conversion",
                      label="Cost Per Conversion", unit="currency",
                      direction="lower_is_better", confidence=0.97),
        ColumnMapping(source_column="Conversion Rate", target_metric="conversion_rate",
                      label="Conversion Rate", unit="percent", confidence=0.97),
        ColumnMapping(source_column="Reactions", target_metric="reactions",
                      label="Reactions", unit="number", confidence=0.95),
        ColumnMapping(source_column="Comments", target_metric="comments",
                      label="Comments", unit="number", confidence=0.95),
        ColumnMapping(source_column="Shares", target_metric="shares",
                      label="Shares", unit="number", confidence=0.95),
        ColumnMapping(source_column="Follows", target_metric="follows",
                      label="Follows", unit="number", confidence=0.95),
    ],
)


# ---------------------------------------------------------------------------
# What the slide must SAY. Values from the fixture's own totals.
# ---------------------------------------------------------------------------

class Expect:
    def __init__(self, label: str, value: float, tolerance: float,
                 symbol: str | None = None, percent: bool = False):
        self.label = label
        self.value = value
        self.tolerance = tolerance
        self.symbol = symbol
        self.percent = percent


EXPECTED: dict[str, Expect] = {
    "impressions": Expect("IMPRESSIONS",        3_884_961,  0),
    "clicks":      Expect("CLICKS",                27_011,  0),
    "ctr":         Expect("CLICK THROUGH RATE",      0.70,  0.01, percent=True),
    "cpc":         Expect("AVERAGE CPC",             7.38,  0.01, symbol="$"),
    "cpm":         Expect("AVERAGE CPM",            51.34,  0.1,  symbol="$"),
    "spend":       Expect("TOTAL SPENT",       199_457.12,  0.01, symbol="$"),
}

# Values the deck rendered before this fix. Their reappearance anywhere on the
# slide means a regression, so they are asserted absent as well.
REGRESSION_STRINGS: dict[str, str] = {
    "1.04%":     "unweighted mean CTR (rates averaged across rows again)",
    "2,053,132": "half-period impressions headline",
    "14,318":    "half-period clicks headline",
    "107,135":   "half-period spend headline",
}


def _rendered_kpi_slide(
    pptx_path: str, tokens: dict[str, tuple[int, int]]
) -> dict[tuple[int, int], str]:
    """
    Text shapes on the deck's CSV KPI slide, keyed by position.

    The slide is identified by geometry, not by its text: the deck also has a
    generic KPI scorecard that says "IMPRESSIONS", and matching on the label
    picked that one instead — which is how the first run of this script
    reported six missing labels while the CSV slide beside it rendered all six
    correctly. Only the CSV slide carries shapes at the template's csv_kpi
    token coordinates, so those coordinates select it unambiguously.

    Token replacement rewrites a shape's text but never moves it, so the same
    coordinates then tie each rendered value back to its label.
    """
    wanted = set(tokens.values())
    if not wanted:
        return {}

    prs = Presentation(pptx_path)
    best: dict[tuple[int, int], str] = {}
    best_hits = 0
    for slide in prs.slides:
        positions = {
            (sh.left, sh.top): sh.text_frame.text
            for sh in slide.shapes if sh.has_text_frame
        }
        hits = len(wanted & set(positions))
        if hits > best_hits:
            best_hits, best = hits, positions
    return best if best_hits == len(wanted) else {}


def _all_deck_text(pptx_path: str) -> str:
    """Every word in the deck — regression guards check the whole file."""
    prs = Presentation(pptx_path)
    return " | ".join(
        sh.text_frame.text
        for slide in prs.slides for sh in slide.shapes if sh.has_text_frame
    )


def _template_token_positions() -> dict[str, tuple[int, int]]:
    """Where each {{csv_kpi_N_label/value}} token sits in the source template."""
    prs = Presentation(TEMPLATE)
    positions: dict[str, tuple[int, int]] = {}
    for slide in prs.slides:
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            text = shape.text_frame.text.strip()
            if re.fullmatch(r"\{\{csv_kpi_\d+_(label|value|change)\}\}", text):
                positions[text] = (shape.left, shape.top)
    return positions


def _parse_number(displayed: str) -> float | None:
    """The number a rendered string shows, ignoring symbols and separators."""
    cleaned = re.sub(r"[^\d.\-]", "", displayed)
    if not cleaned or cleaned in ("-", "."):
        return None
    try:
        return float(cleaned)
    except ValueError:
        return None


def build_source(live: bool) -> dict:
    with open(FIXTURE, "rb") as handle:
        content = handle.read()

    import hashlib
    digest = hashlib.md5(content).hexdigest()
    check("fixture is the expected file", digest == FIXTURE_MD5,
          f"md5 {digest} != {FIXTURE_MD5}")

    profile = profile_file(content, os.path.basename(FIXTURE))[0]

    if live:
        import asyncio
        from services.csv_ingest.mapper import propose_mapping
        mapping = asyncio.run(propose_mapping(profile, os.path.basename(FIXTURE)))
    else:
        mapping = RECORDED_MAPPING

    return normalize(profile, mapping, source_name=os.path.basename(FIXTURE))


def render(source: dict) -> str:
    """Generate a real deck. Account currency is INR on purpose."""
    from services.chart_generator import generate_all_charts
    from services.report_generator import generate_pptx_report

    data = {
        "client_name": "Acme Ltd",
        "period_start": "2026-07-01",
        "period_end": "2026-07-31",
        "csv_sources": [source],
        # The account bills in rupees while the upload states USD. If the deck
        # renders '$' anyway, per-source currency is genuinely winning; if the
        # account default were still in charge this would show '₹' and the
        # currency assertions below would fail.
        "meta_ads": {"currency": "INR"},
    }

    out_dir = tempfile.mkdtemp()
    charts = generate_all_charts(data, out_dir, "#4338CA", "modern_clean")
    pptx = generate_pptx_report(
        data,
        {"executive_summary": "Summary.", "csv_performance": "LinkedIn delivered.",
         "key_wins": ["Leads up"], "concerns": ["CPC rising"],
         "next_steps": ["Scale winners"]},
        charts,
        {"name": "Acme Ltd", "agency_name": "SapienBotics"},
        None, "full", None,
        {"agency_name": "SapienBotics", "brand_color": "#4338CA",
         "agency_logo_url": "", "client_logo_url": "", "powered_by_badge": False,
         "_cover_theme": "modern_clean"},
        "modern_clean", "en",
        {"theme": "modern_clean", "headline": "Acme Ltd", "subtitle": None,
         "brand_primary_color": None, "accent_color": None},
    )
    path = os.path.join(out_dir, "report.pptx")
    with open(path, "wb") as handle:
        handle.write(pptx)
    return path


def verify_rendered(pptx_path: str) -> None:
    print("\nRendered KPI slide")
    tokens = _template_token_positions()
    check("template exposes all 6 KPI slots", len(tokens) >= 12, f"found {len(tokens)}")

    by_position = _rendered_kpi_slide(pptx_path, tokens)
    check("CSV KPI slide located in the deck", bool(by_position),
          "no slide carried shapes at every csv_kpi token position")
    if not by_position:
        return

    slide_text = " | ".join(by_position.values())
    all_text = _all_deck_text(pptx_path)

    # Pair each label with the value shape that sits at the same slot.
    label_to_value: dict[str, str] = {}
    for index in range(6):
        label_pos = tokens.get(f"{{{{csv_kpi_{index}_label}}}}")
        value_pos = tokens.get(f"{{{{csv_kpi_{index}_value}}}}")
        if not label_pos or not value_pos:
            continue
        label = (by_position.get(label_pos) or "").strip().upper()
        value = (by_position.get(value_pos) or "").strip()
        if label:
            label_to_value[label] = value

    check("all six KPI slots rendered", len(label_to_value) == 6,
          f"found {len(label_to_value)}: {sorted(label_to_value)}")

    for key, expect in EXPECTED.items():
        displayed = label_to_value.get(expect.label)
        if displayed is None:
            check(f"{key}: label '{expect.label}' on slide", False,
                  f"labels present: {sorted(label_to_value)}")
            continue

        shown = _parse_number(displayed)
        if shown is None:
            check(f"{key}: displays a number", False, f"rendered {displayed!r}")
            continue

        within = abs(shown - expect.value) <= expect.tolerance
        check(
            f"{key}: deck shows {expect.value:,} (+/-{expect.tolerance})",
            within,
            f"deck shows {displayed!r} -> {shown:,}",
        )

        if expect.symbol:
            check(
                f"{key}: rendered in {expect.symbol}, not the account default",
                expect.symbol in displayed and "₹" not in displayed,
                f"rendered {displayed!r}",
            )
        if expect.percent:
            check(f"{key}: rendered with a % sign", "%" in displayed,
                  f"rendered {displayed!r}")

    print("\nRegression guards (whole deck)")
    for bad, reason in REGRESSION_STRINGS.items():
        check(f"'{bad}' absent from the deck ({reason})", bad not in all_text,
              f"found in: {all_text[:300]}")

    check("no rupee symbol on the CSV slide",
          "₹" not in slide_text, f"slide text: {slide_text[:300]}")


def verify_derivations(source: dict) -> None:
    """Report which metrics recomputed and which fell back to a weighted mean."""
    print("\nHow each metric was derived")
    report = source.get("derivations", {})
    fallbacks = [
        key for key, info in report.items()
        if info["method"] in ("weighted_mean", "unweighted_mean")
    ]
    for key in sorted(report):
        info = report[key]
        detail = info["detail"].replace("Σ", "sum ")
        print(f"    {key:22} {info['method']:16} {detail}")

    check("no metric fell back to an unweighted mean",
          not [k for k in report if report[k]["method"] == "unweighted_mean"],
          str([k for k in report if report[k]["method"] == "unweighted_mean"]))
    if fallbacks:
        print(f"    NOTE weighted-mean fallback used for: {', '.join(sorted(fallbacks))}")
    else:
        print("    NOTE every rate recomputed from its own components")


def verify_entity_breakdown(source: dict) -> None:
    """Per-entity rates must obey the same arithmetic as the headline."""
    print("\nEntity breakdown")
    breakdown = source.get("breakdown") or []
    check("breakdown present", bool(breakdown))
    impossible = [
        (row.get("name"), key, row[key])
        for row in breakdown
        for key in ("ctr", "conversion_rate")
        if isinstance(row.get(key), (int, float)) and row[key] > 100
    ]
    check("no per-entity rate exceeds 100%", not impossible, str(impossible[:3]))

    for row in breakdown[:4]:
        print(f"    {row.get('name', '?')[:44]:44} "
              f"ctr={row.get('ctr')}  cvr={row.get('conversion_rate')}")


def verify_currency(source: dict) -> None:
    print("\nCurrency detection")
    check("currency read from the file", source.get("currency") == "USD",
          f"detected {source.get('currency')!r} via {source.get('currency_source')!r}")
    print(f"    detected {source.get('currency')} from {source.get('currency_source')}")


def main() -> int:
    parser = argparse.ArgumentParser()
    parser.add_argument("--live", action="store_true",
                        help="call the real GPT-4.1 mapper instead of the recorded mapping")
    args = parser.parse_args()

    print(f"Fixture: {os.path.basename(FIXTURE)}")
    print(f"Mapping: {'live GPT-4.1' if args.live else 'recorded'}")

    source = build_source(args.live)
    verify_currency(source)
    verify_derivations(source)
    verify_entity_breakdown(source)

    path = render(source)
    print(f"\nDeck: {path}")
    verify_rendered(path)

    print(f"\n{len(PASSED)} passed, {len(FAILED)} failed")
    for failure in FAILED:
        print(f"  FAILED  {failure}")
    return 1 if FAILED else 0


if __name__ == "__main__":
    sys.exit(main())
