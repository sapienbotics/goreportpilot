"""
Focused audit: extract geometry of the CSV data slide (contains {{csv_source_name}})
across all 6 templates. Report:
  - KPI card shape bottoms
  - Chart placeholder position & size
  - Footer position
  - Overlap diagnosis
"""
from pptx import Presentation
import glob
import os

TPL_DIR = os.path.join(os.path.dirname(os.path.dirname(os.path.abspath(__file__))),
                       "templates", "pptx")

for tpl in sorted(glob.glob(os.path.join(TPL_DIR, "*.pptx"))):
    basename = os.path.basename(tpl)
    prs = Presentation(tpl)
    slide_w = round(prs.slide_width / 914400, 2)
    slide_h = round(prs.slide_height / 914400, 2)

    csv_slide_idx = None
    csv_slide = None
    for idx, slide in enumerate(prs.slides):
        for shape in slide.shapes:
            if shape.has_text_frame and "{{csv_source_name}}" in shape.text_frame.text:
                csv_slide_idx = idx
                csv_slide = slide
                break
        if csv_slide:
            break

    if csv_slide is None:
        print(f"\n{basename}: NO CSV SLIDE FOUND\n")
        continue

    print(f"\n{'=' * 80}")
    print(f"{basename}  —  slide {csv_slide_idx}  ({slide_w}\" x {slide_h}\")")
    print(f"{'=' * 80}")

    kpi_shapes = []
    chart_shape = None
    footer_shape = None
    title_shape = None
    other_shapes = []

    for shape in csv_slide.shapes:
        t = round(shape.top / 914400, 2) if shape.top else 0
        l = round(shape.left / 914400, 2) if shape.left else 0
        w = round(shape.width / 914400, 2) if shape.width else 0
        h = round(shape.height / 914400, 2) if shape.height else 0
        b = round(t + h, 2)
        text = ""
        if shape.has_text_frame:
            text = shape.text_frame.text.strip()

        info = {
            "id": shape.shape_id,
            "left": l, "top": t, "width": w, "height": h, "bottom": b,
            "text": text[:80],
        }

        if "chart_csv" in text:
            chart_shape = info
        elif "csv_kpi_" in text or "csv_source_name" in text:
            kpi_shapes.append(info)
        elif "agency_name" in text.lower() and ("page" in text.lower() or "confidential" in text.lower()):
            footer_shape = info
        elif text:
            # Could be a KPI card bg, header, etc
            other_shapes.append(info)
        else:
            other_shapes.append(info)

    # ── KPI shapes ────────────────────────────────────────────────
    kpi_bottoms = []
    print("\n  KPI/CSV text shapes:")
    for s in sorted(kpi_shapes, key=lambda x: (x["top"], x["left"])):
        print(f"    id={s['id']:3d}  pos=({s['left']:5.2f}, {s['top']:5.2f})  "
              f"size=({s['width']:5.2f} x {s['height']:5.2f})  "
              f"bottom={s['bottom']:5.2f}  text=\"{s['text']}\"")
        kpi_bottoms.append(s["bottom"])

    # ── Chart placeholder ─────────────────────────────────────────
    if chart_shape:
        print(f"\n  Chart placeholder:")
        print(f"    id={chart_shape['id']:3d}  pos=({chart_shape['left']:5.2f}, {chart_shape['top']:5.2f})  "
              f"size=({chart_shape['width']:5.2f} x {chart_shape['height']:5.2f})  "
              f"bottom={chart_shape['bottom']:5.2f}  text=\"{chart_shape['text']}\"")
    else:
        print("\n  Chart placeholder: NOT FOUND")

    # ── Footer ────────────────────────────────────────────────────
    if footer_shape:
        print(f"\n  Footer:")
        print(f"    id={footer_shape['id']:3d}  pos=({footer_shape['left']:5.2f}, {footer_shape['top']:5.2f})  "
              f"size=({footer_shape['width']:5.2f} x {footer_shape['height']:5.2f})  "
              f"bottom={footer_shape['bottom']:5.2f}  text=\"{footer_shape['text']}\"")
    else:
        print("\n  Footer: NOT FOUND")

    # ── Background / decorative shapes ────────────────────────────
    print(f"\n  Other shapes: {len(other_shapes)}")
    for s in sorted(other_shapes, key=lambda x: (x["top"], x["left"])):
        marker = ""
        if s["text"]:
            marker = f"  text=\"{s['text'][:50]}\""
        print(f"    id={s['id']:3d}  pos=({s['left']:5.2f}, {s['top']:5.2f})  "
              f"size=({s['width']:5.2f} x {s['height']:5.2f})  "
              f"bottom={s['bottom']:5.2f}{marker}")

    # ── Diagnosis ────────────────────────────────────────────────
    print(f"\n  === DIAGNOSIS ===")
    max_kpi_bottom = max(kpi_bottoms) if kpi_bottoms else 0
    print(f"  Max KPI bottom edge:       {max_kpi_bottom:.2f}\"")
    if chart_shape:
        print(f"  Chart placeholder top:     {chart_shape['top']:.2f}\"")
        gap = chart_shape["top"] - max_kpi_bottom
        print(f"  Gap (chart top - KPI bottom): {gap:+.2f}\"")
        if gap < 0:
            print(f"  ** OVERLAP: chart starts {abs(gap):.2f}\" ABOVE last KPI bottom **")
        elif gap < 0.1:
            print(f"  ** TIGHT: only {gap:.2f}\" gap — may look cramped **")
        else:
            print(f"  OK: {gap:.2f}\" clearance")
    if footer_shape:
        print(f"  Footer top:                {footer_shape['top']:.2f}\"")
        if chart_shape:
            chart_to_footer = footer_shape["top"] - chart_shape["bottom"]
            print(f"  Chart bottom to footer:    {chart_to_footer:+.2f}\"")
    print()
