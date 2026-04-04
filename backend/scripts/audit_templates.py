"""
Audit all PPTX templates — print every shape's position, size, and text
on every slide so we can see exact geometry (inches).
Focus: CSV data slide (contains {{csv_source_name}}).
"""
from pptx import Presentation
from pptx.util import Emu
import glob
import os

TPL_DIR = os.path.join(os.path.dirname(os.path.dirname(os.path.abspath(__file__))),
                       "templates", "pptx")

for tpl in sorted(glob.glob(os.path.join(TPL_DIR, "*.pptx"))):
    basename = os.path.basename(tpl)
    print(f'\n{"=" * 90}\n{basename}\n{"=" * 90}')
    prs = Presentation(tpl)
    print(f"  Slide dimensions: {round(prs.slide_width / 914400, 2)}\" x {round(prs.slide_height / 914400, 2)}\"")

    for idx, slide in enumerate(prs.slides):
        # Find first text for slide identification
        title = ""
        for s in slide.shapes:
            if s.has_text_frame and s.text_frame.text.strip():
                title = s.text_frame.text.strip()[:60]
                break
        print(f"\n  Slide {idx} — \"{title}\"")

        for shape in slide.shapes:
            l = round(shape.left / 914400, 2) if shape.left else 0
            t = round(shape.top / 914400, 2) if shape.top else 0
            w = round(shape.width / 914400, 2) if shape.width else 0
            h = round(shape.height / 914400, 2) if shape.height else 0
            bottom = round(t + h, 2)
            text = ""
            if shape.has_text_frame:
                text = shape.text_frame.text[:120].replace("\n", " | ")
            print(f"    id={shape.shape_id:3d}  pos=({l:6.2f}, {t:6.2f})  "
                  f"size=({w:6.2f} x {h:6.2f})  bottom={bottom:6.2f}  "
                  f"text=\"{text}\"")
