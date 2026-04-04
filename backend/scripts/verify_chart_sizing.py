"""Verify chart sizing fixes: CSV chart centering + dual-chart solo centering."""
import os, sys, io, tempfile
sys.stdout = io.TextIOWrapper(sys.stdout.buffer, encoding='utf-8')
sys.path.insert(0, os.path.join(os.path.dirname(os.path.dirname(os.path.abspath(__file__)))))

from pptx import Presentation

EMU = 914400

from services.report_generator import generate_pptx_report
from services.chart_generator import generate_all_charts
from services.mock_data import generate_all_mock_data

raw_data = generate_all_mock_data("Test Client", "2026-03-01", "2026-03-30")
# Remove device_breakdown to force solo chart on Website Engagement slide
raw_data["ga4"]["device_breakdown"] = []
raw_data["csv_sources"] = [
    {"source_name": "TikTok Ads", "metrics": [
        {"name": "Impressions", "current_value": 320000, "previous_value": 275000, "unit": "number"},
        {"name": "Clicks", "current_value": 9800, "previous_value": 8100, "unit": "number"},
        {"name": "CTR", "current_value": 3.06, "previous_value": 2.95, "unit": "percent"},
        {"name": "Spend", "current_value": 2400, "previous_value": 2100, "unit": "currency"},
    ]},
]

with tempfile.TemporaryDirectory() as tmpdir:
    charts = generate_all_charts(raw_data, os.path.join(tmpdir, "charts"), "#4338CA", "modern_clean")
    narrative = {
        "executive_summary": "Summary.", "website_performance": "Web OK.",
        "paid_advertising": "Ads OK.", "key_wins": "Win 1", "concerns": "Issue 1",
        "next_steps": "Step 1",
    }
    pptx_bytes = generate_pptx_report(
        data=raw_data, narrative=narrative, charts=charts,
        client_info={"name": "Test", "agency_name": "Agency"},
        visual_template="modern_clean",
    )
    path = os.path.join(tmpdir, "test.pptx")
    with open(path, "wb") as f:
        f.write(pptx_bytes)

    prs = Presentation(path)
    slide_w = prs.slide_width / EMU

    print(f"Slide dimensions: {slide_w:.2f}\" × {prs.slide_height / EMU:.2f}\"")
    print(f"Total slides: {len(prs.slides)}\n")

    for idx, slide in enumerate(prs.slides):
        title = ""
        for shape in slide.shapes:
            if shape.has_text_frame and shape.text_frame.text.strip():
                title = shape.text_frame.text.strip()[:40]
                break

        # Find images (embedded charts)
        for shape in slide.shapes:
            try:
                if shape.image:
                    l = round(shape.left / EMU, 2)
                    t = round(shape.top / EMU, 2)
                    w = round(shape.width / EMU, 2)
                    h = round(shape.height / EMU, 2)
                    ratio = w / h if h > 0 else 0
                    center_offset = abs(l + w / 2 - slide_w / 2)
                    centered = "CENTERED" if center_offset < 0.15 else f"off-center by {center_offset:.2f}\""
                    print(f"  Slide {idx:2d} ({title[:25]:25s})  "
                          f"img at ({l:.2f}\", {t:.2f}\") "
                          f"size {w:.2f}\" × {h:.2f}\" "
                          f"ratio={ratio:.1f}:1 "
                          f"{centered}")
            except Exception:
                pass

    # Specific checks
    print("\n=== CHECKS ===")

    # Check CSV chart sizing
    for idx, slide in enumerate(prs.slides):
        for shape in slide.shapes:
            if shape.has_text_frame and shape.text_frame.text.strip() == "TikTok Ads":
                # Found CSV slide — check its image
                for s2 in slide.shapes:
                    try:
                        if s2.image:
                            w = s2.width / EMU
                            l = s2.left / EMU
                            center_off = abs(l + w / 2 - slide_w / 2)
                            ok = w <= 8.5 and center_off < 0.15
                            print(f"CSV chart: width={w:.2f}\" left={l:.2f}\" "
                                  f"centered={center_off < 0.15} "
                                  f"{'PASS' if ok else 'FAIL (too wide or off-center)'}")
                    except Exception:
                        pass
                break

    # Check solo chart (Website Engagement with no device_breakdown)
    for idx, slide in enumerate(prs.slides):
        for shape in slide.shapes:
            if shape.has_text_frame and "Engagement" in shape.text_frame.text:
                for s2 in slide.shapes:
                    try:
                        if s2.image:
                            w = s2.width / EMU
                            l = s2.left / EMU
                            center_off = abs(l + w / 2 - slide_w / 2)
                            ok = w <= 9.0 and center_off < 0.15
                            print(f"Solo chart (Engagement): width={w:.2f}\" left={l:.2f}\" "
                                  f"centered={center_off < 0.15} "
                                  f"{'PASS' if ok else 'FAIL'}")
                    except Exception:
                        pass
                break
