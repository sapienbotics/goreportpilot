"""
Phase 3 parity v2 — C-fix2. Add subtle non-text chrome to
``minimal_elegant.pptx``'s cover so the theme has recognisable
identity in the Design-tab picker.

Rationale
---------
The v1 chrome-only strip left ``minimal_elegant`` with a single 0.01"
horizontal divider — the theme's entire visual identity had been
carried by its typography + placeholder text shapes. Every other
theme's decorative chrome (header bands, gradient strips, geometric
blocks) survived the strip naturally; minimal_elegant's did not.

This script adds two subtle decorative rules in the editorial style
the theme is named for:

  1. A 0.04" horizontal rule at the top margin (y=0.55, width 2.0").
     Acts as a header mark above where the agency logo sits.
  2. Thickens the existing 0.01" divider at y=4.00 to 0.04" so it
     survives thumbnail rasterisation.

Net effect in picker thumbnails: two short thin horizontal rules
bracketing the whitespace — instantly recognisable as editorial /
serif-era typography without cluttering the cover.

Idempotent: re-running on an already-enhanced template detects the
existing rules by coordinates and skips re-adding.

Run
---

    python backend/scripts/enhance_minimal_elegant_chrome.py

Then re-run ``regenerate_cover_thumbnails.py`` to refresh the PNG.
"""
from __future__ import annotations

import sys
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.util import Emu, Inches

BACKEND = Path(__file__).resolve().parents[1]
TEMPLATE = BACKEND / "templates" / "pptx" / "minimal_elegant.pptx"

# Editorial slate — matches the theme's existing text colour palette.
RULE_COLOUR = (30, 41, 59)   # #1E293B (slate-800)

# Shape positions (inches). Short centered rules, 2" wide.
TOP_RULE    = {"x": 5.65, "y": 0.55, "w": 2.00, "h": 0.04}
MID_RULE_Y  = 4.00           # existing divider y-position
MID_RULE_H  = 0.04           # new height (from 0.01)


def _emu(v: float) -> int:
    return int(v * 914400)


def _has_shape_at(cover, target_x: float, target_y: float, target_w: float,
                  tolerance: float = 0.05) -> bool:
    """True if a shape already sits approximately at (target_x, target_y)."""
    for sh in cover.shapes:
        try:
            L = (sh.left or 0) / 914400
            T = (sh.top  or 0) / 914400
            W = (sh.width or 0) / 914400
            if (abs(L - target_x) < tolerance
                and abs(T - target_y) < tolerance
                and abs(W - target_w) < 0.25):
                return True
        except Exception:
            continue
    return False


def _add_rule(cover, x: float, y: float, w: float, h: float, rgb: tuple) -> None:
    """Add a filled rectangle of the given dimensions + colour."""
    shape = cover.shapes.add_shape(
        1,  # MSO_SHAPE.RECTANGLE
        Emu(_emu(x)), Emu(_emu(y)),
        Emu(_emu(w)), Emu(_emu(h)),
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(*rgb)
    try:
        shape.line.fill.background()
    except Exception:
        pass


def _thicken_existing_divider(cover) -> bool:
    """Thicken the existing ~0.01" divider at y~4.00 to 0.04"."""
    for sh in cover.shapes:
        try:
            T = (sh.top or 0) / 914400
            H = (sh.height or 0) / 914400
            # Match the existing divider by y-position + small height
            if abs(T - MID_RULE_Y) < 0.1 and 0 < H < 0.03:
                new_h = _emu(MID_RULE_H)
                sh.height = Emu(new_h)
                # Also recolour to the editorial slate.
                try:
                    sh.fill.solid()
                    sh.fill.fore_color.rgb = RGBColor(*RULE_COLOUR)
                except Exception:
                    pass
                return True
        except Exception:
            continue
    return False


def main() -> int:
    if not TEMPLATE.exists():
        print(f"[ERR] Template not found: {TEMPLATE}", file=sys.stderr)
        return 2

    prs = Presentation(str(TEMPLATE))
    cover = prs.slides[0]

    print(f"Cover shape count before: {len(cover.shapes)}")

    # 1. Thicken existing mid-slide divider.
    if _thicken_existing_divider(cover):
        print(f"  thickened existing divider at y={MID_RULE_Y} to height {MID_RULE_H}\"")
    else:
        print("  (no existing divider found near y=4.00)")

    # 2. Add top rule (idempotent — skip if already present).
    if _has_shape_at(cover, TOP_RULE["x"], TOP_RULE["y"], TOP_RULE["w"]):
        print(f"  top rule already present at ({TOP_RULE['x']}, {TOP_RULE['y']}) — skipped")
    else:
        _add_rule(cover, **TOP_RULE, rgb=RULE_COLOUR)
        print(f"  added top rule at ({TOP_RULE['x']}, {TOP_RULE['y']}, {TOP_RULE['w']}x{TOP_RULE['h']})")

    prs.save(str(TEMPLATE))
    print(f"Cover shape count after: {len(Presentation(str(TEMPLATE)).slides[0].shapes)}")
    print(f"Saved: {TEMPLATE}")
    return 0


if __name__ == "__main__":
    sys.exit(main())
